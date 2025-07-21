"""
Sistema de Processamento de Arquivos
Processa PDFs e PPTXs para extrair conteúdo e gerar vídeos/cursos
"""

import asyncio
import os
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import aiofiles
import PyPDF2
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import openai
from moviepy.editor import VideoFileClip, AudioFileClip, CompositeVideoClip, ImageClip
import xml.etree.ElementTree as ET

from database import get_db
from models import FileUpload, Video, Project
from logger import logger
from utils import send_notification
from config import settings

# Importar novo serviço TTS
from services.tts_service import tts_service, TTSConfig, TTSProvider, generate_course_narration


class FileProcessor:
    """Processador principal de arquivos"""
    
    def __init__(self):
        self.temp_dir = Path(tempfile.gettempdir()) / "tecnocursos_processing"
        self.temp_dir.mkdir(exist_ok=True)
        
        # Configurar OpenAI se disponível
        if settings.OPENAI_API_KEY:
            openai.api_key = settings.OPENAI_API_KEY
        
        # Configurar TTS padrão
        self.tts_config = TTSConfig(
            provider=TTSProvider.AUTO,  # Auto-seleciona melhor provider
            language="pt",
            voice="pt_speaker_0",  # Voz padrão para Bark
            output_format="mp3"
        )
    
    async def process_file(self, file_id: str) -> Dict:
        """Processa um arquivo baseado no seu tipo"""
        db = next(get_db())
        try:
            file_upload = db.query(FileUpload).filter(FileUpload.id == file_id).first()
            if not file_upload:
                raise ValueError(f"Arquivo {file_id} não encontrado")
            
            logger.info(f"Iniciando processamento do arquivo {file_upload.filename}")
            
            # Atualizar status
            file_upload.status = "processing"
            file_upload.processing_progress = 10
            db.commit()
            
            # Processar baseado no tipo
            file_path = Path(file_upload.file_path)
            extension = file_path.suffix.lower()
            
            result = {}
            
            if extension == '.pdf':
                result = await self._process_pdf(file_upload, db)
            elif extension in ['.pptx', '.ppt']:
                result = await self._process_presentation(file_upload, db)
            else:
                raise ValueError(f"Tipo de arquivo não suportado: {extension}")
            
            # Finalizar processamento
            file_upload.status = "processed"
            file_upload.processing_progress = 100
            file_upload.metadata = result
            db.commit()
            
            # Notificar usuário
            await send_notification(
                file_upload.project.user_id,
                f"Arquivo {file_upload.filename} processado com sucesso!",
                "file_processed"
            )
            
            logger.info(f"Processamento concluído para {file_upload.filename}")
            return result
            
        except Exception as e:
            logger.error(f"Erro no processamento do arquivo {file_id}: {str(e)}")
            if 'file_upload' in locals():
                file_upload.status = "error"
                file_upload.error_message = str(e)
                db.commit()
            raise
        finally:
            db.close()
    
    async def _process_pdf(self, file_upload: FileUpload, db) -> Dict:
        """Processa arquivo PDF"""
        file_path = Path(file_upload.file_path)
        
        # Extrair texto do PDF
        text_content = await self._extract_pdf_text(file_path)
        
        # Extrair imagens
        images = await self._extract_pdf_images(file_path)
        
        # Atualizar progresso
        file_upload.processing_progress = 30
        db.commit()
        
        # Gerar estrutura do curso usando IA
        course_structure = await self._generate_course_structure(text_content)
        
        # Atualizar progresso
        file_upload.processing_progress = 50
        db.commit()
        
        # Gerar áudio para cada seção usando novo TTS
        audio_files = await self._generate_audio_content_advanced(course_structure)
        
        # Atualizar progresso
        file_upload.processing_progress = 70
        db.commit()
        
        # Criar vídeo
        video_path = await self._create_video_from_content(
            course_structure, audio_files, images, file_upload
        )
        
        # Atualizar progresso
        file_upload.processing_progress = 90
        db.commit()
        
        # Salvar vídeo no banco
        video = Video(
            project_id=file_upload.project_id,
            file_upload_id=file_upload.id,
            title=f"Curso - {file_upload.filename}",
            video_path=str(video_path),
            duration=await self._get_video_duration(video_path),
            status="ready"
        )
        db.add(video)
        db.commit()
        
        return {
            "type": "pdf_course",
            "text_length": len(text_content),
            "images_count": len(images),
            "course_sections": len(course_structure),
            "video_id": video.id,
            "video_path": str(video_path),
            "structure": course_structure,
            "tts_provider": "bark_advanced"  # Indicar que usou TTS avançado
        }
    
    async def _process_presentation(self, file_upload: FileUpload, db) -> Dict:
        """Processa apresentação PowerPoint"""
        file_path = Path(file_upload.file_path)
        
        # Carregar apresentação
        presentation = Presentation(file_path)
        
        # Extrair conteúdo dos slides
        slides_content = await self._extract_slides_content(presentation)
        
        # Atualizar progresso
        file_upload.processing_progress = 40
        db.commit()
        
        # Gerar narração para os slides usando novo TTS
        narration = await self._generate_slides_narration_advanced(slides_content)
        
        # Atualizar progresso
        file_upload.processing_progress = 60
        db.commit()
        
        # Exportar slides como imagens
        slide_images = await self._export_slides_as_images(presentation, file_upload.id)
        
        # Atualizar progresso
        file_upload.processing_progress = 80
        db.commit()
        
        # Criar vídeo da apresentação
        video_path = await self._create_presentation_video(
            slide_images, narration, file_upload
        )
        
        # Salvar vídeo no banco
        video = Video(
            project_id=file_upload.project_id,
            file_upload_id=file_upload.id,
            title=f"Apresentação - {file_upload.filename}",
            video_path=str(video_path),
            duration=await self._get_video_duration(video_path),
            status="ready"
        )
        db.add(video)
        db.commit()
        
        return {
            "type": "presentation_video",
            "slides_count": len(slides_content),
            "video_id": video.id,
            "video_path": str(video_path),
            "slides_content": slides_content,
            "narration": narration,
            "tts_provider": "bark_advanced"  # Indicar que usou TTS avançado
        }
    
    async def _extract_pdf_text(self, file_path: Path) -> str:
        """Extrai texto de um PDF"""
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    async def _extract_pdf_images(self, file_path: Path) -> List[str]:
        """Extrai imagens de um PDF"""
        images = []
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        image_path = self.temp_dir / f"pdf_image_{page_num}_{img_index}.png"
                        pix.save(str(image_path))
                        images.append(str(image_path))
                    
                    pix = None
            doc.close()
        except ImportError:
            logger.warning("PyMuPDF não disponível, pulando extração de imagens")
        
        return images
    
    async def _generate_course_structure(self, text_content: str) -> List[Dict]:
        """Gera estrutura do curso usando IA"""
        if not settings.OPENAI_API_KEY:
            # Estrutura básica sem IA
            chunks = self._split_text_into_chunks(text_content, 1000)
            return [
                {
                    "title": f"Seção {i+1}",
                    "content": chunk,
                    "duration": min(len(chunk) / 150, 30)  # ~150 chars por segundo, max 30s
                }
                for i, chunk in enumerate(chunks)
            ]
        
        try:
            response = await openai.ChatCompletion.acreate(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "system",
                        "content": "Você é um especialista em criação de cursos. Analise o texto e crie uma estrutura de curso dividida em seções lógicas."
                    },
                    {
                        "role": "user",
                        "content": f"Crie uma estrutura de curso baseada neste texto:\n\n{text_content[:3000]}"
                    }
                ],
                max_tokens=1500
            )
            
            # Processar resposta da IA
            ai_response = response.choices[0].message.content
            return self._parse_ai_course_structure(ai_response, text_content)
            
        except Exception as e:
            logger.error(f"Erro ao gerar estrutura com IA: {e}")
            # Fallback para estrutura básica
            chunks = self._split_text_into_chunks(text_content, 1000)
            return [
                {
                    "title": f"Seção {i+1}",
                    "content": chunk,
                    "duration": min(len(chunk) / 150, 30)
                }
                for i, chunk in enumerate(chunks)
            ]
    
    def _split_text_into_chunks(self, text: str, chunk_size: int) -> List[str]:
        """Divide texto em chunks"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        
        for word in words:
            if current_size + len(word) > chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = [word]
                current_size = len(word)
            else:
                current_chunk.append(word)
                current_size += len(word) + 1
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks
    
    def _parse_ai_course_structure(self, ai_response: str, original_text: str) -> List[Dict]:
        """Processa resposta da IA para estrutura do curso"""
        # Implementação simples - pode ser melhorada
        sections = []
        chunks = self._split_text_into_chunks(original_text, 800)
        
        for i, chunk in enumerate(chunks):
            sections.append({
                "title": f"Seção {i+1}: {chunk[:50]}...",
                "content": chunk,
                "duration": min(len(chunk) / 150, 30)
            })
        
        return sections
    
    async def _generate_audio_content_advanced(self, course_structure: List[Dict]) -> List[str]:
        """Gera áudio para cada seção do curso usando TTS avançado"""
        logger.info(f"Gerando narração para {len(course_structure)} seções com TTS avançado")
        
        try:
            # Usar o novo serviço TTS unificado
            results = await generate_course_narration(
                course_structure,
                voice=self.tts_config.voice
            )
            
            audio_files = []
            for i, result in enumerate(results):
                if result.success and result.audio_path:
                    audio_files.append(result.audio_path)
                    logger.info(f"✅ Seção {i+1}: {result.provider_used}, {result.duration:.2f}s")
                else:
                    logger.error(f"❌ Erro na seção {i+1}: {result.error}")
                    # Criar áudio em branco como fallback
                    silence_path = await self._create_silence_audio(course_structure[i].get("duration", 5.0))
                    audio_files.append(silence_path)
            
            return audio_files
            
        except Exception as e:
            logger.error(f"Erro na geração avançada de áudio: {e}")
            # Fallback para método antigo
            return await self._generate_audio_content_fallback(course_structure)
    
    async def _generate_audio_content_fallback(self, course_structure: List[Dict]) -> List[str]:
        """Método fallback usando gTTS diretamente"""
        logger.warning("Usando fallback gTTS para geração de áudio")
        
        audio_files = []
        
        for i, section in enumerate(course_structure):
            try:
                # Usar gTTS como fallback
                from gtts import gTTS
                tts = gTTS(text=section["content"], lang='pt', slow=False)
                audio_path = self.temp_dir / f"section_{i}_audio_fallback.mp3"
                tts.save(str(audio_path))
                audio_files.append(str(audio_path))
                
            except Exception as e:
                logger.error(f"Erro ao gerar áudio para seção {i}: {e}")
                # Criar áudio em branco como fallback
                silence_path = await self._create_silence_audio(section.get("duration", 5.0))
                audio_files.append(silence_path)
        
        return audio_files
    
    async def _create_silence_audio(self, duration: float) -> str:
        """Cria áudio em branco"""
        try:
            import pydub
            silence = pydub.AudioSegment.silent(duration=int(duration * 1000))
            audio_path = self.temp_dir / f"silence_{duration}s.mp3"
            silence.export(str(audio_path), format="mp3")
            return str(audio_path)
        except ImportError:
            logger.warning("Pydub não disponível, usando arquivo de áudio vazio")
            return ""
    
    async def _create_video_from_content(
        self, 
        course_structure: List[Dict], 
        audio_files: List[str],
        images: List[str],
        file_upload: FileUpload
    ) -> Path:
        """Cria vídeo a partir do conteúdo"""
        try:
            clips = []
            current_time = 0
            
            for i, (section, audio_file) in enumerate(zip(course_structure, audio_files)):
                if not audio_file or not os.path.exists(audio_file):
                    continue
                
                # Criar slide de texto
                slide_image = await self._create_text_slide(
                    section["title"], 
                    section["content"][:200] + "..." if len(section["content"]) > 200 else section["content"]
                )
                
                # Carregar áudio
                audio_clip = AudioFileClip(audio_file)
                duration = audio_clip.duration
                
                # Criar clip de vídeo
                image_clip = ImageClip(slide_image, duration=duration)
                video_clip = image_clip.set_audio(audio_clip)
                clips.append(video_clip)
                
                current_time += duration
            
            if not clips:
                raise ValueError("Nenhum clip de vídeo foi criado")
            
            # Concatenar clips
            final_video = CompositeVideoClip(clips, method="compose")
            
            # Salvar vídeo
            video_path = Path(settings.UPLOAD_DIR) / "videos" / f"{file_upload.id}_course.mp4"
            video_path.parent.mkdir(exist_ok=True)
            
            final_video.write_videofile(
                str(video_path),
                fps=24,
                codec='libx264',
                audio_codec='aac'
            )
            
            # Limpar clips
            for clip in clips:
                clip.close()
            final_video.close()
            
            return video_path
            
        except Exception as e:
            logger.error(f"Erro ao criar vídeo: {e}")
            # Criar vídeo placeholder
            return await self._create_placeholder_video(file_upload)
    
    async def _create_text_slide(self, title: str, content: str) -> str:
        """Cria imagem de slide com texto"""
        try:
            # Criar imagem
            width, height = 1920, 1080
            image = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(image)
            
            # Tentar carregar fonte
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 48)
                content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 32)
            except:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
            
            # Desenhar título
            title_bbox = draw.textbbox((0, 0), title, font=title_font)
            title_width = title_bbox[2] - title_bbox[0]
            title_x = (width - title_width) // 2
            draw.text((title_x, 100), title, fill='black', font=title_font)
            
            # Desenhar conteúdo
            words = content.split()
            lines = []
            current_line = []
            
            for word in words:
                test_line = " ".join(current_line + [word])
                bbox = draw.textbbox((0, 0), test_line, font=content_font)
                line_width = bbox[2] - bbox[0]
                
                if line_width <= width - 200:  # Margem de 100px cada lado
                    current_line.append(word)
                else:
                    if current_line:
                        lines.append(" ".join(current_line))
                    current_line = [word]
            
            if current_line:
                lines.append(" ".join(current_line))
            
            # Desenhar linhas de conteúdo
            y_offset = 200
            for line in lines[:10]:  # Máximo 10 linhas
                line_bbox = draw.textbbox((0, 0), line, font=content_font)
                line_width = line_bbox[2] - line_bbox[0]
                line_x = (width - line_width) // 2
                draw.text((line_x, y_offset), line, fill='black', font=content_font)
                y_offset += 50
            
            # Salvar imagem
            slide_path = self.temp_dir / f"slide_{hash(title)}.png"
            image.save(str(slide_path))
            return str(slide_path)
            
        except Exception as e:
            logger.error(f"Erro ao criar slide: {e}")
            # Criar imagem em branco como fallback
            blank_path = self.temp_dir / "blank_slide.png"
            blank_image = Image.new('RGB', (1920, 1080), color='white')
            blank_image.save(str(blank_path))
            return str(blank_path)
    
    async def _extract_slides_content(self, presentation: Presentation) -> List[Dict]:
        """Extrai conteúdo dos slides"""
        slides_content = []
        
        for i, slide in enumerate(presentation.slides):
            slide_data = {
                "slide_number": i + 1,
                "title": "",
                "content": "",
                "images": [],
                "notes": ""
            }
            
            # Extrair texto das formas
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if not slide_data["title"]:
                        slide_data["title"] = shape.text.strip()
                    else:
                        slide_data["content"] += shape.text.strip() + "\n"
                
                # Verificar se é imagem
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_data["images"].append(f"Imagem no slide {i+1}")
            
            # Extrair notas do slide
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text"):
                        slide_data["notes"] += shape.text.strip()
            
            slides_content.append(slide_data)
        
        return slides_content
    
    async def _generate_slides_narration_advanced(self, slides_content: List[Dict]) -> List[str]:
        """Gera narração para os slides usando TTS avançado"""
        logger.info(f"Gerando narração para {len(slides_content)} slides com TTS avançado")
        
        try:
            # Preparar textos para narração
            narration_texts = []
            for slide in slides_content:
                narration_text = ""
                if slide["title"]:
                    narration_text += f"Slide {slide['slide_number']}: {slide['title']}. "
                if slide["content"]:
                    narration_text += slide["content"]
                if slide["notes"]:
                    narration_text += " " + slide["notes"]
                
                if not narration_text.strip():
                    narration_text = f"Slide {slide['slide_number']}"
                
                narration_texts.append(narration_text)
            
            # Usar o serviço TTS unificado
            results = await tts_service.generate_batch_speech(narration_texts, self.tts_config)
            
            narration_files = []
            for i, result in enumerate(results):
                if result.success and result.audio_path:
                    narration_files.append(result.audio_path)
                    logger.info(f"✅ Slide {i+1}: {result.provider_used}, {result.duration:.2f}s")
                else:
                    logger.error(f"❌ Erro no slide {i+1}: {result.error}")
                    # Criar áudio em branco
                    silence_path = await self._create_silence_audio(3.0)  # 3 segundos
                    narration_files.append(silence_path)
            
            return narration_files
            
        except Exception as e:
            logger.error(f"Erro na geração avançada de narração: {e}")
            # Fallback para método antigo
            return await self._generate_slides_narration_fallback(slides_content)
    
    async def _generate_slides_narration_fallback(self, slides_content: List[Dict]) -> List[str]:
        """Método fallback para narração usando gTTS"""
        logger.warning("Usando fallback gTTS para narração de slides")
        
        narration_files = []
        
        for i, slide in enumerate(slides_content):
            try:
                # Preparar texto para narração
                narration_text = ""
                if slide["title"]:
                    narration_text += f"Slide {slide['slide_number']}: {slide['title']}. "
                if slide["content"]:
                    narration_text += slide["content"]
                if slide["notes"]:
                    narration_text += " " + slide["notes"]
                
                if not narration_text.strip():
                    narration_text = f"Slide {slide['slide_number']}"
                
                # Gerar áudio com gTTS
                from gtts import gTTS
                tts = gTTS(text=narration_text, lang='pt', slow=False)
                audio_path = self.temp_dir / f"slide_{i}_narration_fallback.mp3"
                tts.save(str(audio_path))
                narration_files.append(str(audio_path))
                
            except Exception as e:
                logger.error(f"Erro ao gerar narração para slide {i}: {e}")
                # Criar áudio em branco
                silence_path = await self._create_silence_audio(3.0)  # 3 segundos
                narration_files.append(silence_path)
        
        return narration_files
    
    async def _export_slides_as_images(self, presentation: Presentation, file_id: str) -> List[str]:
        """Exporta slides como imagens"""
        slide_images = []
        
        # Para PowerPoint, seria necessário usar COM automation ou outro método
        # Por simplicidade, vamos criar slides baseados no conteúdo de texto
        slides_content = await self._extract_slides_content(presentation)
        
        for i, slide in enumerate(slides_content):
            slide_image = await self._create_text_slide(
                slide["title"] or f"Slide {i+1}",
                slide["content"]
            )
            slide_images.append(slide_image)
        
        return slide_images
    
    async def _create_presentation_video(
        self,
        slide_images: List[str],
        narration: List[str],
        file_upload: FileUpload
    ) -> Path:
        """Cria vídeo da apresentação"""
        try:
            clips = []
            
            for slide_image, audio_file in zip(slide_images, narration):
                if not os.path.exists(slide_image) or not os.path.exists(audio_file):
                    continue
                
                # Carregar áudio para obter duração
                audio_clip = AudioFileClip(audio_file)
                duration = max(audio_clip.duration, 3.0)  # Mínimo 3 segundos
                
                # Criar clip de imagem
                image_clip = ImageClip(slide_image, duration=duration)
                video_clip = image_clip.set_audio(audio_clip)
                clips.append(video_clip)
            
            if not clips:
                raise ValueError("Nenhum clip foi criado")
            
            # Concatenar clips
            final_video = CompositeVideoClip(clips, method="compose")
            
            # Salvar vídeo
            video_path = Path(settings.UPLOAD_DIR) / "videos" / f"{file_upload.id}_presentation.mp4"
            video_path.parent.mkdir(exist_ok=True)
            
            final_video.write_videofile(
                str(video_path),
                fps=24,
                codec='libx264',
                audio_codec='aac'
            )
            
            # Limpar clips
            for clip in clips:
                clip.close()
            final_video.close()
            
            return video_path
            
        except Exception as e:
            logger.error(f"Erro ao criar vídeo da apresentação: {e}")
            return await self._create_placeholder_video(file_upload)
    
    async def _create_placeholder_video(self, file_upload: FileUpload) -> Path:
        """Cria vídeo placeholder em caso de erro"""
        try:
            # Criar vídeo simples de 10 segundos
            placeholder_image = await self._create_text_slide(
                "Processamento",
                f"O arquivo {file_upload.filename} está sendo processado.\nTente novamente mais tarde."
            )
            
            video_path = Path(settings.UPLOAD_DIR) / "videos" / f"{file_upload.id}_placeholder.mp4"
            video_path.parent.mkdir(exist_ok=True)
            
            # Criar clip de 10 segundos
            image_clip = ImageClip(placeholder_image, duration=10)
            image_clip.write_videofile(
                str(video_path),
                fps=24,
                codec='libx264'
            )
            
            image_clip.close()
            return video_path
            
        except Exception as e:
            logger.error(f"Erro ao criar vídeo placeholder: {e}")
            # Retornar caminho vazio se tudo falhar
            return Path("placeholder.mp4")
    
    async def _get_video_duration(self, video_path: Path) -> float:
        """Obtém duração do vídeo"""
        try:
            clip = VideoFileClip(str(video_path))
            duration = clip.duration
            clip.close()
            return duration
        except Exception as e:
            logger.error(f"Erro ao obter duração do vídeo: {e}")
            return 0.0
    
    async def cleanup_temp_files(self):
        """Limpa arquivos temporários"""
        try:
            import shutil
            if self.temp_dir.exists():
                shutil.rmtree(self.temp_dir)
                self.temp_dir.mkdir(exist_ok=True)
            logger.info("Arquivos temporários limpos")
            
            # Também limpar arquivos do TTS service
            await tts_service.cleanup_temp_files()
            
        except Exception as e:
            logger.error(f"Erro ao limpar arquivos temporários: {e}")


# Instância global do processador
file_processor = FileProcessor()


async def process_file_async(file_id: str):
    """Função para processar arquivo de forma assíncrona"""
    try:
        await file_processor.process_file(file_id)
    except Exception as e:
        logger.error(f"Erro no processamento assíncrono do arquivo {file_id}: {e}")


def schedule_file_processing(file_id: str):
    """Agenda o processamento de um arquivo"""
    asyncio.create_task(process_file_async(file_id)) 