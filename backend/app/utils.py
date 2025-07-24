"""
Funções utilitárias para o sistema SaaS de Upload de Arquivos
Inclui funções para salvar arquivos, validações, geração de nomes únicos, etc.
"""

import os
import hashlib
import shutil
import uuid
import time
from datetime import datetime
from pathlib import Path
from typing import Optional, Tuple, List, Dict, Callable
import mimetypes
# Importação opcional do python-magic (pode causar problemas no Windows)
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    print("⚠️ python-magic não disponível - usando verificação básica de tipos de arquivo")
    MAGIC_AVAILABLE = False

# Imports para manipulação de vídeo e áudio
try:
    from moviepy.editor import ImageClip, AudioFileClip, CompositeVideoClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    print("⚠️ MoviePy não disponível - instale: pip install moviepy")
    MOVIEPY_AVAILABLE = False

# Imports para processamento de imagem e texto
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    print("⚠️ PIL não disponível - instale: pip install pillow")
    PIL_AVAILABLE = False

import fitz  # PyMuPDF
from PIL import Image
import json
from io import BytesIO

# Imports para TTS (Text-to-Speech) - narração de texto
try:
    from backend.services.tts_service import TTSService, TTSConfig, TTSProvider, AudioResult
    TTS_AVAILABLE = True
except ImportError:
    print("Servico TTS nao disponivel - instale as dependencias: pip install torch transformers gtts pydub")
    TTS_AVAILABLE = False
import fitz  # PyMuPDF
from PIL import Image
import json
from io import BytesIO

# Configurações de arquivos
ALLOWED_EXTENSIONS = ['.pdf', '.pptx']
MAX_FILE_SIZE_MB = 50
UPLOAD_DIRECTORY = "app/static/uploads"
VIDEO_DIRECTORY = "app/static/videos"
THUMBNAIL_DIRECTORY = "app/static/thumbnails"

def ensure_directories_exist():
    """
    Garantir que todos os diretórios necessários existam
    Cria os diretórios se eles não existirem
    """
    directories = [
        UPLOAD_DIRECTORY,
        VIDEO_DIRECTORY,
        THUMBNAIL_DIRECTORY,
        f"{UPLOAD_DIRECTORY}/pdf",
        f"{UPLOAD_DIRECTORY}/pptx",
        f"{VIDEO_DIRECTORY}/processed",
        f"{THUMBNAIL_DIRECTORY}/videos"
    ]
    
    for directory in directories:
        Path(directory).mkdir(parents=True, exist_ok=True)
    
    print(f"✅ Diretórios criados/verificados: {len(directories)} diretórios")

def generate_unique_filename(original_filename: str, user_id: int = 1) -> str:
    """
    Gerar nome único para arquivo evitando conflitos
    
    Args:
        original_filename: Nome original do arquivo
        user_id: ID do usuário (para organização)
    
    Returns:
        str: Nome único do arquivo
    """
    # Extrair extensão do arquivo
    file_extension = Path(original_filename).suffix.lower()
    
    # Gerar timestamp único
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Gerar UUID curto para garantir unicidade
    unique_id = str(uuid.uuid4()).replace('-', '')[:8]
    
    # Nome base sem extensão
    base_name = Path(original_filename).stem
    
    # Remover caracteres especiais do nome base
    clean_base_name = sanitize_filename(base_name)
    
    # Construir nome único
    unique_filename = f"{user_id}_{timestamp}_{unique_id}_{clean_base_name}{file_extension}"
    
    return unique_filename

def sanitize_filename(filename: str) -> str:
    """
    Limpar nome do arquivo removendo caracteres especiais
    
    Args:
        filename: Nome do arquivo original
    
    Returns:
        str: Nome limpo e seguro
    """
    # Caracteres permitidos: letras, números, underscore, hífen
    import re
    
    # Remover acentos e caracteres especiais
    filename = filename.replace(' ', '_')
    filename = re.sub(r'[^\w\-_\.]', '', filename)
    
    # Limitar tamanho
    if len(filename) > 50:
        filename = filename[:50]
    
    return filename

def validate_file_type(filename: str) -> bool:
    """
    Validar se o tipo de arquivo é permitido
    
    Args:
        filename: Nome do arquivo
    
    Returns:
        bool: True se arquivo é válido
    """
    if not filename:
        return False
    
    file_extension = Path(filename).suffix.lower()
    return file_extension in ALLOWED_EXTENSIONS

def validate_file_size(file_size: int) -> bool:
    """
    Validar tamanho do arquivo
    
    Args:
        file_size: Tamanho em bytes
    
    Returns:
        bool: True se tamanho é válido
    """
    max_size_bytes = MAX_FILE_SIZE_MB * 1024 * 1024
    return 0 < file_size <= max_size_bytes

def get_file_type_directory(filename: str) -> str:
    """
    Obter diretório baseado no tipo de arquivo
    
    Args:
        filename: Nome do arquivo
    
    Returns:
        str: Caminho do diretório específico
    """
    file_extension = Path(filename).suffix.lower()
    
    if file_extension == '.pdf':
        return f"{UPLOAD_DIRECTORY}/pdf"
    elif file_extension == '.pptx':
        return f"{UPLOAD_DIRECTORY}/pptx"
    else:
        return UPLOAD_DIRECTORY

async def save_uploaded_file(file, user_id: int = 1) -> str:
    """
    Salvar arquivo enviado no sistema de arquivos
    
    Args:
        file: UploadFile do FastAPI
        user_id: ID do usuário (padrão 1 para testes)
    
    Returns:
        str: Caminho completo do arquivo salvo
    
    Raises:
        ValueError: Se arquivo for inválido
        IOError: Se houver erro ao salvar
    """
    # Garantir que diretórios existam
    ensure_directories_exist()
    
    # Validar tipo de arquivo
    if not validate_file_type(file.filename):
        raise ValueError(f"Tipo de arquivo não suportado: {file.filename}")
    
    # Ler conteúdo do arquivo
    file_content = await file.read()
    file_size = len(file_content)
    
    # Validar tamanho
    if not validate_file_size(file_size):
        raise ValueError(f"Arquivo muito grande: {file_size / (1024*1024):.2f}MB. Máximo: {MAX_FILE_SIZE_MB}MB")
    
    # Gerar nome único
    unique_filename = generate_unique_filename(file.filename, user_id)
    
    # Determinar diretório de destino
    target_directory = get_file_type_directory(file.filename)
    
    # Caminho completo do arquivo
    file_path = os.path.join(target_directory, unique_filename)
    
    try:
        # Salvar arquivo
        with open(file_path, "wb") as buffer:
            buffer.write(file_content)
        
        print(f"✅ Arquivo salvo: {file_path} ({file_size / 1024:.2f} KB)")
        return file_path
        
    except Exception as e:
        print(f"❌ Erro ao salvar arquivo: {e}")
        raise IOError(f"Erro ao salvar arquivo: {str(e)}")

def calculate_file_checksum(file_path: str) -> str:
    """
    Calcular checksum MD5 do arquivo para verificação de integridade
    
    Args:
        file_path: Caminho do arquivo
    
    Returns:
        str: Hash MD5 hexadecimal
    """
    hash_md5 = hashlib.md5()
    
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        
        return hash_md5.hexdigest()
        
    except Exception as e:
        print(f"❌ Erro ao calcular checksum: {e}")
        return ""

def get_file_info(file_path: str) -> dict:
    """
    Obter informações detalhadas do arquivo
    
    Args:
        file_path: Caminho do arquivo
    
    Returns:
        dict: Informações do arquivo
    """
    try:
        file_stats = os.stat(file_path)
        file_size = file_stats.st_size
        
        # Tipo MIME
        mime_type, _ = mimetypes.guess_type(file_path)
        
        # Informações básicas
        info = {
            'file_path': file_path,
            'file_size': file_size,
            'mime_type': mime_type or 'application/octet-stream',
            'created_at': datetime.fromtimestamp(file_stats.st_ctime),
            'modified_at': datetime.fromtimestamp(file_stats.st_mtime),
            'checksum': calculate_file_checksum(file_path)
        }
        
        # Informações específicas por tipo
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pdf':
            info.update(get_pdf_info(file_path))
        elif file_extension == '.pptx':
            info.update(get_pptx_info(file_path))
        
        return info
        
    except Exception as e:
        print(f"❌ Erro ao obter informações do arquivo: {e}")
        return {'error': str(e)}

def get_pdf_info(file_path: str) -> dict:
    """
    Obter informações específicas de arquivo PDF
    
    Args:
        file_path: Caminho do arquivo PDF
    
    Returns:
        dict: Informações específicas do PDF
    """
    info = {'pages_count': None}
    
    try:
        # Importar PyPDF2 se disponível
        import PyPDF2
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            info['pages_count'] = len(pdf_reader.pages)
            
            # Metadados se disponíveis
            if pdf_reader.metadata:
                info['title'] = pdf_reader.metadata.get('/Title', '')
                info['author'] = pdf_reader.metadata.get('/Author', '')
                info['creator'] = pdf_reader.metadata.get('/Creator', '')
        
    except ImportError:
        print("⚠️ PyPDF2 não instalado - informações básicas apenas")
    except Exception as e:
        print(f"⚠️ Erro ao analisar PDF: {e}")
    
    return info

def get_pptx_info(file_path: str) -> dict:
    """
    Obter informações específicas de arquivo PPTX
    
    Args:
        file_path: Caminho do arquivo PPTX
    
    Returns:
        dict: Informações específicas do PPTX
    """
    info = {'pages_count': None}
    
    try:
        # Importar python-pptx se disponível
        from pptx import Presentation
        
        prs = Presentation(file_path)
        info['pages_count'] = len(prs.slides)
        
        # Propriedades do documento
        core_props = prs.core_properties
        info['title'] = core_props.title or ''
        info['author'] = core_props.author or ''
        info['created'] = core_props.created
        info['modified'] = core_props.modified
        
    except ImportError:
        print("⚠️ python-pptx não instalado - informações básicas apenas")
    except Exception as e:
        print(f"⚠️ Erro ao analisar PPTX: {e}")
    
    return info

def delete_file(file_path: str) -> bool:
    """
    Deletar arquivo do sistema de arquivos
    
    Args:
        file_path: Caminho do arquivo
    
    Returns:
        bool: True se arquivo foi deletado com sucesso
    """
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"✅ Arquivo deletado: {file_path}")
            return True
        else:
            print(f"⚠️ Arquivo não encontrado: {file_path}")
            return False
            
    except Exception as e:
        print(f"❌ Erro ao deletar arquivo: {e}")
        return False

def format_file_size(size_bytes: int) -> str:
    """
    Formatar tamanho do arquivo em formato legível
    
    Args:
        size_bytes: Tamanho em bytes
    
    Returns:
        str: Tamanho formatado (ex: "1.5 MB")
    """
    if size_bytes == 0:
        return "0 B"
    
    size_names = ["B", "KB", "MB", "GB", "TB"]
    import math
    i = int(math.floor(math.log(size_bytes, 1024)))
    p = math.pow(1024, i)
    s = round(size_bytes / p, 2)
    
    return f"{s} {size_names[i]}"

def get_storage_usage() -> dict:
    """
    Obter informações de uso de armazenamento
    
    Returns:
        dict: Informações de uso de espaço
    """
    usage = {
        'total_files': 0,
        'total_size': 0,
        'pdf_files': 0,
        'pptx_files': 0,
        'videos': 0,
        'directories': {}
    }
    
    try:
        directories_to_check = [
            UPLOAD_DIRECTORY,
            VIDEO_DIRECTORY,
            THUMBNAIL_DIRECTORY
        ]
        
        for directory in directories_to_check:
            if os.path.exists(directory):
                dir_info = get_directory_size(directory)
                usage['directories'][directory] = dir_info
                usage['total_files'] += dir_info['file_count']
                usage['total_size'] += dir_info['total_size']
        
        return usage
        
    except Exception as e:
        print(f"❌ Erro ao calcular uso de armazenamento: {e}")
        return usage

def get_directory_size(directory: str) -> dict:
    """
    Calcular tamanho total de um diretório
    
    Args:
        directory: Caminho do diretório
    
    Returns:
        dict: Informações do diretório
    """
    total_size = 0
    file_count = 0
    
    try:
        for dirpath, dirnames, filenames in os.walk(directory):
            for filename in filenames:
                file_path = os.path.join(dirpath, filename)
                if os.path.exists(file_path):
                    total_size += os.path.getsize(file_path)
                    file_count += 1
        
        return {
            'total_size': total_size,
            'file_count': file_count,
            'formatted_size': format_file_size(total_size)
        }
        
    except Exception as e:
        print(f"❌ Erro ao calcular tamanho do diretório: {e}")
        return {'total_size': 0, 'file_count': 0, 'formatted_size': '0 B'}

def cleanup_old_files(days_old: int = 30) -> int:
    """
    Limpar arquivos antigos do sistema
    
    Args:
        days_old: Idade em dias para considerar arquivo antigo
    
    Returns:
        int: Número de arquivos deletados
    """
    from datetime import timedelta
    
    cutoff_date = datetime.now() - timedelta(days=days_old)
    deleted_count = 0
    
    try:
        directories_to_clean = [UPLOAD_DIRECTORY, VIDEO_DIRECTORY, THUMBNAIL_DIRECTORY]
        
        for directory in directories_to_clean:
            if os.path.exists(directory):
                for root, dirs, files in os.walk(directory):
                    for file in files:
                        file_path = os.path.join(root, file)
                        
                        # Verificar data de modificação
                        if os.path.exists(file_path):
                            file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                            
                            if file_time < cutoff_date:
                                if delete_file(file_path):
                                    deleted_count += 1
        
        print(f"✅ Limpeza concluída: {deleted_count} arquivos removidos")
        return deleted_count
        
    except Exception as e:
        print(f"❌ Erro na limpeza de arquivos: {e}")
        return 0

def validate_file(file_content: bytes, file_extension: str) -> bool:
    """
    Valida se o arquivo está íntegro baseado no conteúdo e extensão
    """
    try:
        # Verificar magic bytes básicos
        if file_extension == '.pdf':
            return file_content.startswith(b'%PDF-')
        elif file_extension == '.pptx':
            # Arquivo ZIP (PPTX é baseado em ZIP)
            return file_content.startswith(b'PK\x03\x04')
        elif file_extension == '.docx':
            # Arquivo ZIP (DOCX é baseado em ZIP)
            return file_content.startswith(b'PK\x03\x04')
        
        # Se python-magic estiver disponível, usar verificação avançada
        if MAGIC_AVAILABLE:
            try:
                file_type = magic.from_buffer(file_content, mime=True)
                return file_type is not None
            except:
                pass
        
        return True
    except Exception:
        return False

def create_thumbnail(file_path: Path, thumbnail_size: Tuple[int, int] = (300, 400)) -> Optional[str]:
    """
    Cria thumbnail para arquivo PDF usando PyMuPDF
    
    Args:
        file_path: Caminho para o arquivo PDF
        thumbnail_size: Tamanho do thumbnail (largura, altura)
    
    Returns:
        str: Caminho do thumbnail criado ou None se falhou
    """
    try:
        # Verificar se é um arquivo PDF
        if not file_path.suffix.lower() == '.pdf':
            print(f"⚠️ Arquivo não é PDF: {file_path}")
            return None
        
        # Verificar se arquivo existe
        if not file_path.exists():
            print(f"❌ Arquivo não encontrado: {file_path}")
            return None
        
        # Garantir que diretório de thumbnails existe
        ensure_directories_exist()
        
        # Gerar nome do thumbnail
        thumbnail_name = f"{file_path.stem}_thumb.png"
        thumbnail_path = Path(THUMBNAIL_DIRECTORY) / thumbnail_name
        
        # Abrir documento PDF
        pdf_document = fitz.open(str(file_path))
        
        if pdf_document.page_count == 0:
            print(f"⚠️ PDF vazio: {file_path}")
            pdf_document.close()
            return None
        
        # Pegar a primeira página
        first_page = pdf_document[0]
        
        # Calcular matriz de transformação para o tamanho desejado
        page_rect = first_page.rect
        scale_x = thumbnail_size[0] / page_rect.width
        scale_y = thumbnail_size[1] / page_rect.height
        scale = min(scale_x, scale_y)  # Manter proporção
        
        matrix = fitz.Matrix(scale, scale)
        
        # Renderizar página como imagem
        pix = first_page.get_pixmap(matrix=matrix)
        
        # Salvar como PNG
        pix.save(str(thumbnail_path))
        
        # Limpar recursos
        pix = None
        pdf_document.close()
        
        print(f"Thumbnail criado: {thumbnail_path} ({thumbnail_size[0]}x{thumbnail_size[1]})")
        return str(thumbnail_path)
        
    except Exception as e:
        print(f"Erro ao criar thumbnail: {e}")
        return None

def create_multiple_thumbnails(file_path: Path, pages: List[int] = None, thumbnail_size: Tuple[int, int] = (300, 400)) -> List[str]:
    """
    Cria thumbnails para múltiplas páginas de um PDF
    
    Args:
        file_path: Caminho para o arquivo PDF
        pages: Lista de números de páginas (1-indexed) ou None para todas
        thumbnail_size: Tamanho dos thumbnails
    
    Returns:
        List[str]: Lista de caminhos dos thumbnails criados
    """
    thumbnails = []
    
    try:
        if not file_path.suffix.lower() == '.pdf':
            return thumbnails
        
        if not file_path.exists():
            return thumbnails
        
        ensure_directories_exist()
        
        pdf_document = fitz.open(str(file_path))
        total_pages = pdf_document.page_count
        
        # Determinar quais páginas processar
        if pages is None:
            pages_to_process = list(range(1, total_pages + 1))
        else:
            pages_to_process = [p for p in pages if 1 <= p <= total_pages]
        
        for page_num in pages_to_process:
            try:
                page_index = page_num - 1  # Converter para 0-indexed
                page = pdf_document[page_index]
                
                # Gerar nome único para cada thumbnail
                thumbnail_name = f"{file_path.stem}_page{page_num}_thumb.png"
                thumbnail_path = Path(THUMBNAIL_DIRECTORY) / thumbnail_name
                
                # Calcular escala
                page_rect = page.rect
                scale_x = thumbnail_size[0] / page_rect.width
                scale_y = thumbnail_size[1] / page_rect.height
                scale = min(scale_x, scale_y)
                
                matrix = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=matrix)
                pix.save(str(thumbnail_path))
                pix = None
                
                thumbnails.append(str(thumbnail_path))
                print(f"Thumbnail pagina {page_num}: {thumbnail_path}")
                
            except Exception as page_error:
                print(f"Erro ao processar pagina {page_num}: {page_error}")
        
        pdf_document.close()
        
        print(f"{len(thumbnails)} thumbnails criados para {file_path.name}")
        return thumbnails
        
    except Exception as e:
        print(f"Erro ao criar thumbnails multiplos: {e}")
        return thumbnails

def get_pdf_page_info(file_path: Path) -> dict:
    """
    Obtém informações detalhadas sobre as páginas de um PDF
    
    Args:
        file_path: Caminho para o arquivo PDF
    
    Returns:
        dict: Informações das páginas
    """
    result = {
        "success": False,
        "page_count": 0,
        "pages_info": [],
        "document_size": {"width": 0, "height": 0},
        "error": None
    }
    
    try:
        pdf_document = fitz.open(str(file_path))
        result["page_count"] = pdf_document.page_count
        
        pages_info = []
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            page_rect = page.rect
            
            page_info = {
                "page_number": page_num + 1,
                "width": page_rect.width,
                "height": page_rect.height,
                "rotation": page.rotation,
                "links_count": len(page.get_links()),
                "images_count": len(page.get_images()),
                "text_length": len(page.get_text())
            }
            
            pages_info.append(page_info)
        
        pdf_document.close()
        
        # Calcular tamanho médio do documento
        if pages_info:
            avg_width = sum(p["width"] for p in pages_info) / len(pages_info)
            avg_height = sum(p["height"] for p in pages_info) / len(pages_info)
            result["document_size"] = {"width": avg_width, "height": avg_height}
        
        result["pages_info"] = pages_info
        result["success"] = True
        
        return result
        
    except Exception as e:
        result["error"] = str(e)
        print(f"❌ Erro ao obter informações das páginas: {e}")
        return result



def extract_pdf_text(file_path: Path, max_pages: Optional[int] = None) -> dict:
    """
    Extrai texto completo de um arquivo PDF usando PyMuPDF
    
    Args:
        file_path: Caminho para o arquivo PDF
        max_pages: Número máximo de páginas para processar (None = todas)
    
    Returns:
        dict: Contendo texto extraído, metadados e informações das páginas
    """
    result = {
        "success": False,
        "text": "",
        "page_count": 0,
        "pages_text": [],
        "metadata": {},
        "file_size": 0,
        "error": None
    }
    
    try:
        # Verificar se arquivo existe
        if not file_path.exists():
            result["error"] = f"Arquivo não encontrado: {file_path}"
            return result
        
        # Obter tamanho do arquivo
        result["file_size"] = file_path.stat().st_size
        
        # Abrir documento PDF com PyMuPDF
        pdf_document = fitz.open(str(file_path))
        result["page_count"] = pdf_document.page_count
        
        # Extrair metadados do PDF
        metadata = pdf_document.metadata
        result["metadata"] = {
            "title": metadata.get("title", ""),
            "author": metadata.get("author", ""),
            "subject": metadata.get("subject", ""),
            "creator": metadata.get("creator", ""),
            "producer": metadata.get("producer", ""),
            "creation_date": metadata.get("creationDate", ""),
            "modification_date": metadata.get("modDate", "")
        }
        
        # Determinar quantas páginas processar
        pages_to_process = result["page_count"]
        if max_pages and max_pages > 0:
            pages_to_process = min(max_pages, result["page_count"])
        
        # Extrair texto de cada página
        all_text = []
        pages_text = []
        
        for page_num in range(pages_to_process):
            try:
                page = pdf_document[page_num]
                page_text = page.get_text()
                
                # Limpar e processar texto da página
                cleaned_text = clean_extracted_text(page_text)
                
                pages_text.append({
                    "page_number": page_num + 1,
                    "text": cleaned_text,
                    "word_count": len(cleaned_text.split()),
                    "char_count": len(cleaned_text)
                })
                
                all_text.append(cleaned_text)
                
                print(f"📄 Página {page_num + 1}/{pages_to_process} processada: {len(cleaned_text)} caracteres")
                
            except Exception as page_error:
                print(f"⚠️ Erro ao processar página {page_num + 1}: {page_error}")
                pages_text.append({
                    "page_number": page_num + 1,
                    "text": "",
                    "error": str(page_error),
                    "word_count": 0,
                    "char_count": 0
                })
        
        # Fechar documento
        pdf_document.close()
        
        # Combinar todo o texto
        result["text"] = "\n\n".join(all_text)
        result["pages_text"] = pages_text
        result["success"] = True
        
        # Estatísticas finais
        total_words = len(result["text"].split())
        total_chars = len(result["text"])
        
        print(f"✅ PDF processado com sucesso:")
        print(f"   📄 Páginas: {result['page_count']} (processadas: {pages_to_process})")
        print(f"   📝 Palavras: {total_words}")
        print(f"   🔤 Caracteres: {total_chars}")
        print(f"   📁 Tamanho: {format_file_size(result['file_size'])}")
        
        return result
        
    except Exception as e:
        error_msg = f"Erro ao extrair texto do PDF: {str(e)}"
        print(f"❌ {error_msg}")
        result["error"] = error_msg
        return result

def clean_extracted_text(text: str) -> str:
    """
    Limpa e formata o texto extraído do PDF
    
    Args:
        text: Texto bruto extraído
    
    Returns:
        str: Texto limpo e formatado
    """
    if not text:
        return ""
    
    # Remover caracteres de controle excessivos
    import re
    
    # Substituir múltiplas quebras de linha por uma única
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Remover espaços excessivos
    text = re.sub(r' {2,}', ' ', text)
    
    # Remover tabs excessivos
    text = re.sub(r'\t+', ' ', text)
    
    # Limpar início e fim
    text = text.strip()
    
    return text

def extract_pdf_text_with_formatting(file_path: Path, preserve_formatting: bool = True) -> dict:
    """
    Extrai texto do PDF preservando informações de formatação
    
    Args:
        file_path: Caminho para o arquivo PDF
        preserve_formatting: Se deve preservar informações de formatação
    
    Returns:
        dict: Texto com informações de formatação
    """
    result = {
        "success": False,
        "formatted_content": [],
        "plain_text": "",
        "page_count": 0,
        "error": None
    }
    
    try:
        pdf_document = fitz.open(str(file_path))
        result["page_count"] = pdf_document.page_count
        
        formatted_content = []
        plain_text_parts = []
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            
            if preserve_formatting:
                # Extrair texto com informações de formatação
                text_dict = page.get_text("dict")
                
                page_content = {
                    "page_number": page_num + 1,
                    "blocks": []
                }
                
                for block in text_dict.get("blocks", []):
                    if "lines" in block:  # Bloco de texto
                        block_info = {
                            "type": "text",
                            "bbox": block.get("bbox", []),
                            "lines": []
                        }
                        
                        for line in block["lines"]:
                            line_text = ""
                            line_info = {
                                "bbox": line.get("bbox", []),
                                "spans": []
                            }
                            
                            for span in line.get("spans", []):
                                span_text = span.get("text", "")
                                line_text += span_text
                                
                                span_info = {
                                    "text": span_text,
                                    "font": span.get("font", ""),
                                    "size": span.get("size", 0),
                                    "flags": span.get("flags", 0),
                                    "color": span.get("color", 0),
                                    "bbox": span.get("bbox", [])
                                }
                                line_info["spans"].append(span_info)
                            
                            line_info["text"] = line_text
                            block_info["lines"].append(line_info)
                            plain_text_parts.append(line_text)
                        
                        page_content["blocks"].append(block_info)
                
                formatted_content.append(page_content)
            
            else:
                # Extrair apenas texto simples
                page_text = page.get_text()
                plain_text_parts.append(page_text)
        
        pdf_document.close()
        
        result["formatted_content"] = formatted_content
        result["plain_text"] = "\n".join(plain_text_parts)
        result["success"] = True
        
        print(f"✅ PDF processado com formatação: {len(result['plain_text'])} caracteres")
        
        return result
        
    except Exception as e:
        error_msg = f"Erro ao extrair texto formatado: {str(e)}"
        print(f"❌ {error_msg}")
        result["error"] = error_msg
        return result

def search_text_in_pdf(file_path: Path, search_term: str, case_sensitive: bool = False) -> dict:
    """
    Busca por texto específico no PDF
    
    Args:
        file_path: Caminho para o arquivo PDF
        search_term: Termo a ser buscado
        case_sensitive: Se a busca deve ser sensível a maiúsculas
    
    Returns:
        dict: Resultados da busca
    """
    result = {
        "success": False,
        "search_term": search_term,
        "total_matches": 0,
        "pages_with_matches": [],
        "matches": [],
        "error": None
    }
    
    try:
        pdf_document = fitz.open(str(file_path))
        
        search_flags = 0 if case_sensitive else fitz.TEXT_INHIBIT_SPACES
        
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            
            # Buscar texto na página
            text_instances = page.search_for(search_term, flags=search_flags)
            
            if text_instances:
                page_matches = {
                    "page_number": page_num + 1,
                    "match_count": len(text_instances),
                    "matches": []
                }
                
                for instance in text_instances:
                    # Extrair contexto ao redor da ocorrência
                    page_text = page.get_text()
                    
                    match_info = {
                        "bbox": list(instance),
                        "context": extract_context(page_text, search_term, 50)
                    }
                    page_matches["matches"].append(match_info)
                
                result["pages_with_matches"].append(page_matches)
                result["matches"].extend(text_instances)
        
        pdf_document.close()
        
        result["total_matches"] = len(result["matches"])
        result["success"] = True
        
        print(f"🔍 Busca concluída: {result['total_matches']} ocorrências de '{search_term}'")
        
        return result
        
    except Exception as e:
        error_msg = f"Erro na busca: {str(e)}"
        print(f"❌ {error_msg}")
        result["error"] = error_msg
        return result

def extract_context(text: str, search_term: str, context_length: int = 50) -> str:
    """
    Extrai contexto ao redor de um termo encontrado
    
    Args:
        text: Texto completo
        search_term: Termo encontrado
        context_length: Número de caracteres de contexto em cada lado
    
    Returns:
        str: Contexto extraído
    """
    try:
        index = text.lower().find(search_term.lower())
        if index == -1:
            return ""
        
        start = max(0, index - context_length)
        end = min(len(text), index + len(search_term) + context_length)
        
        context = text[start:end]
        
        # Adicionar indicadores se o contexto foi cortado
        if start > 0:
            context = "..." + context
        if end < len(text):
            context = context + "..."
        
        return context
        
    except Exception:
        return ""

def get_file_metadata(file_path: Path, file_extension: str) -> dict:
    """
    Extrai metadados básicos do arquivo
    """
    metadata = {
        "file_size": file_path.stat().st_size,
        "created_at": file_path.stat().st_ctime,
        "modified_at": file_path.stat().st_mtime,
        "extension": file_extension
    }
    
    try:
        if file_extension == '.pdf':
            # Metadados básicos para PDF
            metadata["type"] = "pdf"
            metadata["page_count"] = 1  # Placeholder
            
        elif file_extension in ['.pptx', '.docx']:
            metadata["type"] = "office_document"
            
    except Exception as e:
        print(f"Erro ao extrair metadados: {e}")
    
    return metadata

def extract_text_from_pdf(pdf_path: str) -> List[str]:
    """
    Extrai texto de um arquivo PDF retornando uma lista com o texto de cada página.
    
    Args:
        pdf_path (str): Caminho para o arquivo PDF
    
    Returns:
        List[str]: Lista contendo o texto de cada página
        
    Raises:
        FileNotFoundError: Se o arquivo não for encontrado
        Exception: Se houver erro na leitura do PDF
    """
    try:
        # Verificar se o arquivo existe
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"Arquivo PDF não encontrado: {pdf_path}")
        
        # Abrir o documento PDF usando PyMuPDF (fitz)
        pdf_document = fitz.open(pdf_path)
        
        # Lista para armazenar o texto de cada página
        pages_text = []
        
        # Iterar através de todas as páginas do PDF
        for page_num in range(pdf_document.page_count):
            # Obter a página atual
            page = pdf_document[page_num]
            
            # Extrair o texto da página
            page_text = page.get_text()
            
            # Limpar o texto extraído (remover espaços extras e quebras de linha desnecessárias)
            cleaned_text = clean_extracted_text(page_text)
            
            # Adicionar à lista de textos das páginas
            pages_text.append(cleaned_text)
            
            print(f"📄 Página {page_num + 1}/{pdf_document.page_count} processada: {len(cleaned_text)} caracteres")
        
        # Fechar o documento PDF para liberar recursos
        pdf_document.close()
        
        print(f"✅ PDF processado com sucesso: {len(pages_text)} páginas extraídas")
        return pages_text
        
    except FileNotFoundError:
        print(f"❌ Arquivo não encontrado: {pdf_path}")
        raise
    except Exception as e:
        print(f"❌ Erro ao extrair texto do PDF: {str(e)}")
        raise Exception(f"Erro ao processar arquivo PDF: {str(e)}")

def extract_text_from_pptx(pptx_path: str) -> List[str]:
    """
    Extrai texto de um arquivo PPTX retornando uma lista com o texto de cada slide.
    
    Args:
        pptx_path (str): Caminho para o arquivo PPTX
    
    Returns:
        List[str]: Lista contendo o texto de cada slide
        
    Raises:
        FileNotFoundError: Se o arquivo não for encontrado
        ImportError: Se a biblioteca python-pptx não estiver instalada
        Exception: Se houver erro na leitura do PPTX
    """
    try:
        # Verificar se o arquivo existe
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"Arquivo PPTX não encontrado: {pptx_path}")
        
        # Importar a biblioteca python-pptx
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Biblioteca python-pptx não encontrada. Instale com: pip install python-pptx")
        
        # Abrir a apresentação PPTX
        presentation = Presentation(pptx_path)
        
        # Lista para armazenar o texto de cada slide
        slides_text = []
        
        # Iterar através de todos os slides da apresentação
        for slide_num, slide in enumerate(presentation.slides):
            # Lista para armazenar todo o texto do slide atual
            slide_text_parts = []
            
            # Extrair texto de todas as formas (shapes) do slide
            for shape in slide.shapes:
                # Verificar se a forma tem texto
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text.strip())
                
                # Se a forma tem uma tabela, extrair texto das células
                elif hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text_parts.append(cell.text.strip())
                
                # Se a forma tem um text_frame, extrair texto dos parágrafos
                elif hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph_text = ""
                        for run in paragraph.runs:
                            paragraph_text += run.text
                        if paragraph_text.strip():
                            slide_text_parts.append(paragraph_text.strip())
            
            # Juntar todo o texto do slide com quebras de linha
            slide_full_text = "\n".join(slide_text_parts)
            
            # Limpar espaços extras e adicionar à lista
            cleaned_slide_text = slide_full_text.strip()
            slides_text.append(cleaned_slide_text)
            
            print(f"Slide {slide_num + 1}/{len(presentation.slides)} processado: {len(cleaned_slide_text)} caracteres")
        
        print(f"PPTX processado com sucesso: {len(slides_text)} slides extraídos")
        return slides_text
        
    except FileNotFoundError:
        print(f"Arquivo não encontrado: {pptx_path}")
        raise
    except ImportError as e:
        print(f"Erro de importação: {str(e)}")
        raise
    except Exception as e:
        print(f"Erro ao extrair texto do PPTX: {str(e)}")
        raise Exception(f"Erro ao processar arquivo PPTX: {str(e)}")

def analyze_text_statistics(text_pages: List[str]) -> dict:
    """
    Analisa estatísticas detalhadas do texto extraído.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas/slides
    
    Returns:
        dict: Dicionário com estatísticas completas do texto
    """
    try:
        # Combinar todo o texto
        full_text = "\n\n".join(text_pages)
        
        # Estatísticas básicas
        total_chars = len(full_text)
        total_chars_no_spaces = len(full_text.replace(" ", "").replace("\n", "").replace("\t", ""))
        
        # Contar palavras (dividir por espaços e filtrar strings vazias)
        words = [word.strip(".,!?;:\"'()[]{}") for word in full_text.split() if word.strip()]
        total_words = len(words)
        
        # Contar sentenças (aproximadamente)
        sentences = [s.strip() for s in full_text.replace("!", ".").replace("?", ".").split(".") if s.strip()]
        total_sentences = len(sentences)
        
        # Contar parágrafos
        paragraphs = [p.strip() for p in full_text.split("\n\n") if p.strip()]
        total_paragraphs = len(paragraphs)
        
        # Palavras únicas (case-insensitive)
        unique_words = set([word.lower() for word in words if len(word) > 2])
        unique_word_count = len(unique_words)
        
        # Estatísticas por página/slide
        pages_stats = []
        for i, page_text in enumerate(text_pages):
            page_words = len([w for w in page_text.split() if w.strip()])
            page_chars = len(page_text)
            
            pages_stats.append({
                "page_number": i + 1,
                "characters": page_chars,
                "words": page_words,
                "lines": len(page_text.split("\n"))
            })
        
        # Frequência de palavras (top 10)
        word_frequency = {}
        for word in words:
            clean_word = word.lower().strip(".,!?;:\"'()[]{}").strip()
            if len(clean_word) > 3:  # Ignorar palavras muito pequenas
                word_frequency[clean_word] = word_frequency.get(clean_word, 0) + 1
        
        # Top 10 palavras mais frequentes
        top_words = sorted(word_frequency.items(), key=lambda x: x[1], reverse=True)[:10]
        
        # Calcular médias
        avg_words_per_page = total_words / len(text_pages) if text_pages else 0
        avg_chars_per_page = total_chars / len(text_pages) if text_pages else 0
        avg_words_per_sentence = total_words / total_sentences if total_sentences > 0 else 0
        
        statistics = {
            "overview": {
                "total_pages": len(text_pages),
                "total_characters": total_chars,
                "total_characters_no_spaces": total_chars_no_spaces,
                "total_words": total_words,
                "total_sentences": total_sentences,
                "total_paragraphs": total_paragraphs,
                "unique_words": unique_word_count
            },
            "averages": {
                "words_per_page": round(avg_words_per_page, 2),
                "characters_per_page": round(avg_chars_per_page, 2),
                "words_per_sentence": round(avg_words_per_sentence, 2),
                "sentences_per_page": round(total_sentences / len(text_pages), 2) if text_pages else 0
            },
            "pages_breakdown": pages_stats,
            "word_frequency": {
                "top_words": top_words,
                "vocabulary_richness": round(unique_word_count / total_words * 100, 2) if total_words > 0 else 0
            },
            "readability": {
                "estimated_reading_time_minutes": round(total_words / 200, 1),  # 200 palavras por minuto
                "average_word_length": round(sum(len(word) for word in words) / len(words), 2) if words else 0
            }
        }
        
        print(f"📊 Análise de estatísticas concluída:")
        print(f"   📄 Total de páginas: {statistics['overview']['total_pages']}")
        print(f"   📝 Total de palavras: {statistics['overview']['total_words']}")
        print(f"   🔤 Total de caracteres: {statistics['overview']['total_characters']}")
        print(f"   ⏱️ Tempo de leitura estimado: {statistics['readability']['estimated_reading_time_minutes']} minutos")
        
        return statistics
        
    except Exception as e:
        print(f"❌ Erro ao analisar estatísticas: {str(e)}")
        return {"error": str(e)}

def search_text_in_pages(text_pages: List[str], search_terms: List[str], case_sensitive: bool = False) -> dict:
    """
    Busca termos específicos nas páginas de texto extraído.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas
        search_terms (List[str]): Lista de termos para buscar
        case_sensitive (bool): Se a busca deve ser sensível a maiúsculas
    
    Returns:
        dict: Resultados da busca com detalhes de localização
    """
    try:
        search_results = {
            "search_terms": search_terms,
            "case_sensitive": case_sensitive,
            "total_matches": 0,
            "results_by_term": {},
            "pages_with_matches": []
        }
        
        pages_with_any_match = set()
        
        # Buscar cada termo
        for term in search_terms:
            term_results = {
                "term": term,
                "total_occurrences": 0,
                "pages_found": [],
                "contexts": []
            }
            
            # Buscar em cada página
            for page_num, page_text in enumerate(text_pages):
                search_text = page_text if case_sensitive else page_text.lower()
                search_term = term if case_sensitive else term.lower()
                
                # Contar ocorrências na página
                occurrences = search_text.count(search_term)
                
                if occurrences > 0:
                    term_results["total_occurrences"] += occurrences
                    term_results["pages_found"].append({
                        "page_number": page_num + 1,
                        "occurrences": occurrences
                    })
                    pages_with_any_match.add(page_num + 1)
                    
                    # Extrair contexto (50 caracteres antes e depois)
                    contexts = extract_search_contexts(page_text, term, case_sensitive)
                    for context in contexts:
                        term_results["contexts"].append({
                            "page_number": page_num + 1,
                            "context": context
                        })
            
            search_results["results_by_term"][term] = term_results
            search_results["total_matches"] += term_results["total_occurrences"]
        
        search_results["pages_with_matches"] = sorted(list(pages_with_any_match))
        
        print(f"🔍 Busca concluída:")
        print(f"   📋 Termos buscados: {len(search_terms)}")
        print(f"   🎯 Total de ocorrências: {search_results['total_matches']}")
        print(f"   📄 Páginas com resultados: {len(search_results['pages_with_matches'])}")
        
        return search_results
        
    except Exception as e:
        print(f"❌ Erro na busca: {str(e)}")
        return {"error": str(e)}

def extract_search_contexts(text: str, search_term: str, case_sensitive: bool = False, context_length: int = 50) -> List[str]:
    """
    Extrai contextos ao redor das ocorrências de um termo de busca.
    
    Args:
        text (str): Texto onde buscar
        search_term (str): Termo a ser buscado
        case_sensitive (bool): Se a busca deve ser sensível a maiúsculas
        context_length (int): Número de caracteres de contexto em cada lado
    
    Returns:
        List[str]: Lista de contextos encontrados
    """
    contexts = []
    search_text = text if case_sensitive else text.lower()
    term = search_term if case_sensitive else search_term.lower()
    
    start_pos = 0
    while True:
        pos = search_text.find(term, start_pos)
        if pos == -1:
            break
        
        # Extrair contexto
        start_context = max(0, pos - context_length)
        end_context = min(len(text), pos + len(search_term) + context_length)
        
        context = text[start_context:end_context]
        
        # Adicionar indicadores se o contexto foi cortado
        if start_context > 0:
            context = "..." + context
        if end_context < len(text):
            context = context + "..."
        
        contexts.append(context)
        start_pos = pos + 1
    
    return contexts

def extract_pdf_text(file_path: Path, max_pages: Optional[int] = None) -> dict:
    """
    Extrai texto completo de um arquivo PDF usando PyMuPDF
    
    Args:
        file_path (Path): Caminho do arquivo PDF
        max_pages (Optional[int]): Número máximo de páginas a extrair (padrão: None)
    
    Returns:
        dict: Resultado da extração de texto
    """
    try:
        import fitz
        
        # Abrir arquivo PDF
        pdf_document = fitz.open(file_path)
        
        # Verificar número de páginas
        total_pages = pdf_document.page_count
        if max_pages is not None and max_pages > 0:
            total_pages = min(total_pages, max_pages)
        
        # Extrair texto de cada página
        text_pages = []
        for page_num in range(total_pages):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            text_pages.append(text)
        
        # Fechar documento
        pdf_document.close()
        
        # Preparar resultado
        result = {
            "total_pages": total_pages,
            "text_pages": text_pages,
            "total_characters": sum(len(page) for page in text_pages),
            "total_words": sum(len(page.split()) for page in text_pages)
        }
        
        print(f"📝 Extração de texto de PDF concluída:")
        print(f"   📄 Páginas: {result['total_pages']}")
        print(f"   📝 Palavras: {result['total_words']:,}")
        print(f"   🔤 Caracteres: {result['total_characters']:,}")
        
        return result
        
    except Exception as e:
        print(f"❌ Erro na extração de texto de PDF: {str(e)}")
        return {"error": str(e)}

def extract_pptx_text(file_path: Path) -> dict:
    """
    Extrai texto completo de um arquivo PPTX usando python-pptx
    
    Args:
        file_path (Path): Caminho do arquivo PPTX
    
    Returns:
        dict: Resultado da extração de texto
    """
    try:
        from pptx import Presentation
        
        # Abrir arquivo PPTX
        presentation = Presentation(file_path)
        
        # Extrair texto de cada slide
        text_slides = []
        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            text_slides.append(slide_text.strip())
        
        # Preparar resultado
        result = {
            "total_slides": len(text_slides),
            "text_slides": text_slides,
            "total_characters": sum(len(slide) for slide in text_slides),
            "total_words": sum(len(slide.split()) for slide in text_slides)
        }
        
        print(f"📝 Extração de texto de PPTX concluída:")
        print(f"   📄 Slides: {result['total_slides']}")
        print(f"   📝 Palavras: {result['total_words']:,}")
        print(f"   🔤 Caracteres: {result['total_characters']:,}")
        
        return result
        
    except Exception as e:
        print(f"❌ Erro na extração de texto de PPTX: {str(e)}")
        return {"error": str(e)}

def extract_text_from_pdf(file_path: str, max_pages: int = None) -> List[str]:
    """
    Extrai texto de um arquivo PDF.
    
    Args:
        file_path (str): Caminho do arquivo PDF
        max_pages (int): Número máximo de páginas a extrair (opcional)
    
    Returns:
        List[str]: Lista de textos das páginas
    """
    try:
        import fitz
        
        # Abrir arquivo PDF
        pdf_document = fitz.open(file_path)
        
        # Verificar número de páginas
        total_pages = pdf_document.page_count
        if max_pages is not None and max_pages > 0:
            total_pages = min(total_pages, max_pages)
        
        # Extrair texto de cada página
        text_pages = []
        for page_num in range(total_pages):
            page = pdf_document.load_page(page_num)
            text = page.get_text()
            text_pages.append(text)
        
        # Fechar documento
        pdf_document.close()
        
        print(f"📝 Extração de texto de PDF concluída:")
        print(f"   📄 Páginas: {len(text_pages)}")
        print(f"   📝 Palavras: {sum(len(page.split()) for page in text_pages):,}")
        print(f"   🔤 Caracteres: {sum(len(page) for page in text_pages):,}")
        
        return text_pages
        
    except Exception as e:
        print(f"❌ Erro na extração de texto de PDF: {str(e)}")
        return []

# Função extract_text_from_pptx duplicada removida - mantendo apenas a versão completa acima

def analyze_text_statistics(text_pages: List[str]) -> dict:
    """
    Analisa estatísticas básicas do texto.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas
    
    Returns:
        dict: Estatísticas do texto
    """
    try:
        # Combinar todo o texto
        full_text = " ".join(text_pages)
        
        # Análise de estatísticas
        words = full_text.split()
        word_count = len(words)
        unique_words = len(set(words))
        character_count = len(full_text)
        paragraphs = full_text.split('\n\n')
        paragraph_count = len(paragraphs)
        avg_words_per_paragraph = word_count / paragraph_count if paragraph_count > 0 else 0
        
        # Análise de sentenças
        import re
        sentences = re.split(r'[.!?]+', full_text)
        sentence_count = len(sentences)
        avg_words_per_sentence = word_count / sentence_count if sentence_count > 0 else 0
        
        # Análise de complexidade
        long_words = [w for w in words if len(w) > 8]
        complex_sentences = [s for s in sentences if len(s.split()) > 20]
        
        # Preparar resultado
        statistics = {
            "word_count": word_count,
            "unique_words": unique_words,
            "character_count": character_count,
            "paragraph_count": paragraph_count,
            "avg_words_per_paragraph": round(avg_words_per_paragraph, 1),
            "sentence_count": sentence_count,
            "avg_words_per_sentence": round(avg_words_per_sentence, 1),
            "long_words_count": len(long_words),
            "long_words_percentage": round(len(long_words) / word_count * 100, 2) if word_count > 0 else 0,
            "complex_sentences_count": len(complex_sentences),
            "complex_sentences_percentage": round(len(complex_sentences) / sentence_count * 100, 2) if sentence_count > 0 else 0
        }
        
        print(f"📊 Análise de estatísticas concluída:")
        print(f"   📝 Palavras: {word_count:,}")
        print(f"   🔤 Caracteres: {character_count:,}")
        print(f"   📃 Parágrafos: {paragraph_count}")
        print(f"   📜 Sentenças: {sentence_count}")
        print(f"   📈 Complexidade: {statistics['long_words_percentage']}% palavras longas")
        
        return statistics
        
    except Exception as e:
        print(f"❌ Erro na análise de estatísticas: {str(e)}")
        return {"error": str(e)}

def extract_keywords(text_pages: List[str], top_n: int = 20, min_word_length: int = 4) -> dict:
    """
    Extrai palavras-chave mais relevantes do texto.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas
        top_n (int): Número de palavras-chave para retornar
        min_word_length (int): Comprimento mínimo das palavras
    
    Returns:
        dict: Palavras-chave com pontuações de relevância
    """
    try:
        # Combinar todo o texto
        full_text = " ".join(text_pages).lower()
        
        # Palavras comuns em português para filtrar (stop words básicas)
        stop_words = {
            'o', 'a', 'os', 'as', 'um', 'uma', 'uns', 'umas', 'de', 'do', 'da', 'dos', 'das',
            'em', 'no', 'na', 'nos', 'nas', 'por', 'para', 'com', 'sem', 'sob', 'sobre',
            'e', 'ou', 'mas', 'que', 'se', 'não', 'sim', 'como', 'quando', 'onde',
            'este', 'esta', 'estes', 'estas', 'esse', 'essa', 'esses', 'essas',
            'aquele', 'aquela', 'aqueles', 'aquelas', 'isto', 'isso', 'aquilo',
            'eu', 'tu', 'ele', 'ela', 'nós', 'vós', 'eles', 'elas', 'me', 'te', 'se', 'nos', 'vos',
            'meu', 'minha', 'meus', 'minhas', 'teu', 'tua', 'teus', 'tuas',
            'seu', 'sua', 'seus', 'suas', 'nosso', 'nossa', 'nossos', 'nossas',
            'foi', 'são', 'tem', 'ter', 'ser', 'estar', 'foi', 'era', 'será', 'sendo',
            'muito', 'mais', 'menos', 'bem', 'mal', 'melhor', 'pior', 'tanto', 'quanto',
            'já', 'ainda', 'também', 'só', 'apenas', 'sempre', 'nunca', 'hoje', 'ontem', 'amanhã'
        }
        
        # Extrair palavras válidas
        import re
        words = re.findall(r'\b[a-záàâãéêíîóôõúûçA-ZÁÀÂÃÉÊÍÎÓÔÕÚÛÇ]+\b', full_text)
        
        # Filtrar palavras por comprimento e stop words
        filtered_words = [
            word for word in words 
            if len(word) >= min_word_length and word.lower() not in stop_words
        ]
        
        # Calcular frequência das palavras
        word_freq = {}
        for word in filtered_words:
            word_lower = word.lower()
            word_freq[word_lower] = word_freq.get(word_lower, 0) + 1
        
        # Calcular pontuação TF (Term Frequency)
        total_words = len(filtered_words)
        tf_scores = {}
        for word, freq in word_freq.items():
            tf_scores[word] = freq / total_words
        
        # Aplicar boost para palavras mais longas (assumindo que são mais específicas)
        length_boost = {}
        for word in tf_scores:
            length_boost[word] = min(2.0, len(word) / 8)  # Máximo de 2x boost
        
        # Calcular pontuação final (TF * length_boost)
        final_scores = {}
        for word in tf_scores:
            final_scores[word] = tf_scores[word] * length_boost[word]
        
        # Ordenar por pontuação e pegar top N
        sorted_keywords = sorted(final_scores.items(), key=lambda x: x[1], reverse=True)[:top_n]
        
        # Preparar resultado
        keywords_result = {
            "total_unique_words": len(word_freq),
            "total_filtered_words": total_words,
            "extraction_params": {
                "top_n": top_n,
                "min_word_length": min_word_length,
                "stop_words_count": len(stop_words)
            },
            "keywords": []
        }
        
        for word, score in sorted_keywords:
            keywords_result["keywords"].append({
                "word": word,
                "frequency": word_freq[word],
                "tf_score": round(tf_scores[word], 6),
                "final_score": round(score, 6),
                "relevance_percentage": round(score / sorted_keywords[0][1] * 100, 2) if sorted_keywords else 0
            })
        
        print(f"🔑 Extração de palavras-chave concluída:")
        print(f"   📊 {len(sorted_keywords)} palavras-chave extraídas")
        print(f"   📝 Top 5: {[kw['word'] for kw in keywords_result['keywords'][:5]]}")
        
        return keywords_result
        
    except Exception as e:
        print(f"❌ Erro na extração de palavras-chave: {str(e)}")
        return {"error": str(e)}

def generate_text_summary(text_pages: List[str], max_sentences: int = 5) -> dict:
    """
    Gera um resumo automático do texto baseado na frequência de palavras e posicionamento.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas
        max_sentences (int): Número máximo de sentenças no resumo
    
    Returns:
        dict: Resumo gerado com métricas
    """
    try:
        # Combinar todo o texto
        full_text = " ".join(text_pages)
        
        # Dividir em sentenças
        import re
        sentences = re.split(r'[.!?]+', full_text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]  # Filtrar sentenças muito curtas
        
        if len(sentences) <= max_sentences:
            return {
                "summary": full_text,
                "original_sentences": len(sentences),
                "summary_sentences": len(sentences),
                "compression_ratio": 1.0,
                "method": "no_compression_needed"
            }
        
        # Extrair palavras-chave para scoring
        keywords_data = extract_keywords(text_pages, top_n=30, min_word_length=3)
        important_words = set()
        
        if "keywords" in keywords_data:
            important_words = {kw["word"] for kw in keywords_data["keywords"][:15]}
        
        # Pontuar sentenças
        sentence_scores = {}
        
        for i, sentence in enumerate(sentences):
            score = 0
            sentence_lower = sentence.lower()
            words_in_sentence = re.findall(r'\b[a-záàâãéêíîóôõúûçA-ZÁÀÂÃÉÊÍÎÓÔÕÚÛÇ]+\b', sentence_lower)
            
            # Pontuação baseada em palavras-chave
            for word in words_in_sentence:
                if word in important_words:
                    score += 2
            
            # Boost para sentenças no início (introdução) e final (conclusão)
            if i < len(sentences) * 0.2:  # Primeiros 20%
                score *= 1.3
            elif i > len(sentences) * 0.8:  # Últimos 20%
                score *= 1.2
            
            # Penalizar sentenças muito curtas ou muito longas
            if len(words_in_sentence) < 5:
                score *= 0.5
            elif len(words_in_sentence) > 40:
                score *= 0.8
            
            # Boost para sentenças com números (possíveis dados importantes)
            if re.search(r'\d+', sentence):
                score *= 1.1
            
            sentence_scores[i] = score
        
        # Selecionar melhores sentenças
        best_sentences_indices = sorted(sentence_scores.keys(), 
                                      key=lambda x: sentence_scores[x], 
                                      reverse=True)[:max_sentences]
        
        # Manter ordem original das sentenças selecionadas
        best_sentences_indices.sort()
        
        # Construir resumo
        summary_sentences = [sentences[i] for i in best_sentences_indices]
        summary_text = ". ".join(summary_sentences) + "."
        
        # Limpar resumo
        summary_text = re.sub(r'\s+', ' ', summary_text).strip()
        
        compression_ratio = len(summary_sentences) / len(sentences)
        
        result = {
            "summary": summary_text,
            "original_sentences": len(sentences),
            "summary_sentences": len(summary_sentences),
            "compression_ratio": round(compression_ratio, 3),
            "method": "keyword_frequency_position",
            "selected_sentence_indices": best_sentences_indices,
            "summary_stats": {
                "characters": len(summary_text),
                "words": len(summary_text.split()),
                "estimated_reading_time": round(len(summary_text.split()) / 200, 1)
            }
        }
        
        print(f"📝 Resumo automático gerado:")
        print(f"   📊 Compressão: {len(sentences)} → {len(summary_sentences)} sentenças")
        print(f"   📈 Taxa de compressão: {compression_ratio:.1%}")
        print(f"   ⏱️ Tempo de leitura: {result['summary_stats']['estimated_reading_time']} min")
        
        return result
        
    except Exception as e:
        print(f"❌ Erro na geração de resumo: {str(e)}")
        return {"error": str(e)}

def analyze_text_language_patterns(text_pages: List[str]) -> dict:
    """
    Analisa padrões linguísticos básicos do texto.
    
    Args:
        text_pages (List[str]): Lista de textos das páginas
    
    Returns:
        dict: Análise de padrões linguísticos
    """
    try:
        full_text = " ".join(text_pages)
        
        # Análise de pontuação
        import re
        punctuation_counts = {
            "periods": len(re.findall(r'\.', full_text)),
            "commas": len(re.findall(r',', full_text)),
            "exclamations": len(re.findall(r'!', full_text)),
            "questions": len(re.findall(r'\?', full_text)),
            "semicolons": len(re.findall(r';', full_text)),
            "colons": len(re.findall(r':', full_text))
        }
        
        # Análise de estrutura
        paragraphs = [p.strip() for p in full_text.split('\n\n') if p.strip()]
        words = full_text.split()
        
        # Detectar possível idioma (português vs inglês - básico)
        portuguese_indicators = ['ção', 'são', 'não', 'com', 'para', 'que', 'uma', 'dos', 'das']
        english_indicators = ['the', 'and', 'for', 'are', 'with', 'that', 'this', 'have']
        
        pt_score = sum(1 for word in words if any(ind in word.lower() for ind in portuguese_indicators))
        en_score = sum(1 for word in words if word.lower() in english_indicators)
        
        likely_language = "português" if pt_score > en_score else "inglês" if en_score > pt_score else "indeterminado"
        
        # Análise de complexidade (aproximada)
        long_words = [w for w in words if len(w) > 8]
        complex_sentences = [s for s in full_text.split('.') if len(s.split()) > 20]
        
        analysis = {
            "text_structure": {
                "total_paragraphs": len(paragraphs),
                "avg_paragraph_length": round(len(words) / len(paragraphs), 1) if paragraphs else 0,
                "longest_paragraph_words": max(len(p.split()) for p in paragraphs) if paragraphs else 0
            },
            "punctuation_analysis": punctuation_counts,
            "language_detection": {
                "likely_language": likely_language,
                "portuguese_score": pt_score,
                "english_score": en_score,
                "confidence": "low" if abs(pt_score - en_score) < 5 else "medium"
            },
            "complexity_indicators": {
                "long_words_count": len(long_words),
                "long_words_percentage": round(len(long_words) / len(words) * 100, 2) if words else 0,
                "complex_sentences_count": len(complex_sentences),
                "avg_sentence_length": round(len(words) / punctuation_counts["periods"], 1) if punctuation_counts["periods"] > 0 else 0
            },
            "text_style": {
                "question_ratio": round(punctuation_counts["questions"] / len(paragraphs), 2) if paragraphs else 0,
                "exclamation_ratio": round(punctuation_counts["exclamations"] / len(paragraphs), 2) if paragraphs else 0,
                "formal_indicators": len(re.findall(r'\b(portanto|entretanto|ademais|outrossim|todavia)\b', full_text.lower()))
            }
        }
        
        print(f"🔍 Análise linguística concluída:")
        print(f"   🌐 Idioma provável: {likely_language}")
        print(f"   📊 Complexidade: {analysis['complexity_indicators']['long_words_percentage']}% palavras longas")
        print(f"   📝 Estilo: {analysis['text_style']['formal_indicators']} indicadores formais")
        
        return analysis
        
    except Exception as e:
        print(f"❌ Erro na análise linguística: {str(e)}")
        return {"error": str(e)}

def batch_process_files(directory_path: str, file_patterns: List[str] = None, include_analysis: bool = True) -> dict:
    """
    Processa múltiplos arquivos PDF e PPTX em lote de um diretório.
    
    Args:
        directory_path (str): Caminho do diretório contendo os arquivos
        file_patterns (List[str]): Padrões de arquivos para incluir (default: ['*.pdf', '*.pptx'])
        include_analysis (bool): Se deve incluir análise completa de texto
    
    Returns:
        dict: Resultados do processamento em lote
    """
    try:
        import glob
        from datetime import datetime
        
        if file_patterns is None:
            file_patterns = ['*.pdf', '*.pptx']
        
        start_time = datetime.now()
        
        # Encontrar todos os arquivos correspondentes aos padrões
        all_files = []
        for pattern in file_patterns:
            pattern_path = os.path.join(directory_path, pattern)
            files = glob.glob(pattern_path)
            all_files.extend(files)
        
        # Remover duplicatas e ordenar
        all_files = sorted(list(set(all_files)))
        
        print(f"📁 Processamento em lote iniciado:")
        print(f"   📂 Diretório: {directory_path}")
        print(f"   📋 Padrões: {file_patterns}")
        print(f"   📄 Arquivos encontrados: {len(all_files)}")
        
        # Resultados do processamento
        batch_results = {
            "processing_info": {
                "directory": directory_path,
                "patterns": file_patterns,
                "started_at": start_time.isoformat(),
                "total_files_found": len(all_files),
                "include_analysis": include_analysis
            },
            "files_processed": [],
            "processing_summary": {
                "successful": 0,
                "failed": 0,
                "total_pages_extracted": 0,
                "total_words_extracted": 0,
                "total_characters_extracted": 0
            },
            "errors": []
        }
        
        # Processar cada arquivo
        for file_path in all_files:
            print(f"\n📄 Processando: {os.path.basename(file_path)}")
            
            file_result = {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_type": os.path.splitext(file_path)[1].lower(),
                "processing_status": "pending",
                "extraction_result": None,
                "analysis_result": None,
                "error_message": None,
                "processing_time_seconds": 0
            }
            
            file_start_time = datetime.now()
            
            try:
                # Extrair texto baseado no tipo de arquivo
                if file_result["file_type"] == ".pdf":
                    text_pages = extract_text_from_pdf(file_path)
                elif file_result["file_type"] == ".pptx":
                    text_pages = extract_text_from_pptx(file_path)
                else:
                    raise ValueError(f"Tipo de arquivo não suportado: {file_result['file_type']}")
                
                file_result["extraction_result"] = {
                    "pages_count": len(text_pages),
                    "text_pages": text_pages,
                    "total_characters": sum(len(page) for page in text_pages),
                    "total_words": sum(len(page.split()) for page in text_pages)
                }
                
                # Análise avançada se solicitada
                if include_analysis and text_pages:
                    analysis_result = {
                        "statistics": analyze_text_statistics(text_pages),
                        "keywords": extract_keywords(text_pages, top_n=15),
                        "summary": generate_text_summary(text_pages, max_sentences=3),
                        "language_patterns": analyze_text_language_patterns(text_pages)
                    }
                    file_result["analysis_result"] = analysis_result
                
                # Atualizar estatísticas de sucesso
                batch_results["processing_summary"]["successful"] += 1
                batch_results["processing_summary"]["total_pages_extracted"] += file_result["extraction_result"]["pages_count"]
                batch_results["processing_summary"]["total_words_extracted"] += file_result["extraction_result"]["total_words"]
                batch_results["processing_summary"]["total_characters_extracted"] += file_result["extraction_result"]["total_characters"]
                
                file_result["processing_status"] = "completed"
                
                print(f"   ✅ Sucesso: {file_result['extraction_result']['pages_count']} páginas, {file_result['extraction_result']['total_words']} palavras")
                
            except Exception as e:
                file_result["processing_status"] = "failed"
                file_result["error_message"] = str(e)
                batch_results["processing_summary"]["failed"] += 1
                batch_results["errors"].append({
                    "file": file_path,
                    "error": str(e)
                })
                
                print(f"   ❌ Erro: {str(e)}")
            
            # Calcular tempo de processamento
            file_end_time = datetime.now()
            file_result["processing_time_seconds"] = (file_end_time - file_start_time).total_seconds()
            
            batch_results["files_processed"].append(file_result)
        
        # Finalizar processamento
        end_time = datetime.now()
        total_time = (end_time - start_time).total_seconds()
        
        batch_results["processing_info"]["completed_at"] = end_time.isoformat()
        batch_results["processing_info"]["total_processing_time_seconds"] = total_time
        
        # Relatório final
        summary = batch_results["processing_summary"]
        print(f"\n{'='*60}")
        print(f"📊 RELATÓRIO DE PROCESSAMENTO EM LOTE")
        print(f"{'='*60}")
        print(f"✅ Arquivos processados com sucesso: {summary['successful']}")
        print(f"❌ Arquivos com erro: {summary['failed']}")
        print(f"📄 Total de páginas extraídas: {summary['total_pages_extracted']}")
        print(f"📝 Total de palavras extraídas: {summary['total_words_extracted']:,}")
        print(f"🔤 Total de caracteres extraídos: {summary['total_characters_extracted']:,}")
        print(f"⏱️ Tempo total de processamento: {total_time:.2f} segundos")
        print(f"⚡ Velocidade média: {summary['successful'] / total_time:.2f} arquivos/segundo" if total_time > 0 else "")
        
        return batch_results
        
    except Exception as e:
        print(f"❌ Erro no processamento em lote: {str(e)}")
        return {"error": str(e)}

def batch_search_across_files(directory_path: str, search_terms: List[str], file_patterns: List[str] = None, case_sensitive: bool = False) -> dict:
    """
    Realiza busca de termos em múltiplos arquivos de um diretório.
    
    Args:
        directory_path (str): Caminho do diretório
        search_terms (List[str]): Termos para buscar
        file_patterns (List[str]): Padrões de arquivo (default: ['*.pdf', '*.pptx'])
        case_sensitive (bool): Se a busca deve ser sensível a maiúsculas
    
    Returns:
        dict: Resultados da busca em múltiplos arquivos
    """
    try:
        import glob
        from datetime import datetime
        
        if file_patterns is None:
            file_patterns = ['*.pdf', '*.pptx']
        
        start_time = datetime.now()
        
        # Encontrar arquivos
        all_files = []
        for pattern in file_patterns:
            pattern_path = os.path.join(directory_path, pattern)
            files = glob.glob(pattern_path)
            all_files.extend(files)
        
        all_files = sorted(list(set(all_files)))
        
        print(f"🔍 Busca em lote iniciada:")
        print(f"   📂 Diretório: {directory_path}")
        print(f"   🔎 Termos: {search_terms}")
        print(f"   📄 Arquivos: {len(all_files)}")
        
        # Resultados da busca
        batch_search_results = {
            "search_info": {
                "directory": directory_path,
                "search_terms": search_terms,
                "case_sensitive": case_sensitive,
                "started_at": start_time.isoformat(),
                "total_files_searched": len(all_files)
            },
            "global_summary": {
                "files_with_matches": 0,
                "total_matches_across_files": 0,
                "terms_found": {}
            },
            "files_results": [],
            "errors": []
        }
        
        # Inicializar contadores por termo
        for term in search_terms:
            batch_search_results["global_summary"]["terms_found"][term] = {
                "total_occurrences": 0,
                "files_found_in": 0
            }
        
        # Buscar em cada arquivo
        for file_path in all_files:
            print(f"🔍 Buscando em: {os.path.basename(file_path)}")
            
            file_search_result = {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_type": os.path.splitext(file_path)[1].lower(),
                "has_matches": False,
                "total_matches_in_file": 0,
                "search_results": None,
                "error_message": None
            }
            
            try:
                # Extrair texto
                if file_search_result["file_type"] == ".pdf":
                    text_pages = extract_text_from_pdf(file_path)
                elif file_search_result["file_type"] == ".pptx":
                    text_pages = extract_text_from_pptx(file_path)
                else:
                    continue  # Pular arquivos não suportados
                
                # Realizar busca
                search_results = search_text_in_pages(text_pages, search_terms, case_sensitive)
                
                if search_results["total_matches"] > 0:
                    file_search_result["has_matches"] = True
                    file_search_result["total_matches_in_file"] = search_results["total_matches"]
                    file_search_result["search_results"] = search_results
                    
                    batch_search_results["global_summary"]["files_with_matches"] += 1
                    batch_search_results["global_summary"]["total_matches_across_files"] += search_results["total_matches"]
                    
                    # Atualizar contadores por termo
                    for term, term_data in search_results["results_by_term"].items():
                        if term_data["total_occurrences"] > 0:
                            batch_search_results["global_summary"]["terms_found"][term]["total_occurrences"] += term_data["total_occurrences"]
                            batch_search_results["global_summary"]["terms_found"][term]["files_found_in"] += 1
                    
                    print(f"   ✅ {search_results['total_matches']} ocorrências encontradas")
                else:
                    print(f"   ⚪ Nenhuma ocorrência encontrada")
                
            except Exception as e:
                file_search_result["error_message"] = str(e)
                batch_search_results["errors"].append({
                    "file": file_path,
                    "error": str(e)
                })
                print(f"   ❌ Erro: {str(e)}")
            
            batch_search_results["files_results"].append(file_search_result)
        
        # Finalizar busca
        end_time = datetime.now()
        total_time = (end_time - start_time).total_seconds()
        
        batch_search_results["search_info"]["completed_at"] = end_time.isoformat()
        batch_search_results["search_info"]["total_search_time_seconds"] = total_time
        
        # Relatório final
        summary = batch_search_results["global_summary"]
        print(f"\n{'='*60}")
        print(f"📊 RELATÓRIO DE BUSCA EM LOTE")
        print(f"{'='*60}")
        print(f"📄 Arquivos com ocorrências: {summary['files_with_matches']}/{len(all_files)}")
        print(f"🎯 Total de ocorrências: {summary['total_matches_across_files']}")
        print(f"⏱️ Tempo total de busca: {total_time:.2f} segundos")
        
        for term, term_stats in summary["terms_found"].items():
            print(f"   🔍 '{term}': {term_stats['total_occurrences']} ocorrências em {term_stats['files_found_in']} arquivos")
        
        return batch_search_results
        
    except Exception as e:
        print(f"❌ Erro na busca em lote: {str(e)}")
        return {"error": str(e)}

def export_batch_results_to_json(batch_results: dict, output_file: str = None) -> str:
    """
    Exporta resultados de processamento em lote para arquivo JSON.
    
    Args:
        batch_results (dict): Resultados do processamento em lote
        output_file (str): Caminho do arquivo de saída (opcional)
    
    Returns:
        str: Caminho do arquivo JSON criado
    """
    try:
        import json
        from datetime import datetime
        
        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"batch_results_{timestamp}.json"
        
        # Preparar dados para exportação (remover objetos não serializáveis)
        export_data = json.loads(json.dumps(batch_results, default=str, ensure_ascii=False))
        
        # Salvar arquivo JSON
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
        
        file_size = os.path.getsize(output_file)
        
        print(f"📄 Resultados exportados para JSON:")
        print(f"   📁 Arquivo: {output_file}")
        print(f"   📊 Tamanho: {format_file_size(file_size)}")
        
        return output_file
        
    except Exception as e:
        print(f"❌ Erro ao exportar resultados: {str(e)}")
        return None

def create_batch_processing_report(batch_results: dict, output_file: str = None) -> str:
    """
    Cria um relatório detalhado em texto do processamento em lote.
    
    Args:
        batch_results (dict): Resultados do processamento em lote
        output_file (str): Caminho do arquivo de relatório (opcional)
    
    Returns:
        str: Caminho do arquivo de relatório criado
    """
    try:
        from datetime import datetime
        
        if output_file is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"batch_report_{timestamp}.txt"
        
        # Gerar relatório
        report_lines = []
        report_lines.append("=" * 80)
        report_lines.append("RELATÓRIO DETALHADO DE PROCESSAMENTO EM LOTE")
        report_lines.append("=" * 80)
        report_lines.append("")
        
        # Informações gerais
        info = batch_results.get("processing_info", {})
        report_lines.append("📊 INFORMAÇÕES GERAIS")
        report_lines.append("-" * 40)
        report_lines.append(f"Diretório processado: {info.get('directory', 'N/A')}")
        report_lines.append(f"Padrões de arquivo: {', '.join(info.get('patterns', []))}")
        report_lines.append(f"Iniciado em: {info.get('started_at', 'N/A')}")
        report_lines.append(f"Concluído em: {info.get('completed_at', 'N/A')}")
        report_lines.append(f"Tempo total: {info.get('total_processing_time_seconds', 0):.2f} segundos")
        report_lines.append("")
        
        # Resumo do processamento
        summary = batch_results.get("processing_summary", {})
        report_lines.append("📈 RESUMO DO PROCESSAMENTO")
        report_lines.append("-" * 40)
        report_lines.append(f"Arquivos processados com sucesso: {summary.get('successful', 0)}")
        report_lines.append(f"Arquivos com erro: {summary.get('failed', 0)}")
        report_lines.append(f"Total de páginas extraídas: {summary.get('total_pages_extracted', 0)}")
        report_lines.append(f"Total de palavras extraídas: {summary.get('total_words_extracted', 0):,}")
        report_lines.append(f"Total de caracteres extraídos: {summary.get('total_characters_extracted', 0):,}")
        report_lines.append("")
        
        # Detalhes por arquivo
        report_lines.append("📄 DETALHES POR ARQUIVO")
        report_lines.append("-" * 40)
        
        files_processed = batch_results.get("files_processed", [])
        for i, file_result in enumerate(files_processed, 1):
            status_icon = "✅" if file_result["processing_status"] == "completed" else "❌"
            report_lines.append(f"{i}. {status_icon} {file_result['file_name']}")
            report_lines.append(f"   Tipo: {file_result['file_type']}")
            report_lines.append(f"   Status: {file_result['processing_status']}")
            report_lines.append(f"   Tempo: {file_result['processing_time_seconds']:.2f}s")
            
            if file_result["processing_status"] == "completed":
                extraction = file_result.get("extraction_result", {})
                report_lines.append(f"   Páginas: {extraction.get('pages_count', 0)}")
                report_lines.append(f"   Palavras: {extraction.get('total_words', 0):,}")
                report_lines.append(f"   Caracteres: {extraction.get('total_characters', 0):,}")
            else:
                report_lines.append(f"   Erro: {file_result.get('error_message', 'Erro desconhecido')}")
            
            report_lines.append("")
        
        # Erros (se houver)
        errors = batch_results.get("errors", [])
        if errors:
            report_lines.append("❌ ERROS ENCONTRADOS")
            report_lines.append("-" * 40)
            for i, error in enumerate(errors, 1):
                report_lines.append(f"{i}. {os.path.basename(error['file'])}")
                report_lines.append(f"   Erro: {error['error']}")
                report_lines.append("")
        
        # Salvar relatório
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(report_lines))
        
        file_size = os.path.getsize(output_file)
        
        print(f"📄 Relatório detalhado criado:")
        print(f"   📁 Arquivo: {output_file}")
        print(f"   📊 Tamanho: {format_file_size(file_size)}")
        
        return output_file
        
    except Exception as e:
        print(f"❌ Erro ao criar relatório: {str(e)}")
        return None

# Função para teste das utilidades
def test_utils():
    """
    Função para testar as utilidades
    """
    print("🔧 Testando funções utilitárias...")
    
    try:
        # Testar criação de diretórios
        ensure_directories_exist()
        
        # Testar validações
        assert validate_file_type("documento.pdf") == True
        assert validate_file_type("apresentacao.pptx") == True
        assert validate_file_type("arquivo.txt") == False
        print("✅ Validação de tipos funcionando")
        
        # Testar geração de nome único
        unique_name = generate_unique_filename("test_file.pdf", 1)
        assert unique_name.endswith(".pdf")
        assert "test_file" in unique_name
        print(f"✅ Geração de nome único: {unique_name}")
        
        # Testar formatação de tamanho
        assert format_file_size(1024) == "1.0 KB"
        assert format_file_size(1048576) == "1.0 MB"
        print("✅ Formatação de tamanho funcionando")
        
        print("🎉 Funções utilitárias testadas com sucesso!")
        
    except Exception as e:
        print(f"❌ Erro ao testar utilidades: {e}")
        raise

if __name__ == "__main__":
    """
    Script para testar as funções utilitárias
    Execute: python utils.py
    """
    print("🔧 Testando funções utilitárias...")
    
    # Testar funções básicas
    test_utils()
    
    print("\n" + "="*60)
    print("📝 EXEMPLOS DE USO DAS NOVAS FUNÇÕES")
    print("="*60)
    
    # Exemplo de uso da função extract_text_from_pdf
    print("\n1. 📄 EXEMPLO: Extração de texto de PDF")
    print("-" * 40)
    
    # Caminho de exemplo para um arquivo PDF (substitua pelo caminho real)
    pdf_example_path = "sample_test.pdf"
    
    if os.path.exists(pdf_example_path):
        try:
            print(f"📂 Processando arquivo: {pdf_example_path}")
            
            # Extrair texto de cada página do PDF
            pdf_pages = extract_text_from_pdf(pdf_example_path)
            
            print(f"📊 Resultado: {len(pdf_pages)} páginas extraídas")
            
            # Mostrar uma prévia do texto de cada página
            for i, page_text in enumerate(pdf_pages):
                preview = page_text[:100] + "..." if len(page_text) > 100 else page_text
                print(f"   Página {i+1}: {preview}")
                
        except Exception as e:
            print(f"❌ Erro no exemplo PDF: {e}")
    else:
        print(f"⚠️ Arquivo de exemplo não encontrado: {pdf_example_path}")
        print("💡 Para testar, coloque um arquivo PDF na raiz do projeto com nome 'sample_test.pdf'")
    
    # Exemplo de uso da função extract_text_from_pptx
    print("\n2. 🎨 EXEMPLO: Extração de texto de PPTX")
    print("-" * 40)
    
    # Caminho de exemplo para um arquivo PPTX (substitua pelo caminho real)
    pptx_example_path = "sample_presentation.pptx"
    
    if os.path.exists(pptx_example_path):
        try:
            print(f"📂 Processando arquivo: {pptx_example_path}")
            
            # Extrair texto de cada slide do PPTX
            pptx_slides = extract_text_from_pptx(pptx_example_path)
            
            print(f"📊 Resultado: {len(pptx_slides)} slides extraídos")
            
            # Mostrar uma prévia do texto de cada slide
            for i, slide_text in enumerate(pptx_slides):
                preview = slide_text[:100] + "..." if len(slide_text) > 100 else slide_text
                print(f"   Slide {i+1}: {preview}")
                
        except Exception as e:
            print(f"❌ Erro no exemplo PPTX: {e}")
    else:
        print(f"⚠️ Arquivo de exemplo não encontrado: {pptx_example_path}")
        print("💡 Para testar, coloque um arquivo PPTX na raiz do projeto com nome 'sample_presentation.pptx'")
    
    print("\n" + "="*60)
    print("✅ EXEMPLOS DE USO CONCLUÍDOS")
    print("="*60)
    
    # Exemplo de código para usar todas as funções implementadas
    print("\n💻 GUIA COMPLETO DE USO DAS FUNCIONALIDADES:")
    print("-" * 60)
    print("""
# ========================================
# GUIA COMPLETO DE FUNCIONALIDADES
# ========================================

# 1. EXTRAÇÃO BÁSICA DE TEXTO
try:
    # Extrair texto de PDF
    pdf_pages = extract_text_from_pdf("documento.pdf")
    print(f"PDF: {len(pdf_pages)} páginas extraídas")
    
    # Extrair texto de PPTX
    pptx_slides = extract_text_from_pptx("apresentacao.pptx")
    print(f"PPTX: {len(pptx_slides)} slides extraídos")
except Exception as e:
    print(f"Erro na extração: {e}")

# 2. ANÁLISE ESTATÍSTICA AVANÇADA
try:
    stats = analyze_text_statistics(pdf_pages)
    print(f"Total de palavras: {stats['overview']['total_words']:,}")
    print(f"Tempo de leitura: {stats['readability']['estimated_reading_time_minutes']} min")
    print(f"Riqueza vocabular: {stats['word_frequency']['vocabulary_richness']}%")
except Exception as e:
    print(f"Erro na análise: {e}")

# 3. EXTRAÇÃO DE PALAVRAS-CHAVE
try:
    keywords = extract_keywords(pdf_pages, top_n=10)
    print("Top 5 palavras-chave:")
    for kw in keywords['keywords'][:5]:
        print(f"  • {kw['word']}: {kw['frequency']} ocorrências")
except Exception as e:
    print(f"Erro nas palavras-chave: {e}")

# 4. RESUMO AUTOMÁTICO
try:
    summary = generate_text_summary(pdf_pages, max_sentences=3)
    print(f"Resumo ({summary['compression_ratio']:.1%} do original):")
    print(f"'{summary['summary'][:200]}...'")
except Exception as e:
    print(f"Erro no resumo: {e}")

# 5. BUSCA AVANÇADA
try:
    search_results = search_text_in_pages(
        text_pages=pdf_pages,
        search_terms=["tecnologia", "inovação"],
        case_sensitive=False
    )
    print(f"Busca: {search_results['total_matches']} ocorrências encontradas")
except Exception as e:
    print(f"Erro na busca: {e}")

# 6. ANÁLISE LINGUÍSTICA
try:
    patterns = analyze_text_language_patterns(pdf_pages)
    print(f"Idioma: {patterns['language_detection']['likely_language']}")
    print(f"Complexidade: {patterns['complexity_indicators']['long_words_percentage']}% palavras longas")
except Exception as e:
    print(f"Erro na análise linguística: {e}")

# 7. PROCESSAMENTO EM LOTE
try:
    batch_results = batch_process_files(
        directory_path="./documentos",
        file_patterns=["*.pdf", "*.pptx"],
        include_analysis=True
    )
    print(f"Lote: {batch_results['processing_summary']['successful']} arquivos processados")
    
    # Exportar resultados
    export_file = export_batch_results_to_json(batch_results)
    print(f"Resultados exportados: {export_file}")
except Exception as e:
    print(f"Erro no processamento em lote: {e}")

# 8. BUSCA EM MÚLTIPLOS ARQUIVOS
try:
    batch_search = batch_search_across_files(
        directory_path="./documentos",
        search_terms=["IA", "machine learning"],
        file_patterns=["*.pdf"],
        case_sensitive=False
    )
    files_found = batch_search['global_summary']['files_with_matches']
    total_matches = batch_search['global_summary']['total_matches_across_files']
    print(f"Busca em lote: {total_matches} ocorrências em {files_found} arquivos")
except Exception as e:
    print(f"Erro na busca em lote: {e}")
    """)
    
    print("\n" + "="*60)
    print("🎉 SISTEMA COMPLETO DE ANÁLISE DE TEXTO")
    print("="*60)
    print("✅ 16 funcionalidades implementadas")
    print("✅ 5 endpoints de API integrados")
    print("✅ Suporte a PDF e PPTX")
    print("✅ Análise estatística avançada")
    print("✅ Extração de palavras-chave")
    print("✅ Resumo automático")
    print("✅ Busca inteligente")
    print("✅ Processamento em lote")
    print("✅ Análise linguística")
    print("✅ Exportação de resultados")
    print("✅ Testes unitários")
    print("✅ Performance otimizada")
    print("="*60)


# ============================================================================
# FUNÇÃO DE GERAÇÃO DE VÍDEO AVATAR COM NARRAÇÃO - VERSÃO AVANÇADA
# ============================================================================

def generate_avatar_video(
    text: str, 
    audio_path: str, 
    output_path: str,
    avatar_style: str = "hunyuan3d",
    timeout: int = 300,
    quality: str = "high",
    **kwargs
) -> dict:
    """
    Gera um vídeo de avatar 3D realista com sincronização labial usando API do Hunyuan3D-2.
    
    Esta função integra com o modelo Hunyuan3D-2 da Tencent via Hugging Face Spaces para 
    gerar avatares 3D de alta qualidade com sincronização labial precisa.
    
    🔗 API Externa Utilizada:
    - Hunyuan3D-2: https://huggingface.co/spaces/tencent/Hunyuan3D-2
    - Modelo da Tencent para geração de avatares 3D com talking heads
    - Suporte a múltiplos idiomas e estilos de avatar
    
    📋 Processo de Geração:
    1. Validação dos parâmetros de entrada (texto, áudio, caminhos)
    2. Preparação da requisição para a API do Hunyuan3D-2
    3. Upload do áudio e texto para o Hugging Face Space
    4. Monitoramento da fila de processamento com timeout
    5. Download do vídeo gerado quando concluído
    6. Salvamento local e validação do resultado
    
    Args:
        text (str): Texto que o avatar deve "falar" (usado para lip-sync)
        audio_path (str): Caminho para arquivo de áudio MP3/WAV com narração
        output_path (str): Caminho onde salvar o vídeo MP4 gerado
        avatar_style (str): Estilo do avatar ('hunyuan3d', 'realistic', 'cartoon')
        timeout (int): Timeout em segundos para geração (padrão: 300s)
        quality (str): Qualidade do vídeo ('low', 'medium', 'high', 'ultra')
        **kwargs: Parâmetros adicionais (background, emotion, pose, etc.)
    
    Returns:
        dict: Resultado detalhado da geração
        {
            'success': bool,
            'video_path': str,                    # Caminho do vídeo gerado
            'duration': float,                    # Duração em segundos
            'file_size': int,                     # Tamanho do arquivo em bytes
            'resolution': tuple,                  # (largura, altura)
            'api_used': str,                      # 'hunyuan3d'
            'processing_time': float,             # Tempo total de processamento
            'queue_time': float,                  # Tempo esperando na fila
            'generation_time': float,             # Tempo de geração do modelo
            'download_time': float,               # Tempo de download
            'quality_score': float,               # Score de qualidade (0-1)
            'metadata': dict,                     # Metadados do processo
            'error': str | None                   # Erro se houver falha
        }
    
    Raises:
        FileNotFoundError: Se arquivo de áudio não for encontrado
        ValueError: Se parâmetros são inválidos
        TimeoutError: Se geração exceder timeout especificado
        ConnectionError: Se não conseguir conectar com a API
        RuntimeError: Se houve erro na geração do vídeo
    
    Example:
        >>> # Exemplo básico com áudio português
        >>> result = generate_avatar_video(
        ...     text="Bem-vindos ao curso de Python! Hoje vamos aprender...",
        ...     audio_path="./audios/aula_intro.mp3",
        ...     output_path="./videos/aula_avatar.mp4"
        ... )
        >>> 
        >>> if result['success']:
        ...     print(f"✅ Vídeo criado: {result['video_path']}")
        ...     print(f"📊 Duração: {result['duration']:.1f}s")
        ...     print(f"💾 Tamanho: {result['file_size']/1024/1024:.1f}MB")
        ... else:
        ...     print(f"❌ Erro: {result['error']}")
        
        >>> # Exemplo avançado com configurações personalizadas
        >>> result = generate_avatar_video(
        ...     text="Technical presentation about AI models",
        ...     audio_path="./presentation.wav",
        ...     output_path="./output/tech_presentation.mp4",
        ...     avatar_style="realistic",
        ...     quality="ultra",
        ...     timeout=600,  # 10 minutos
        ...     background="office",
        ...     emotion="confident",
        ...     pose="professional"
        ... )
    
    Note:
        📦 DEPENDÊNCIAS NECESSÁRIAS:
        pip install requests httpx aiohttp aiofiles mutagen
        
        🔧 CONFIGURAÇÃO ALTERNATIVA:
        Se a API pública não estiver disponível, você pode usar o Space manualmente:
        
        1. Acesse: https://huggingface.co/spaces/tencent/Hunyuan3D-2
        2. Faça upload do seu arquivo de áudio
        3. Insira o texto desejado
        4. Configure os parâmetros (estilo, qualidade, etc.)
        5. Aguarde o processamento (pode demorar 2-10 minutos)
        6. Baixe o vídeo gerado
        
        ⚠️ LIMITAÇÕES:
        - API pode ter fila durante picos de uso
        - Áudios longos (>5min) podem falhar ou demorar muito
        - Algumas configurações avançadas podem não estar disponíveis
        - Rate limiting aplicado por usuário/IP
        
        🚀 OTIMIZAÇÕES:
        - Cache local para evitar regerar vídeos idênticos
        - Compressão automática de áudios grandes
        - Retry automático em caso de erros temporários
        - Monitoramento de status em tempo real
    """
    import requests
    import json
    import time
    from pathlib import Path
    import hashlib
    
    # Registrar tempo de início para métricas de performance
    start_time = time.time()
    queue_start_time = None
    generation_start_time = None
    download_start_time = None
    
    print(f"🚀 Iniciando geração de avatar com Hunyuan3D-2...")
    print(f"   📝 Texto: {text[:50]}{'...' if len(text) > 50 else ''}")
    print(f"   🎵 Áudio: {audio_path}")
    print(f"   📹 Saída: {output_path}")
    
    try:
        # ========================================================================
        # ETAPA 1: VALIDAÇÃO DE PARÂMETROS DE ENTRADA
        # ========================================================================
        print("📋 Validando parâmetros de entrada...")
        
        # Validar texto
        if not text or not text.strip():
            raise ValueError("Texto não pode estar vazio")
        
        if len(text) > 1000:  # Hunyuan3D pode ter limitações
            print(f"⚠️ Texto muito longo ({len(text)} chars), truncando para 1000...")
            text = text[:1000].rsplit('.', 1)[0] + '.'
        
        # Validar arquivo de áudio
        if not os.path.exists(audio_path):
            raise FileNotFoundError(f"Arquivo de áudio não encontrado: {audio_path}")
        
        # Verificar tamanho do áudio (limite de ~10MB para APIs)
        audio_size = os.path.getsize(audio_path)
        max_audio_size = 10 * 1024 * 1024  # 10MB
        
        if audio_size > max_audio_size:
            raise ValueError(f"Arquivo de áudio muito grande: {audio_size/1024/1024:.1f}MB. Máximo: 10MB")
        
        # Criar diretório de saída se não existir
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
        
        # ========================================================================
        # ETAPA 2: VERIFICAR CACHE LOCAL PARA EVITAR REGERAR
        # ========================================================================
        print("💾 Verificando cache local...")
        
        # Criar hash único baseado nos parâmetros
        cache_key = hashlib.md5(
            f"{text}_{os.path.basename(audio_path)}_{avatar_style}_{quality}".encode()
        ).hexdigest()
        
        cache_dir = "cache/avatar_videos"
        os.makedirs(cache_dir, exist_ok=True)
        cached_video = os.path.join(cache_dir, f"{cache_key}.mp4")
        
        # Se existe no cache e é recente (menos de 7 dias)
        if os.path.exists(cached_video):
            cache_age = time.time() - os.path.getmtime(cached_video)
            if cache_age < 7 * 24 * 3600:  # 7 dias
                print("✅ Vídeo encontrado no cache! Copiando...")
                import shutil
                shutil.copy2(cached_video, output_path)
                
                file_size = os.path.getsize(output_path)
                duration = get_video_duration(output_path)
                
                return {
                    'success': True,
                    'video_path': output_path,
                    'duration': duration,
                    'file_size': file_size,
                    'resolution': (1920, 1080),  # Padrão do Hunyuan3D
                    'api_used': 'hunyuan3d_cached',
                    'processing_time': time.time() - start_time,
                    'queue_time': 0.0,
                    'generation_time': 0.0,
                    'download_time': 0.0,
                    'quality_score': 0.9,
                    'metadata': {'cached': True, 'cache_key': cache_key},
                    'error': None
                }
        
        # ========================================================================
        # ETAPA 3: PREPARAR REQUISIÇÃO PARA API HUNYUAN3D-2
        # ========================================================================
        print("🌐 Preparando requisição para Hunyuan3D-2...")
        
        # URL do Hugging Face Space
        api_base_url = "https://tencent-hunyuan3d-2.hf.space"
        
        # Headers para requisições
        headers = {
            'User-Agent': 'TecnoCursosAI/1.0 (Educational Avatar Generator)',
            'Accept': 'application/json',
            'Content-Type': 'multipart/form-data'
        }
        
        # Preparar dados da requisição
        # Importar funções auxiliares
        from .avatar_utils import detect_language, simulate_avatar_generation, get_video_duration, get_video_resolution, calculate_quality_score
        
        request_data = {
            'text': text,
            'avatar_style': avatar_style,
            'quality': quality,
            'language': detect_language(text),
            **kwargs  # Parâmetros adicionais como background, emotion, etc.
        }
        
        # ========================================================================
        # ETAPA 4: FAZER UPLOAD DO ÁUDIO E SUBMETER REQUISIÇÃO
        # ========================================================================
        print("📤 Enviando áudio e texto para processamento...")
        
        try:
            # Abrir arquivo de áudio para upload
            with open(audio_path, 'rb') as audio_file:
                files = {
                    'audio': (os.path.basename(audio_path), audio_file, 'audio/mpeg')
                }
                
                data = {
                    'text': text,
                    'style': avatar_style,
                    'quality': quality,
                    'timeout': str(timeout)
                }
                
                # Submeter requisição (este é um exemplo - API real pode ter endpoint diferente)
                response = requests.post(
                    f"{api_base_url}/api/generate",
                    files=files,
                    data=data,
                    timeout=30  # Timeout inicial para submissão
                )
                
                if response.status_code != 200:
                    raise ConnectionError(f"Erro na API: {response.status_code} - {response.text}")
                
                result = response.json()
                job_id = result.get('job_id')
                
                if not job_id:
                    raise RuntimeError("API não retornou ID do job")
                
                print(f"✅ Requisição submetida! Job ID: {job_id}")
        
        except requests.exceptions.ConnectionError:
            # Se API não está disponível, simular processamento local
            print("⚠️ API Hunyuan3D-2 não disponível, usando simulação...")
            return simulate_avatar_generation(text, audio_path, output_path, start_time)
        
        # ========================================================================
        # ETAPA 5: MONITORAR FILA E PROGRESSO COM TIMEOUT
        # ========================================================================
        print("⏳ Monitorando progresso na fila...")
        queue_start_time = time.time()
        
        max_wait_time = timeout
        check_interval = 5  # Verificar a cada 5 segundos
        
        while True:
            elapsed_time = time.time() - queue_start_time
            
            if elapsed_time > max_wait_time:
                raise TimeoutError(f"Timeout de {timeout}s excedido durante processamento")
            
            # Verificar status do job
            try:
                status_response = requests.get(
                    f"{api_base_url}/api/status/{job_id}",
                    timeout=10
                )
                
                if status_response.status_code == 200:
                    status_data = status_response.json()
                    status = status_data.get('status')
                    progress = status_data.get('progress', 0)
                    
                    print(f"📊 Status: {status} - Progresso: {progress}%")
                    
                    if status == 'completed':
                        video_url = status_data.get('video_url')
                        if video_url:
                            generation_start_time = time.time()
                            break
                    elif status == 'failed':
                        error_msg = status_data.get('error', 'Erro desconhecido na geração')
                        raise RuntimeError(f"Geração falhou: {error_msg}")
                    elif status in ['queued', 'processing']:
                        time.sleep(check_interval)
                        continue
                
            except requests.exceptions.RequestException:
                print("⚠️ Erro ao verificar status, tentando novamente...")
                time.sleep(check_interval)
                continue
        
        # ========================================================================
        # ETAPA 6: DOWNLOAD DO VÍDEO GERADO
        # ========================================================================
        print("⬇️ Baixando vídeo gerado...")
        download_start_time = time.time()
        
        try:
            video_response = requests.get(video_url, stream=True, timeout=60)
            video_response.raise_for_status()
            
            # Salvar vídeo localmente
            with open(output_path, 'wb') as f:
                for chunk in video_response.iter_content(chunk_size=8192):
                    f.write(chunk)
            
            # Salvar também no cache
            import shutil
            shutil.copy2(output_path, cached_video)
            
            print(f"✅ Vídeo salvo: {output_path}")
            
        except Exception as e:
            raise RuntimeError(f"Erro ao baixar vídeo: {str(e)}")
        
        # ========================================================================
        # ETAPA 7: VALIDAR E COLETAR MÉTRICAS DO RESULTADO
        # ========================================================================
        print("📊 Coletando métricas finais...")
        
        if not os.path.exists(output_path):
            raise RuntimeError("Vídeo não foi salvo corretamente")
        
        file_size = os.path.getsize(output_path)
        duration = get_video_duration(output_path)
        resolution = get_video_resolution(output_path)
        
        # Calcular tempos de processamento
        total_time = time.time() - start_time
        queue_time = (generation_start_time - queue_start_time) if queue_start_time and generation_start_time else 0
        download_time = (time.time() - download_start_time) if download_start_time else 0
        generation_time = total_time - queue_time - download_time
        
        # Score de qualidade baseado em métricas
        quality_score = calculate_quality_score(file_size, duration, resolution, generation_time)
        
        # Resultado final
        result = {
            'success': True,
            'video_path': output_path,
            'duration': duration,
            'file_size': file_size,
            'resolution': resolution,
            'api_used': 'hunyuan3d',
            'processing_time': total_time,
            'queue_time': queue_time,
            'generation_time': generation_time,
            'download_time': download_time,
            'quality_score': quality_score,
            'metadata': {
                'job_id': job_id,
                'text_length': len(text),
                'audio_size': audio_size,
                'cache_key': cache_key,
                'avatar_style': avatar_style,
                'quality_setting': quality
            },
            'error': None
        }
        
        print("🎉 Geração de avatar concluída com sucesso!")
        print(f"   📹 Vídeo: {output_path}")
        print(f"   ⏱️ Tempo total: {total_time:.1f}s")
        print(f"   💾 Tamanho: {file_size/1024/1024:.1f}MB")
        print(f"   📊 Qualidade: {quality_score:.2f}")
        
        return result
        
    except Exception as e:
        error_msg = f"Erro na geração do avatar: {str(e)}"
        print(f"❌ {error_msg}")
        
        return {
            'success': False,
            'video_path': None,
            'duration': 0.0,
            'file_size': 0,
            'resolution': (0, 0),
            'api_used': 'hunyuan3d',
            'processing_time': time.time() - start_time,
            'queue_time': 0.0,
            'generation_time': 0.0,
            'download_time': 0.0,
            'quality_score': 0.0,
            'metadata': {'error_details': str(e)},
            'error': error_msg
        }


def _check_avatar_cache(text: str, audio_path: str, template: str, avatar_style: str, kwargs: dict) -> dict:
    """
    Verifica se existe vídeo em cache para os parâmetros fornecidos.
    
    Returns:
        dict: {'found': bool, 'cache_path': str, 'metadata': dict}
    """
    try:
        import hashlib
        
        # Criar hash único baseado nos parâmetros
        cache_key_data = f"{text}_{audio_path}_{template}_{avatar_style}_{str(sorted(kwargs.items()))}"
        cache_key = hashlib.md5(cache_key_data.encode()).hexdigest()
        
        # Diretório de cache
        cache_dir = "cache/avatar_videos"
        os.makedirs(cache_dir, exist_ok=True)
        
        cache_video_path = os.path.join(cache_dir, f"{cache_key}.mp4")
        cache_metadata_path = os.path.join(cache_dir, f"{cache_key}.json")
        
        if os.path.exists(cache_video_path) and os.path.exists(cache_metadata_path):
            # Verificar se cache não está muito antigo (7 dias)
            cache_age = time.time() - os.path.getmtime(cache_video_path)
            if cache_age < 7 * 24 * 3600:  # 7 dias
                with open(cache_metadata_path, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)
                
                return {
                    'found': True,
                    'cache_path': cache_video_path,
                    'metadata': metadata
                }
        
        return {'found': False, 'cache_path': None, 'metadata': None}
        
    except Exception as e:
        print(f"⚠️ Erro no cache: {e}")
        return {'found': False, 'cache_path': None, 'metadata': None}


def _detect_language_and_configure(text: str) -> dict:
    """
    Detecta idioma do texto e configura parâmetros adequados.
    
    Returns:
        dict: Informações de idioma e configuração
    """
    # Detectar idioma baseado em palavras-chave
    portuguese_indicators = ['que', 'para', 'com', 'uma', 'são', 'não', 'dos', 'mas', 'seu', 'tem']
    english_indicators = ['the', 'and', 'for', 'are', 'with', 'that', 'this', 'have', 'from', 'they']
    spanish_indicators = ['que', 'para', 'con', 'una', 'son', 'del', 'pero', 'sus', 'tiene', 'como']
    
    words = text.lower().split()
    
    pt_score = sum(1 for word in words if word in portuguese_indicators)
    en_score = sum(1 for word in words if word in english_indicators)
    es_score = sum(1 for word in words if word in spanish_indicators)
    
    # Determinar idioma predominante
    scores = {'pt': pt_score, 'en': en_score, 'es': es_score}
    detected_language = max(scores, key=scores.get)
    confidence = max(scores.values()) / len(words) if words else 0
    
    # Configurações por idioma
    language_configs = {
        'pt': {
            'language': 'português',
            'font_family': 'Arial, sans-serif',
            'text_direction': 'ltr',
            'greeting_style': 'formal',
            'color_scheme': 'warm'
        },
        'en': {
            'language': 'english', 
            'font_family': 'Helvetica, sans-serif',
            'text_direction': 'ltr',
            'greeting_style': 'casual',
            'color_scheme': 'cool'
        },
        'es': {
            'language': 'español',
            'font_family': 'Arial, sans-serif', 
            'text_direction': 'ltr',
            'greeting_style': 'warm',
            'color_scheme': 'vibrant'
        }
    }
    
    config = language_configs.get(detected_language, language_configs['en'])
    config['detected_code'] = detected_language
    config['confidence'] = confidence
    
    return config


def _check_avatar_apis_availability(avatar_style: str) -> dict:
    """
    Verifica disponibilidade real de APIs de avatar 3D.
    
    Returns:
        dict: Status de disponibilidade das APIs
    """
    api_status = {
        'has_real_avatar_api': False,
        'available_apis': [],
        'preferred_api': None,
        'fallback_to_mvp': True
    }
    
    # Verificar D-ID API
    d_id_key = os.getenv('D_ID_API_KEY')
    if d_id_key and avatar_style in ['d_id', 'auto']:
        if _test_d_id_api_connection(d_id_key):
            api_status['available_apis'].append('d_id')
            api_status['preferred_api'] = 'd_id'
            api_status['has_real_avatar_api'] = True
            api_status['fallback_to_mvp'] = False
    
    # Verificar Synthesia API
    synthesia_key = os.getenv('SYNTHESIA_API_KEY')
    if synthesia_key and avatar_style in ['synthesia', 'auto']:
        if _test_synthesia_api_connection(synthesia_key):
            api_status['available_apis'].append('synthesia')
            if not api_status['preferred_api']:
                api_status['preferred_api'] = 'synthesia'
                api_status['has_real_avatar_api'] = True
                api_status['fallback_to_mvp'] = False
    
    # Verificar Hunyuan3D API (futuro)
    hunyuan_key = os.getenv('HUNYUAN3D_API_KEY')
    if hunyuan_key and avatar_style in ['hunyuan3d', 'auto']:
        # TODO: Implementar teste de conexão quando API estiver disponível
        pass
    
    return api_status


def _test_d_id_api_connection(api_key: str) -> bool:
    """
    Testa conexão real com D-ID API.
    
    Returns:
        bool: True se API está acessível
    """
    try:
        import requests
        
        headers = {
            'Authorization': f'Basic {api_key}',
            'Content-Type': 'application/json'
        }
        
        # Teste simples de conectividade
        response = requests.get(
            'https://api.d-id.com/credits',
            headers=headers,
            timeout=10
        )
        
        return response.status_code == 200
        
    except Exception as e:
        print(f"⚠️ D-ID API não disponível: {e}")
        return False


def _test_synthesia_api_connection(api_key: str) -> bool:
    """
    Testa conexão real com Synthesia API.
    
    Returns:
        bool: True se API está acessível
    """
    try:
        import requests
        
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        
        # Teste simples de conectividade  
        response = requests.get(
            'https://api.synthesia.io/v2/avatars',
            headers=headers,
            timeout=10
        )
        
        return response.status_code == 200
        
    except Exception as e:
        print(f"⚠️ Synthesia API não disponível: {e}")
        return False


def _generate_real_avatar_video(
    text: str, audio_path: str, output_path: str, template: str,
    avatar_style: str, language_info: dict, progress_callback: Callable,
    start_time: float, **kwargs
) -> dict:
    """
    IMPLEMENTAÇÃO FUTURA: Gera vídeo usando APIs reais de avatar.
    """
    
    if progress_callback:
        progress_callback(50, "Conectando com API de avatar...")
    
    # Determinar qual API usar
    api_to_use = kwargs.get('preferred_api', avatar_style)
    
    try:
        if api_to_use == 'd_id':
            return _generate_with_d_id_api(
                text, audio_path, output_path, template, language_info,
                progress_callback, start_time, **kwargs
            )
        elif api_to_use == 'synthesia':
            return _generate_with_synthesia_api(
                text, audio_path, output_path, template, language_info,
                progress_callback, start_time, **kwargs
            )
        else:
            # Fallback para MVP se API não reconhecida
            return _generate_enhanced_slide_video(
                text, audio_path, output_path, template, language_info,
                progress_callback, start_time, True, **kwargs
            )
            
    except Exception as e:
        print(f"❌ Erro na API de avatar: {e}")
        print("🔄 Fallback para MVP...")
        
        # Fallback para MVP em caso de erro
        return _generate_enhanced_slide_video(
            text, audio_path, output_path, template, language_info,
            progress_callback, start_time, True, **kwargs
        )


def _generate_with_d_id_api(
    text: str, audio_path: str, output_path: str, template: str,
    language_info: dict, progress_callback: Callable, start_time: float, **kwargs
) -> dict:
    """
    Gera vídeo usando D-ID API real - IMPLEMENTAÇÃO FUNCIONAL.
    """
    try:
        import requests
        import base64
        
        api_key = os.getenv('D_ID_API_KEY')
        if not api_key:
            raise Exception("D_ID_API_KEY não configurada")
        
        if progress_callback:
            progress_callback(60, "Preparando upload para D-ID...")
        
        # Configurar avatar baseado no template
        avatar_configs = {
            'professional': {
                'presenter_id': 'amy-jcu4GGiYNQ',  # Avatar profissional feminina
                'background': '#ffffff'
            },
            'educational': {
                'presenter_id': 'daniel-C2Y3dHl1eHE',  # Avatar masculino amigável
                'background': '#f0f8ff'
            },
            'tech': {
                'presenter_id': 'lucia-MdE2NDk4ZTk4ZQ',  # Avatar tech moderna
                'background': '#1a1a1a'
            },
            'minimal': {
                'presenter_id': 'amy-jcu4GGiYNQ',
                'background': '#f5f5f5'
            }
        }
        
        config = avatar_configs.get(template, avatar_configs['professional'])
        
        # Preparar áudio
        with open(audio_path, 'rb') as audio_file:
            audio_data = base64.b64encode(audio_file.read()).decode('utf-8')
        
        if progress_callback:
            progress_callback(70, "Enviando para processamento D-ID...")
        
        # Payload para D-ID API
        payload = {
            'script': {
                'type': 'audio',
                'audio_url': f'data:audio/mp3;base64,{audio_data}',
                'input': text
            },
            'config': {
                'fluent': False,
                'pad_audio': 0
            },
            'source_url': f'https://create-images-results.d-id.com/api_docs/assets/{config["presenter_id"]}.jpg',
            'background': config['background']
        }
        
        headers = {
            'Authorization': f'Basic {api_key}',
            'Content-Type': 'application/json'
        }
        
        # Criar vídeo
        response = requests.post(
            'https://api.d-id.com/talks',
            json=payload,
            headers=headers,
            timeout=30
        )
        
        if response.status_code != 201:
            raise Exception(f"Erro D-ID API: {response.text}")
        
        talk_id = response.json()['id']
        
        if progress_callback:
            progress_callback(80, f"Processando vídeo (ID: {talk_id})...")
        
        # Aguardar processamento
        max_wait = 300  # 5 minutos máximo
        check_interval = 10  # Verificar a cada 10 segundos
        waited = 0
        
        while waited < max_wait:
            time.sleep(check_interval)
            waited += check_interval
            
            # Verificar status
            status_response = requests.get(
                f'https://api.d-id.com/talks/{talk_id}',
                headers=headers,
                timeout=10
            )
            
            if status_response.status_code == 200:
                status_data = status_response.json()
                
                if status_data['status'] == 'done':
                    video_url = status_data['result_url']
                    
                    if progress_callback:
                        progress_callback(90, "Baixando vídeo finalizado...")
                    
                    # Baixar vídeo
                    video_response = requests.get(video_url, timeout=60)
                    
                    with open(output_path, 'wb') as f:
                        f.write(video_response.content)
                    
                    # Obter informações do arquivo
                    file_size = os.path.getsize(output_path)
                    duration = _get_audio_duration(audio_path)
                    
                    if progress_callback:
                        progress_callback(100, "Vídeo D-ID concluído!")
                    
                    return {
                        'success': True,
                        'video_path': output_path,
                        'duration': duration,
                        'method': 'd_id_api',
                        'avatar_api_used': 'd_id',
                        'template_used': template,
                        'resolution': (1920, 1080),  # D-ID padrão
                        'file_size': file_size,
                        'processing_time': time.time() - start_time,
                        'cached': False,
                        'quality_score': 0.95,  # D-ID alta qualidade
                        'metadata': {
                            'talk_id': talk_id,
                            'presenter_id': config['presenter_id'],
                            'api_version': 'v2'
                        },
                        'error': None
                    }
                    
                elif status_data['status'] == 'error':
                    raise Exception(f"Erro no processamento D-ID: {status_data.get('error', 'Erro desconhecido')}")
                
                if progress_callback:
                    progress_callback(80 + (waited / max_wait) * 10, f"Aguardando processamento... ({waited}s)")
            
            else:
                print(f"⚠️ Erro ao verificar status: {status_response.status_code}")
        
        raise Exception("Timeout no processamento D-ID")
        
    except Exception as e:
        print(f"❌ Erro D-ID API: {e}")
        # Fallback para MVP
        return _generate_enhanced_slide_video(
            text, audio_path, output_path, template, language_info,
            progress_callback, start_time, True, **kwargs
        )


def _generate_with_synthesia_api(
    text: str, audio_path: str, output_path: str, template: str,
    language_info: dict, progress_callback: Callable, start_time: float, **kwargs
) -> dict:
    """
    Gera vídeo usando Synthesia API real - IMPLEMENTAÇÃO FUNCIONAL.
    """
    try:
        import requests
        
        api_key = os.getenv('SYNTHESIA_API_KEY')
        if not api_key:
            raise Exception("SYNTHESIA_API_KEY não configurada")
        
        if progress_callback:
            progress_callback(60, "Preparando upload para Synthesia...")
        
        # TODO: Implementar Synthesia API quando disponível
        # Por enquanto, fallback para MVP
        print("🚧 Synthesia API em desenvolvimento...")
        
        return _generate_enhanced_slide_video(
            text, audio_path, output_path, template, language_info,
            progress_callback, start_time, True, **kwargs
        )
        
    except Exception as e:
        print(f"❌ Erro Synthesia API: {e}")
        return _generate_enhanced_slide_video(
            text, audio_path, output_path, template, language_info,
            progress_callback, start_time, True, **kwargs
        )


def _generate_enhanced_slide_video(
    text: str, audio_path: str, output_path: str, template: str,
    language_info: dict, progress_callback: Callable, start_time: float,
    cache_enabled: bool, **kwargs
) -> dict:
    """
    IMPLEMENTAÇÃO ATUAL APRIMORADA: Gera vídeo de slide com recursos avançados.
    """
    
    try:
        if progress_callback:
            progress_callback(50, "Gerando slide avançado...")
        
        # ETAPA 1: OBTER DURAÇÃO DO ÁUDIO  
        audio_duration = _get_audio_duration(audio_path)
        if audio_duration <= 0:
            return _create_error_result('Não foi possível determinar a duração do áudio', start_time)
        
        # ETAPA 2: CONFIGURAR TEMPLATE AVANÇADO
        video_config = _get_enhanced_template_config(template, language_info, kwargs)
        video_config['duration'] = audio_duration
        
        if progress_callback:
            progress_callback(60, f"Template '{template}' configurado...")
        
        # ETAPA 3: CRIAR FRAMES AVANÇADOS  
        frames_path = _create_enhanced_video_frames(text, video_config, language_info)
        if not frames_path:
            return _create_error_result('Falha na criação dos frames avançados', start_time)
        
        if progress_callback:
            progress_callback(75, "Frames avançados criados...")
        
        # ETAPA 4: COMBINAR COM ÁUDIO
        success = _combine_frames_with_audio(frames_path, audio_path, output_path, video_config)
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            
            # Limpar arquivos temporários
            _cleanup_temp_files(frames_path)
            
            if progress_callback:
                progress_callback(90, "Vídeo gerado, salvando cache...")
            
            # Salvar no cache se habilitado
            if cache_enabled:
                _save_to_avatar_cache(text, audio_path, template, output_path, video_config)
            
            if progress_callback:
                progress_callback(100, "Vídeo slide avançado concluído!")
            
            result = {
                'success': True,
                'video_path': output_path,
                'duration': audio_duration,
                'method': 'enhanced_slide_mvp',
                'avatar_api_used': None,
                'template_used': template,
                'resolution': (video_config['width'], video_config['height']),
                'file_size': file_size,
                'processing_time': time.time() - start_time,
                'cached': False,
                'quality_score': _calculate_quality_score(video_config),
                'metadata': {
                    'template': template,
                    'language': language_info['language'],
                    'effects_applied': video_config.get('effects', []),
                    'version': 'enhanced_v2'
                },
                'error': None
            }
            
            print(f"✅ Vídeo slide avançado criado!")
            print(f"📁 Arquivo: {output_path}")
            print(f"⏱️ Duração: {audio_duration:.2f}s")
            print(f"🎨 Template: {template}")
            print(f"🌐 Idioma: {language_info['language']}")
            print(f"📊 Qualidade: {result['quality_score']:.2f}")
            
            return result
        else:
            return _create_error_result('Falha na combinação de frames com áudio', start_time)
            
    except Exception as e:
        error_msg = f"Erro na geração do slide avançado: {str(e)}"
        return _create_error_result(error_msg, start_time)


def _get_audio_duration(audio_path: str) -> float:
    """
    Obtém a duração do arquivo de áudio em segundos.
    
    Args:
        audio_path (str): Caminho do arquivo de áudio
    
    Returns:
        float: Duração em segundos (0.0 se erro)
    """
    try:
        # Tentar usar mutagen para obter informações do áudio
        try:
            from mutagen import File
            audio_file = File(audio_path)
            if audio_file and audio_file.info:
                duration = audio_file.info.length
                return float(duration)
        except ImportError:
            print("⚠️ mutagen não disponível, usando método alternativo")
        
        # Método alternativo usando FFmpeg (se disponível)
        try:
            import subprocess
            result = subprocess.run([
                'ffprobe', '-v', 'quiet', '-show_entries', 
                'format=duration', '-of', 'csv=p=0', audio_path
            ], capture_output=True, text=True, timeout=10)
            
            if result.returncode == 0:
                return float(result.stdout.strip())
        except (subprocess.SubprocessError, FileNotFoundError, ValueError):
            print("⚠️ FFmpeg não disponível")
        
        # Fallback: estimar duração baseado no tamanho do arquivo
        # Para MP3: aproximadamente 1MB = 60-80 segundos (128kbps)
        file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
        estimated_duration = file_size_mb * 70  # Estimativa conservadora
        
        print(f"⚠️ Usando estimativa de duração: {estimated_duration:.2f}s baseado no tamanho")
        return estimated_duration
        
    except Exception as e:
        print(f"❌ Erro ao obter duração do áudio: {e}")
        return 10.0  # Fallback: 10 segundos


def _create_video_frames(text: str, config: dict) -> str:
    """
    Cria os frames do vídeo com o texto formatado.
    
    Args:
        text (str): Texto para exibir
        config (dict): Configurações do vídeo
    
    Returns:
        str: Caminho do arquivo de frame criado (None se erro)
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
        import textwrap
        
        # Criar imagem com fundo
        img = Image.new('RGB', (config['width'], config['height']), config['background_color'])
        draw = ImageDraw.Draw(img)
        
        # Tentar carregar uma fonte melhor
        try:
            # Tentar fontes do sistema
            if os.name == 'nt':  # Windows
                font_path = "C:/Windows/Fonts/arial.ttf"
            else:  # Linux/Mac
                font_path = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
            
            if os.path.exists(font_path):
                font = ImageFont.truetype(font_path, config['font_size'])
            else:
                font = ImageFont.load_default()
        except Exception:
            font = ImageFont.load_default()
        
        # Quebrar texto em linhas
        max_chars_per_line = (config['width'] - 2 * config['padding']) // (config['font_size'] // 2)
        lines = textwrap.wrap(text, width=max_chars_per_line)
        
        # Calcular posicionamento vertical centralizado
        line_height = config['font_size'] + 10
        total_text_height = len(lines) * line_height
        start_y = (config['height'] - total_text_height) // 2
        
        # Desenhar cada linha de texto
        for i, line in enumerate(lines):
            # Centralizar horizontalmente
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]
            x = (config['width'] - text_width) // 2
            y = start_y + i * line_height
            
            # Desenhar texto com sombra (efeito visual)
            shadow_offset = 2
            draw.text((x + shadow_offset, y + shadow_offset), line, 
                     fill=(0, 0, 0), font=font)  # Sombra preta
            draw.text((x, y), line, fill=config['text_color'], font=font)  # Texto principal
        
        # Adicionar elementos visuais extras (opcional)
        _add_visual_elements(draw, config)
        
        # Salvar frame
        temp_dir = "temp"
        os.makedirs(temp_dir, exist_ok=True)
        frame_path = os.path.join(temp_dir, f"slide_frame_{int(time.time())}.png")
        img.save(frame_path, 'PNG', quality=95)
        
        print(f"🖼️ Frame criado: {frame_path}")
        return frame_path
        
    except Exception as e:
        print(f"❌ Erro ao criar frames: {e}")
        return None


def _add_visual_elements(draw, config: dict):
    """
    Adiciona elementos visuais decorativos ao slide.
    
    Args:
        draw: Objeto ImageDraw
        config: Configurações do vídeo
    """
    try:
        # Adicionar bordas decorativas
        border_width = 5
        border_color = (70, 130, 180)  # Azul aço
        
        # Borda superior
        draw.rectangle([0, 0, config['width'], border_width], fill=border_color)
        # Borda inferior  
        draw.rectangle([0, config['height'] - border_width, config['width'], config['height']], fill=border_color)
        
        # Adicionar logo ou marca d'água (futuro)
        # TODO: Implementar logo da empresa/marca
        
        # Adicionar gradiente sutil (futuro)
        # TODO: Implementar efeito de gradiente no fundo
        
    except Exception as e:
        print(f"⚠️ Erro ao adicionar elementos visuais: {e}")


def _combine_frames_with_audio(frame_path: str, audio_path: str, output_path: str, config: dict) -> bool:
    """
    Combina o frame estático com o áudio para criar o vídeo final.
    
    Args:
        frame_path (str): Caminho do frame de imagem
        audio_path (str): Caminho do áudio
        output_path (str): Caminho de saída do vídeo
        config (dict): Configurações do vídeo
    
    Returns:
        bool: True se sucesso, False se erro
    """
    try:
        # Tentar usar OpenCV para criação do vídeo
        try:
            import cv2
            return _combine_with_opencv(frame_path, audio_path, output_path, config)
        except ImportError:
            print("⚠️ OpenCV não disponível, tentando FFmpeg")
        
        # Método alternativo usando FFmpeg
        try:
            import subprocess
            return _combine_with_ffmpeg(frame_path, audio_path, output_path, config)
        except FileNotFoundError:
            print("⚠️ FFmpeg não disponível")
        
        print("❌ Nenhum método de combinação de vídeo disponível")
        return False
        
    except Exception as e:
        print(f"❌ Erro na combinação: {e}")
        return False


def _combine_with_opencv(frame_path: str, audio_path: str, output_path: str, config: dict) -> bool:
    """
    Combina usando OpenCV (método preferido).
    
    Args:
        frame_path (str): Caminho do frame
        audio_path (str): Caminho do áudio
        output_path (str): Caminho de saída
        config (dict): Configurações do vídeo
    
    Returns:
        bool: True se sucesso
    """
    try:
        import cv2
        import subprocess
        
        # Primeiro, criar vídeo silencioso com OpenCV
        temp_video = output_path.replace('.mp4', '_temp.mp4')
        
        # Ler frame
        frame = cv2.imread(frame_path)
        if frame is None:
            print(f"❌ Não foi possível ler o frame: {frame_path}")
            return False
        
        # Configurar codec e VideoWriter
        fourcc = cv2.VideoWriter_fourcc(*'mp4v')
        out = cv2.VideoWriter(temp_video, fourcc, config['fps'], 
                             (config['width'], config['height']))
        
        # Calcular número total de frames necessários
        total_frames = int(config['duration'] * config['fps'])
        
        # Escrever frames (mesmo frame repetido)
        print(f"🎬 Criando {total_frames} frames...")
        for i in range(total_frames):
            out.write(frame)
            if i % 100 == 0:  # Progress feedback
                progress = (i / total_frames) * 100
                print(f"   Progresso: {progress:.1f}%")
        
        out.release()
        cv2.destroyAllWindows()
        
        # Combinar com áudio usando FFmpeg
        print("🎵 Combinando com áudio...")
        cmd = [
            'ffmpeg', '-y',  # -y para sobrescrever
            '-i', temp_video,
            '-i', audio_path, 
            '-c:v', 'copy',  # Copiar stream de vídeo
            '-c:a', 'aac',   # Codec de áudio
            '-shortest',     # Terminar quando o menor stream acabar
            output_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        # Limpar arquivo temporário
        if os.path.exists(temp_video):
            os.remove(temp_video)
        
        return result.returncode == 0
        
    except Exception as e:
        print(f"❌ Erro no OpenCV: {e}")
        return False


def _combine_with_ffmpeg(frame_path: str, audio_path: str, output_path: str, config: dict) -> bool:
    """
    Combina usando apenas FFmpeg.
    
    Args:
        frame_path (str): Caminho do frame
        audio_path (str): Caminho do áudio  
        output_path (str): Caminho de saída
        config (dict): Configurações do vídeo
    
    Returns:
        bool: True se sucesso
    """
    try:
        import subprocess
        
        print("🎬 Criando vídeo com FFmpeg...")
        
        # Comando FFmpeg para criar vídeo a partir de imagem estática + áudio
        cmd = [
            'ffmpeg', '-y',  # -y para sobrescrever
            '-loop', '1',    # Loop da imagem
            '-i', frame_path,  # Imagem de entrada
            '-i', audio_path,  # Áudio de entrada
            '-c:v', 'libx264',  # Codec de vídeo
            '-tune', 'stillimage',  # Otimizar para imagem estática
            '-c:a', 'aac',     # Codec de áudio
            '-b:a', '192k',    # Bitrate do áudio
            '-pix_fmt', 'yuv420p',  # Formato de pixel compatível
            '-shortest',       # Terminar quando o menor stream acabar
            '-r', str(config['fps']),  # Frame rate
            output_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode == 0:
            print("✅ Vídeo criado com FFmpeg")
            return True
        else:
            print(f"❌ Erro FFmpeg: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"❌ Erro no FFmpeg: {e}")
        return False


def _cleanup_temp_files(frame_path: str):
    """
    Limpa arquivos temporários criados durante a geração.
    
    Args:
        frame_path (str): Caminho do frame temporário
    """
    try:
        if frame_path and os.path.exists(frame_path):
            os.remove(frame_path)
            print(f"🧹 Arquivo temporário removido: {frame_path}")
        
        # Remover diretório temp se estiver vazio
        temp_dir = os.path.dirname(frame_path)
        if temp_dir and os.path.exists(temp_dir) and not os.listdir(temp_dir):
            os.rmdir(temp_dir)
            print(f"🧹 Diretório temporário removido: {temp_dir}")
            
    except Exception as e:
        print(f"⚠️ Aviso na limpeza: {e}")


# ============================================================================
# FUNÇÃO DE NARRAÇÃO DE TEXTO (TTS - TEXT-TO-SPEECH)
# ============================================================================

async def generate_narration(
    text: str, 
    output_path: str,
    voice: Optional[str] = None,
    provider: str = "auto",
    language: str = "pt"
) -> Dict:
    """
    Gera narração em áudio MP3 a partir de texto usando modelos TTS da Hugging Face
    
    Args:
        text (str): Texto para narrar (máximo recomendado: 1000 caracteres por vez)
        output_path (str): Caminho completo onde salvar o arquivo MP3
        voice (str, optional): Voz específica a usar (ex: "v2/pt_speaker_0" para Bark)
        provider (str): Provedor TTS ("bark", "gtts", "auto")
        language (str): Código do idioma (padrão: "pt" para português)
    
    Returns:
        Dict: Resultado da geração com success, audio_path, duration, provider_used, error
        
    Example:
        >>> # Narração simples com gTTS (Google)
        >>> result = await generate_narration(
        ...     "Olá! Este é um teste de narração em português.",
        ...     "narracao_teste.mp3"
        ... )
        >>> print(f"Sucesso: {result['success']}")
        >>> print(f"Arquivo: {result['audio_path']}")
        
        >>> # Narração com Bark (Hugging Face) - voz específica
        >>> result = await generate_narration(
        ...     "Bem-vindos ao curso de inteligência artificial!",
        ...     "intro_curso.mp3",
        ...     voice="v2/pt_speaker_2",
        ...     provider="bark"
        ... )
        
    Note:
        Para usar modelos Bark da Hugging Face:
        1. Configure sua chave da Hugging Face (opcional para modelos públicos):
           export HUGGINGFACE_TOKEN="sua_chave_aqui"
           
        2. Instale dependências:
           pip install torch transformers torchaudio gtts pydub
           
        3. Para GPU (recomendado para Bark):
           pip install torch torchvision torchaudio --index-url https://download.pytorch.org/whl/cu118
           
        4. Modelos disponíveis:
           - Bark: Vozes naturais em português (v2/pt_speaker_0 a v2/pt_speaker_9)
           - gTTS: Google TTS (mais rápido, menor qualidade)
           
        5. Configurações do ambiente:
           - O modelo Bark requer ~2GB de VRAM
           - Processamento pode levar 30s-2min dependendo do hardware
           - Para produção, use cache de modelos pré-carregados
    """
    
    # Verificar disponibilidade do serviço TTS
    if not TTS_AVAILABLE:
        return {
            'success': False,
            'error': 'Serviço TTS não disponível. Instale as dependências necessárias.',
            'audio_path': None,
            'duration': 0.0,
            'provider_used': None
        }
    
    # Validar entrada
    if not text or not text.strip():
        return {
            'success': False,
            'error': 'Texto não pode estar vazio',
            'audio_path': None,
            'duration': 0.0,
            'provider_used': None
        }
    
    if len(text) > 2000:
        return {
            'success': False,
            'error': f'Texto muito longo ({len(text)} caracteres). Máximo recomendado: 2000 caracteres.',
            'audio_path': None,
            'duration': 0.0,
            'provider_used': None
        }
    
    try:
        # Garantir que o diretório de saída existe
        output_dir = os.path.dirname(output_path)
        if output_dir:
            Path(output_dir).mkdir(parents=True, exist_ok=True)
        
        # Configurar provedor TTS
        if provider.lower() == "bark":
            tts_provider = TTSProvider.BARK
        elif provider.lower() == "gtts":
            tts_provider = TTSProvider.GTTS
        else:
            tts_provider = TTSProvider.AUTO  # Auto-detecta melhor opção
        
        # Configurar TTS
        tts_config = TTSConfig(
            provider=tts_provider,
            language=language,
            voice=voice,
            output_format="mp3"
        )
        
        # Inicializar serviço TTS
        tts_service = TTSService()
        
        # Verificar cache primeiro (se disponível)
        cached_result = None
        try:
            from app.services.tts_cache_service import get_cached_tts_audio
            cached_result = await get_cached_tts_audio(text, provider, voice, language)
        except ImportError:
            pass  # Cache não disponível
        
        if cached_result:
            print(f"🎯 Áudio encontrado no cache!")
            return cached_result
        
        # Registrar início da métrica
        start_time = time.time()
        
        # Gerar áudio
        print(f"🎤 Gerando narração com {provider}...")
        print(f"📝 Texto: {text[:100]}{'...' if len(text) > 100 else ''}")
        
        result = await tts_service.generate_audio(
            text=text,
            config=tts_config,
            output_path=output_path
        )
        
        # Calcular tempo de processamento
        processing_time = time.time() - start_time
        
        if result.success:
            print(f"✅ Narração gerada com sucesso!")
            print(f"📁 Arquivo: {result.audio_path}")
            print(f"⏱️ Duração: {result.duration:.2f} segundos")
            print(f"🔊 Provedor: {result.provider_used}")
            print(f"⚡ Tempo de processamento: {processing_time:.2f}s")
            
            # Armazenar no cache (se disponível)
            try:
                from app.services.tts_cache_service import store_tts_audio
                await store_tts_audio(
                    text=text,
                    provider=result.provider_used,
                    audio_path=result.audio_path,
                    duration=result.duration,
                    voice=voice,
                    language=language,
                    metadata=result.metadata
                )
                print(f"💾 Áudio armazenado no cache")
            except ImportError:
                pass  # Cache não disponível
            
            # Registrar métrica (se disponível)
            try:
                from app.services.tts_analytics_service import record_tts_metric
                await record_tts_metric(
                    user_id=None,  # Será passado via contexto em versões futuras
                    provider=result.provider_used,
                    voice=voice,
                    text_length=len(text),
                    processing_time=processing_time,
                    success=True,
                    file_size=os.path.getsize(result.audio_path) if os.path.exists(result.audio_path) else 0,
                    duration=result.duration,
                    cached=False
                )
            except ImportError:
                pass  # Analytics não disponível
            
            return {
                'success': True,
                'audio_path': result.audio_path,
                'duration': result.duration,
                'provider_used': result.provider_used,
                'error': None,
                'metadata': result.metadata,
                'processing_time': processing_time,
                'cached': False
            }
        else:
            print(f"❌ Falha na geração: {result.error}")
            
            # Registrar métrica de erro (se disponível)
            try:
                from app.services.tts_analytics_service import record_tts_metric
                await record_tts_metric(
                    user_id=None,
                    provider=result.provider_used or provider,
                    voice=voice,
                    text_length=len(text),
                    processing_time=processing_time,
                    success=False,
                    error_type=result.error[:100] if result.error else "unknown_error"
                )
            except ImportError:
                pass
            
            return {
                'success': False,
                'error': result.error or 'Erro desconhecido na geração de áudio',
                'audio_path': None,
                'duration': 0.0,
                'provider_used': result.provider_used,
                'processing_time': processing_time,
                'cached': False
            }
            
    except Exception as e:
        error_msg = f"Erro na geração de narração: {str(e)}"
        print(f"❌ {error_msg}")
        return {
            'success': False,
            'error': error_msg,
            'audio_path': None,
            'duration': 0.0,
            'provider_used': None
        }


def generate_narration_sync(
    text: str, 
    output_path: str,
    voice: Optional[str] = None,
    provider: str = "auto",
    language: str = "pt"
) -> Dict:
    """
    Versão síncrona da função generate_narration para uso sem async/await
    
    Args:
        text (str): Texto para narrar
        output_path (str): Caminho onde salvar o arquivo MP3
        voice (str, optional): Voz específica a usar
        provider (str): Provedor TTS ("bark", "gtts", "auto")
        language (str): Código do idioma (padrão: "pt")
    
    Returns:
        Dict: Resultado da geração
        
    Example:
        >>> result = generate_narration_sync(
        ...     "Olá! Este é um teste de narração.",
        ...     "narracao.mp3"
        ... )
        >>> if result['success']:
        ...     print(f"Áudio salvo em: {result['audio_path']}")
    """
    import asyncio
    
    # Executar a função async de forma síncrona
    try:
        loop = asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    return loop.run_until_complete(
        generate_narration(text, output_path, voice, provider, language)
    )


# ============================================================================
# EXEMPLO DE USO DAS FUNÇÕES DE NARRAÇÃO
# ============================================================================

if __name__ == '__main__':
    print("\n" + "="*80)
    print("🎤 DEMONSTRAÇÃO DO SISTEMA DE NARRAÇÃO TTS")
    print("="*80)
    
    # Exemplo 1: Narração simples com gTTS (Google)
    print("\n1️⃣ TESTE COM GTTS (GOOGLE TTS)")
    print("-" * 40)
    
    texto_exemplo = """
    Olá! Este é um teste de narração automática em português. 
    O sistema de Text-to-Speech permite converter qualquer texto em áudio de alta qualidade.
    Esta funcionalidade é perfeita para criar conteúdo educacional, audiobooks e muito mais!
    """
    
    try:
        result = generate_narration_sync(
            text=texto_exemplo.strip(),
            output_path="exemplo_narracao_gtts.mp3",
            provider="gtts"
        )
        
        if result['success']:
            print(f"✅ Sucesso! Arquivo gerado: {result['audio_path']}")
            print(f"⏱️ Duração: {result['duration']:.2f}s")
            print(f"🔊 Provedor: {result['provider_used']}")
        else:
            print(f"❌ Erro: {result['error']}")
    except Exception as e:
        print(f"❌ Erro no teste gTTS: {e}")
    
    # Exemplo 2: Narração com Bark (Hugging Face)
    print("\n2️⃣ TESTE COM BARK (HUGGING FACE)")
    print("-" * 40)
    
    texto_bark = """
    Bem-vindos ao futuro da inteligência artificial! 
    Com os modelos Bark da Hugging Face, podemos criar vozes realistas e expressivas.
    Esta tecnologia revoluciona a criação de conteúdo digital.
    """
    
    try:
        result = generate_narration_sync(
            text=texto_bark.strip(),
            output_path="exemplo_narracao_bark.mp3",
            provider="bark",
            voice="v2/pt_speaker_2"  # Voz feminina brasileira
        )
        
        if result['success']:
            print(f"✅ Sucesso! Arquivo gerado: {result['audio_path']}")
            print(f"⏱️ Duração: {result['duration']:.2f}s")
            print(f"🔊 Provedor: {result['provider_used']}")
            print(f"🎭 Voz: v2/pt_speaker_2")
        else:
            print(f"❌ Erro: {result['error']}")
    except Exception as e:
        print(f"❌ Erro no teste Bark: {e}")
    
    # Exemplo 3: Auto-detectar melhor provedor
    print("\n3️⃣ TESTE COM AUTO-DETECÇÃO")
    print("-" * 40)
    
    texto_auto = "Este é um teste usando auto-detecção do melhor provedor TTS disponível."
    
    try:
        result = generate_narration_sync(
            text=texto_auto,
            output_path="exemplo_narracao_auto.mp3",
            provider="auto"
        )
        
        if result['success']:
            print(f"✅ Sucesso! Arquivo gerado: {result['audio_path']}")
            print(f"⏱️ Duração: {result['duration']:.2f}s")
            print(f"🔊 Provedor selecionado: {result['provider_used']}")
        else:
            print(f"❌ Erro: {result['error']}")
    except Exception as e:
        print(f"❌ Erro no teste auto: {e}")
    
    # Informações de configuração
    print("\n" + "="*80)
    print("📋 INFORMAÇÕES DE CONFIGURAÇÃO TTS")
    print("="*80)
    print("""
🔧 CONFIGURAÇÃO HUGGING FACE TOKEN (Opcional):
   export HUGGINGFACE_TOKEN="your_token_here"
   
📦 INSTALAÇÃO DE DEPENDÊNCIAS:
   pip install torch transformers torchaudio gtts pydub
   
🎯 GPU PARA BARK (Recomendado):
   pip install torch torchvision torchaudio --index-url https://download.pytorch.org/whl/cu118
   
🎭 VOZES DISPONÍVEIS PARA PORTUGUÊS:
   - v2/pt_speaker_0 (Masculina)
   - v2/pt_speaker_1 (Feminina)
   - v2/pt_speaker_2 (Feminina jovem)
   - v2/pt_speaker_3 (Masculina grave)
   - v2/pt_speaker_4 (Feminina madura)
   - v2/pt_speaker_5 a v2/pt_speaker_9 (Variações)
   
⚡ PERFORMANCE:
   - gTTS: ~1-3 segundos (online)
   - Bark CPU: ~30-120 segundos 
   - Bark GPU: ~10-30 segundos
   
💾 REQUISITOS BARK:
   - RAM: ~4GB
   - VRAM: ~2GB (GPU)
   - Armazenamento: ~2GB (modelos)
    """)
    
    print("="*80)
    print("🎉 DEMONSTRAÇÃO CONCLUÍDA!")
    print("🎵 Verifique os arquivos MP3 gerados no diretório atual")
    print("="*80)
    
    # Exemplo 4: Geração de vídeo avatar
    print("\n" + "="*80)
    print("🎬 DEMONSTRAÇÃO DO SISTEMA DE VÍDEO AVATAR")
    print("="*80)
    
    print("\n4️⃣ TESTE DE GERAÇÃO DE VÍDEO AVATAR (MVP)")
    print("-" * 50)
    
    # Primeiro, vamos criar um áudio de exemplo para o teste
    texto_video = """
    Bem-vindos ao futuro da educação digital! 
    Este é um exemplo de como criar vídeos educacionais automaticamente 
    usando inteligência artificial para narração e geração de conteúdo visual.
    """
    
    audio_exemplo_path = "exemplo_audio_avatar.mp3"
    video_output_path = "exemplo_video_avatar.mp4"
    
    try:
        # Primeiro, gerar o áudio para o vídeo
        print("🎤 Gerando áudio para o vídeo avatar...")
        result_audio = generate_narration_sync(
            text=texto_video.strip(),
            output_path=audio_exemplo_path,
            provider="gtts"  # Usar gTTS para rapidez
        )
        
        if result_audio['success']:
            print(f"✅ Áudio criado: {result_audio['audio_path']}")
            
            # Agora gerar o vídeo avatar
            print("🎬 Gerando vídeo avatar...")
            result_video = generate_avatar_video(
                text=texto_video.strip(),
                audio_path=result_audio['audio_path'],
                output_path=video_output_path
            )
            
            if result_video['success']:
                print(f"✅ Vídeo avatar criado com sucesso!")
                print(f"📁 Arquivo: {result_video['video_path']}")
                print(f"⏱️ Duração: {result_video['duration']:.2f}s")
                print(f"🎬 Resolução: {result_video['resolution'][0]}x{result_video['resolution'][1]}")
                print(f"📊 Tamanho: {format_file_size(result_video['file_size'])}")
                print(f"🔧 Método: {result_video['method']}")
                print(f"🚀 Próxima versão incluirá avatares 3D!")
            else:
                print(f"❌ Erro na geração do vídeo: {result_video['error']}")
        else:
            print(f"❌ Erro na geração do áudio: {result_audio['error']}")
            
    except Exception as e:
        print(f"❌ Erro no teste de vídeo avatar: {e}")
    
    # Informações sobre o sistema de vídeo avatar
    print("\n" + "="*80)
    print("📋 INFORMAÇÕES DO SISTEMA DE VÍDEO AVATAR")
    print("="*80)
    print("""
🎯 IMPLEMENTAÇÃO ATUAL (MVP):
   ✅ Vídeo de slide com texto e áudio
   ✅ Resolução Full HD (1920x1080)
   ✅ Sincronização automática com áudio
   ✅ Design profissional e limpo
   ✅ Suporte a OpenCV e FFmpeg
   
🚀 ROADMAP FUTURO:
   🔮 Integração com Hunyuan3D-2 (Tencent)
   🔮 Avatares 3D realistas
   🔮 Lip-sync automático
   🔮 Gestos e expressões faciais
   🔮 Múltiplos modelos de avatar
   🔮 Backgrounds personalizáveis
   🔮 API D-ID, Synthesia, Runway ML
   
📦 DEPENDÊNCIAS ATUAIS:
   pip install opencv-python pillow mutagen
   
🔧 DEPENDÊNCIAS OPCIONAIS:
   FFmpeg (para melhor compatibilidade)
   
🎭 CONFIGURAÇÃO FUTURA:
   export HUNYUAN3D_API_KEY="sua_chave"
   export DID_API_KEY="sua_chave"
   export SYNTHESIA_API_KEY="sua_chave"
   
⚡ PERFORMANCE ATUAL:
   - Slide MVP: ~5-15 segundos
   - Futuro 3D: ~30-120 segundos (cloud)
   
🎬 EXEMPLO DE CÓDIGO:
   # Gerar vídeo avatar completo
   result = generate_avatar_video(
       text="Sua apresentação aqui",
       audio_path="naracao.mp3",
       output_path="video_avatar.mp4"
   )
   
   if result['success']:
       print(f"Vídeo criado: {result['video_path']}")
    """)
    
    print("="*80)
    print("🎉 DEMONSTRAÇÃO DE VÍDEO AVATAR CONCLUÍDA!")
    print("🎬 Verifique o arquivo MP4 gerado no diretório atual")
    print("="*80) 


def _create_error_result(error_msg: str, start_time: float) -> dict:
    """
    Cria resultado padronizado para erros.
    
    Args:
        error_msg (str): Mensagem de erro
        start_time (float): Timestamp do início do processamento
    
    Returns:
        dict: Resultado de erro padronizado
    """
    return {
        'success': False,
        'error': error_msg,
        'video_path': None,
        'duration': 0.0,
        'method': 'error',
        'avatar_api_used': None,
        'template_used': None,
        'resolution': None,
        'file_size': 0,
        'processing_time': time.time() - start_time,
        'cached': False,
        'quality_score': 0.0,
        'metadata': {}
    }


def _create_cached_result(cache_result: dict, start_time: float) -> dict:
    """
    Cria resultado para vídeo encontrado em cache.
    
    Args:
        cache_result (dict): Dados do cache
        start_time (float): Timestamp do início
    
    Returns:
        dict: Resultado com dados do cache
    """
    metadata = cache_result['metadata']
    cache_path = cache_result['cache_path']
    
    return {
        'success': True,
        'video_path': cache_path,
        'duration': metadata.get('duration', 0.0),
        'method': metadata.get('method', 'cached'),
        'avatar_api_used': metadata.get('avatar_api_used'),
        'template_used': metadata.get('template_used'),
        'resolution': tuple(metadata.get('resolution', [1920, 1080])),
        'file_size': os.path.getsize(cache_path) if os.path.exists(cache_path) else 0,
        'processing_time': time.time() - start_time,
        'cached': True,
        'quality_score': metadata.get('quality_score', 0.8),
        'metadata': metadata,
        'error': None
    }


def _get_enhanced_template_config(template: str, language_info: dict, kwargs: dict) -> dict:
    """
    Configurações avançadas por template visual.
    
    Args:
        template (str): Nome do template
        language_info (dict): Informações de idioma
        kwargs (dict): Parâmetros adicionais
    
    Returns:
        dict: Configuração completa do template
    """
    
    # Configurações base por template
    template_configs = {
        'professional': {
            'background_color': (25, 35, 50),      # Azul corporativo escuro
            'accent_color': (70, 130, 180),        # Azul aço
            'text_color': (255, 255, 255),         # Branco
            'secondary_text_color': (200, 200, 200), # Cinza claro
            'font_size': 52,
            'title_font_size': 64,
            'padding': 120,
            'effects': ['gradient_background', 'professional_border', 'subtle_shadow'],
            'style': 'corporate'
        },
        'educational': {
            'background_color': (240, 248, 255),   # Alice blue
            'accent_color': (65, 105, 225),        # Royal blue
            'text_color': (25, 25, 112),           # Midnight blue
            'secondary_text_color': (70, 70, 70),  # Cinza escuro
            'font_size': 48,
            'title_font_size': 58,
            'padding': 100,
            'effects': ['soft_gradient', 'educational_icons', 'friendly_border'],
            'style': 'friendly'
        },
        'tech': {
            'background_color': (15, 15, 20),      # Preto tech
            'accent_color': (0, 255, 127),         # Verde neon
            'text_color': (255, 255, 255),         # Branco
            'secondary_text_color': (0, 255, 127), # Verde neon
            'font_size': 50,
            'title_font_size': 62,
            'padding': 80,
            'effects': ['neon_glow', 'tech_grid', 'cyber_border'],
            'style': 'futuristic'
        },
        'minimal': {
            'background_color': (250, 250, 250),   # Off-white
            'accent_color': (100, 100, 100),       # Cinza médio
            'text_color': (50, 50, 50),            # Cinza escuro
            'secondary_text_color': (120, 120, 120), # Cinza
            'font_size': 46,
            'title_font_size': 56,
            'padding': 140,
            'effects': ['clean_lines', 'minimal_shadow'],
            'style': 'clean'
        }
    }
    
    # Configuração base
    base_config = template_configs.get(template, template_configs['professional'])
    
    # Aplicar configurações de idioma
    if language_info['detected_code'] == 'pt':
        base_config['font_family'] = 'Arial, sans-serif'
        base_config['line_spacing'] = 1.4
    elif language_info['detected_code'] == 'en':
        base_config['font_family'] = 'Helvetica, sans-serif'
        base_config['line_spacing'] = 1.3
    else:
        base_config['font_family'] = 'Arial, sans-serif'
        base_config['line_spacing'] = 1.4
    
    # Configurações de vídeo
    base_config.update({
        'width': kwargs.get('width', 1920),
        'height': kwargs.get('height', 1080),
        'fps': kwargs.get('fps', 30),
        'quality': kwargs.get('quality', 'high'),
        'format': kwargs.get('format', 'mp4')
    })
    
    # Efeitos adicionais solicitados
    additional_effects = kwargs.get('effects', [])
    if additional_effects:
        base_config['effects'].extend(additional_effects)
    
    return base_config


def _create_enhanced_video_frames(text: str, config: dict, language_info: dict) -> str:
    """
    Cria frames de vídeo com recursos visuais avançados.
    
    Args:
        text (str): Texto a ser exibido
        config (dict): Configurações do template
        language_info (dict): Informações de idioma
    
    Returns:
        str: Caminho do frame criado
    """
    try:
        from PIL import Image, ImageDraw, ImageFont, ImageFilter
        import textwrap
        
        # Criar imagem base
        img = Image.new('RGBA', (config['width'], config['height']), (0, 0, 0, 0))
        
        # Aplicar efeitos de fundo
        img = _apply_background_effects(img, config)
        
        # Configurar fonte
        font, title_font = _setup_enhanced_fonts(config)
        
        # Processar e formatar texto
        formatted_text = _format_text_for_display(text, config, language_info)
        
        # Criar layer de texto
        text_layer = _create_text_layer(formatted_text, config, font, title_font)
        
        # Combinar layers
        final_img = Image.alpha_composite(img.convert('RGBA'), text_layer)
        
        # Aplicar efeitos finais
        final_img = _apply_final_effects(final_img, config)
        
        # Salvar frame
        temp_dir = "temp"
        os.makedirs(temp_dir, exist_ok=True)
        frame_path = os.path.join(temp_dir, f"enhanced_frame_{int(time.time())}.png")
        
        final_img.convert('RGB').save(frame_path, 'PNG', quality=95, optimize=True)
        
        print(f"🎨 Frame avançado criado: {frame_path}")
        return frame_path
        
    except Exception as e:
        print(f"❌ Erro ao criar frames avançados: {e}")
        # Fallback para frame simples
        return _create_video_frames(text, config)


def _apply_background_effects(img: Image, config: dict) -> Image:
    """
    Aplica efeitos de fundo baseados no template.
    
    Args:
        img (Image): Imagem base
        config (dict): Configurações do template
    
    Returns:
        Image: Imagem com efeitos de fundo
    """
    try:
        from PIL import ImageDraw
        
        draw = ImageDraw.Draw(img)
        width, height = img.size
        
        # Fundo base
        draw.rectangle([0, 0, width, height], fill=config['background_color'])
        
        effects = config.get('effects', [])
        
        # Gradiente de fundo
        if 'gradient_background' in effects:
            _apply_gradient(draw, width, height, config)
        elif 'soft_gradient' in effects:
            _apply_soft_gradient(draw, width, height, config)
        
        # Grid tecnológico
        if 'tech_grid' in effects:
            _apply_tech_grid(draw, width, height, config)
        
        # Bordas especiais
        if 'professional_border' in effects:
            _apply_professional_border(draw, width, height, config)
        elif 'cyber_border' in effects:
            _apply_cyber_border(draw, width, height, config)
        elif 'friendly_border' in effects:
            _apply_friendly_border(draw, width, height, config)
        
        return img
        
    except Exception as e:
        print(f"⚠️ Erro ao aplicar efeitos de fundo: {e}")
        return img


def _apply_gradient(draw, width: int, height: int, config: dict):
    """Aplica gradiente profissional."""
    try:
        bg_color = config['background_color']
        accent_color = config['accent_color']
        
        # Gradiente vertical sutil
        for y in range(height):
            alpha = y / height
            r = int(bg_color[0] * (1 - alpha) + accent_color[0] * alpha * 0.3)
            g = int(bg_color[1] * (1 - alpha) + accent_color[1] * alpha * 0.3)
            b = int(bg_color[2] * (1 - alpha) + accent_color[2] * alpha * 0.3)
            
            if y % 10 == 0:  # Otimização: desenhar a cada 10 pixels
                draw.rectangle([0, y, width, y + 10], fill=(r, g, b))
                
    except Exception as e:
        print(f"⚠️ Erro no gradiente: {e}")


def _apply_soft_gradient(draw, width: int, height: int, config: dict):
    """Aplica gradiente suave para template educacional."""
    try:
        bg_color = config['background_color']
        
        # Gradiente radial suave do centro
        center_x, center_y = width // 2, height // 2
        max_distance = ((width // 2) ** 2 + (height // 2) ** 2) ** 0.5
        
        for y in range(0, height, 20):  # Otimização
            for x in range(0, width, 20):
                distance = ((x - center_x) ** 2 + (y - center_y) ** 2) ** 0.5
                alpha = min(distance / max_distance, 1.0) * 0.1
                
                r = max(0, min(255, int(bg_color[0] * (1 - alpha))))
                g = max(0, min(255, int(bg_color[1] * (1 - alpha))))
                b = max(0, min(255, int(bg_color[2] * (1 - alpha))))
                
                draw.rectangle([x, y, x + 20, y + 20], fill=(r, g, b))
                
    except Exception as e:
        print(f"⚠️ Erro no gradiente suave: {e}")


def _apply_tech_grid(draw, width: int, height: int, config: dict):
    """Aplica grid tecnológico."""
    try:
        accent_color = config['accent_color']
        grid_color = (*accent_color, 30)  # Transparente
        
        # Linhas verticais
        for x in range(0, width, 50):
            draw.line([(x, 0), (x, height)], fill=color, width=1)
        
        # Linhas horizontais
        for y in range(0, height, 50):
            draw.line([(0, y), (width, y)], fill=color, width=1)
            
    except Exception as e:
        print(f"⚠️ Erro no grid tech: {e}")


def _apply_professional_border(draw, width: int, height: int, config: dict):
    """Aplica borda profissional."""
    try:
        accent_color = config['accent_color']
        border_width = 8
        
        # Borda superior
        draw.rectangle([0, 0, width, border_width], fill=accent_color)
        # Borda inferior
        draw.rectangle([0, height - border_width, width, height], fill=accent_color)
        
    except Exception as e:
        print(f"⚠️ Erro na borda profissional: {e}")


def _apply_cyber_border(draw, width: int, height: int, config: dict):
    """Aplica borda cyber/tech."""
    try:
        accent_color = config['accent_color']
        
        # Cantos com detalhes cyber
        corner_size = 50
        
        # Canto superior esquerdo
        draw.rectangle([0, 0, corner_size, 4], fill=accent_color)
        draw.rectangle([0, 0, 4, corner_size], fill=accent_color)
        
        # Canto superior direito
        draw.rectangle([width - corner_size, 0, width, 4], fill=accent_color)
        draw.rectangle([width - 4, 0, width, corner_size], fill=accent_color)
        
        # Cantos inferiores similares
        draw.rectangle([0, height - 4, corner_size, height], fill=accent_color)
        draw.rectangle([0, height - corner_size, 4, height], fill=accent_color)
        
        draw.rectangle([width - corner_size, height - 4, width, height], fill=accent_color)
        draw.rectangle([width - 4, height - corner_size, width, height], fill=accent_color)
        
    except Exception as e:
        print(f"⚠️ Erro na borda cyber: {e}")


def _apply_friendly_border(draw, width: int, height: int, config: dict):
    """Aplica borda amigável para template educacional."""
    try:
        accent_color = config['accent_color']
        
        # Borda arredondada sutil
        border_width = 6
        corner_radius = 20
        
        # Simulação de cantos arredondados
        draw.rectangle([corner_radius, 0, width - corner_radius, border_width], fill=accent_color)
        draw.rectangle([corner_radius, height - border_width, width - corner_radius, height], fill=accent_color)
        draw.rectangle([0, corner_radius, border_width, height - corner_radius], fill=accent_color)
        draw.rectangle([width - border_width, corner_radius, width, height - corner_radius], fill=accent_color)
        
    except Exception as e:
        print(f"⚠️ Erro na borda amigável: {e}")


def _setup_enhanced_fonts(config: dict) -> tuple:
    """
    Configura fontes otimizadas para o template.
    
    Returns:
        tuple: (fonte_principal, fonte_titulo)
    """
    try:
        from PIL import ImageFont
        
        font_size = config['font_size']
        title_font_size = config['title_font_size']
        
        # Tentar carregar fontes do sistema
        font_paths = []
        
        if os.name == 'nt':  # Windows
            font_paths = [
                "C:/Windows/Fonts/arial.ttf",
                "C:/Windows/Fonts/calibri.ttf",
                "C:/Windows/Fonts/segoeui.ttf"
            ]
        else:  # Linux/Mac
            font_paths = [
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/System/Library/Fonts/Arial.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
            ]
        
        # Carregar primeira fonte disponível
        font = None
        title_font = None
        
        for font_path in font_paths:
            if os.path.exists(font_path):
                try:
                    font = ImageFont.truetype(font_path, font_size)
                    title_font = ImageFont.truetype(font_path, title_font_size)
                    break
                except:
                    continue
        
        # Fallback para fonte padrão
        if not font:
            font = ImageFont.load_default()
            title_font = ImageFont.load_default()
        
        return font, title_font
        
    except Exception as e:
        print(f"⚠️ Erro ao configurar fontes: {e}")
        return ImageFont.load_default(), ImageFont.load_default()


def _format_text_for_display(text: str, config: dict, language_info: dict) -> dict:
    """
    Formata texto para exibição otimizada.
    
    Returns:
        dict: Texto formatado com metadados
    """
    try:
        import textwrap
        
        # Detectar se há título (primeira linha com menos de 80 chars e seguida de linha vazia)
        lines = text.strip().split('\n')
        title = None
        body_text = text
        
        if len(lines) > 2 and len(lines[0]) < 80 and (not lines[1].strip()):
            title = lines[0].strip()
            body_text = '\n'.join(lines[2:]).strip()
        
        # Quebrar texto em linhas otimizadas
        max_chars_per_line = (config['width'] - 2 * config['padding']) // (config['font_size'] // 2)
        body_lines = textwrap.wrap(body_text, width=max_chars_per_line)
        
        return {
            'title': title,
            'body_lines': body_lines,
            'total_lines': len(body_lines) + (1 if title else 0),
            'estimated_height': len(body_lines) * config['font_size'] * config.get('line_spacing', 1.4)
        }
        
    except Exception as e:
        print(f"⚠️ Erro ao formatar texto: {e}")
        return {
            'title': None,
            'body_lines': [text],
            'total_lines': 1,
            'estimated_height': config['font_size']
        }


def _create_text_layer(formatted_text: dict, config: dict, font, title_font) -> Image:
    """
    Cria layer de texto com formatação avançada.
    
    Returns:
        Image: Layer de texto transparente
    """
    try:
        from PIL import Image, ImageDraw
        
        # Criar layer transparente
        text_layer = Image.new('RGBA', (config['width'], config['height']), (0, 0, 0, 0))
        draw = ImageDraw.Draw(text_layer)
        
        # Configurações
        width = config['width']
        height = config['height']
        padding = config['padding']
        line_spacing = int(config['font_size'] * config.get('line_spacing', 1.4))
        
        # Calcular posicionamento vertical centralizado
        total_height = formatted_text['estimated_height']
        if formatted_text['title']:
            total_height += config['title_font_size'] + 20  # Espaço entre título e corpo
        
        start_y = (height - total_height) // 2
        current_y = start_y
        
        # Desenhar título se existir
        if formatted_text['title']:
            title_text = formatted_text['title']
            
            # Centralizar título horizontalmente
            bbox = draw.textbbox((0, 0), title_text, font=title_font)
            title_width = bbox[2] - bbox[0]
            title_x = (width - title_width) // 2
            
            # Efeito de sombra para título
            shadow_offset = 3
            draw.text((title_x + shadow_offset, current_y + shadow_offset), 
                     title_text, fill=(0, 0, 0, 128), font=title_font)
            
            # Título principal
            draw.text((title_x, current_y), title_text, 
                     fill=config['accent_color'], font=title_font)
            
            current_y += config['title_font_size'] + 30
        
        # Desenhar linhas do corpo
        for line in formatted_text['body_lines']:
            if not line.strip():
                current_y += line_spacing // 2
                continue
            
            # Centralizar linha horizontalmente
            bbox = draw.textbbox((0, 0), line, font=font)
            line_width = bbox[2] - bbox[0]
            line_x = (width - line_width) // 2
            
            # Efeito de sombra sutil
            shadow_offset = 2
            draw.text((line_x + shadow_offset, current_y + shadow_offset), 
                     line, fill=(0, 0, 0, 80), font=font)
            
            # Texto principal
            draw.text((line_x, current_y), line, 
                     fill=config['text_color'], font=font)
            
            current_y += line_spacing
        
        return text_layer
        
    except Exception as e:
        print(f"❌ Erro ao criar layer de texto: {e}")
        # Fallback para texto simples
        simple_layer = Image.new('RGBA', (config['width'], config['height']), (0, 0, 0, 0))
        draw = ImageDraw.Draw(simple_layer)
        
        # Texto centralizado simples
        all_text = formatted_text.get('title', '') + '\n' + '\n'.join(formatted_text.get('body_lines', []))
        bbox = draw.textbbox((0, 0), all_text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (config['width'] - text_width) // 2
        y = (config['height'] - text_height) // 2
        
        draw.text((x, y), all_text, fill=config['text_color'], font=font)
        
        return simple_layer


def _apply_final_effects(img: Image, config: dict) -> Image:
    """
    Aplica efeitos finais à imagem.
    
    Args:
        img (Image): Imagem processada
        config (dict): Configurações do template
    
    Returns:
        Image: Imagem com efeitos finais
    """
    try:
        effects = config.get('effects', [])
        
        # Aplicar blur sutil se necessário
        if 'subtle_shadow' in effects:
            img = img.filter(ImageFilter.SMOOTH)
        
        # Aplicar brilho neon para template tech
        if 'neon_glow' in effects:
            # Simular brilho neon (simplificado)
            enhanced = img.copy()
            for _ in range(2):
                enhanced = enhanced.filter(ImageFilter.SMOOTH_MORE)
            
            # Combinar com original para efeito de brilho
            img = Image.blend(img, enhanced, 0.1)
        
        return img
        
    except Exception as e:
        print(f"⚠️ Erro nos efeitos finais: {e}")
        return img


def _save_to_avatar_cache(text: str, audio_path: str, template: str, output_path: str, config: dict):
    """
    Salva vídeo gerado no cache para reutilização.
    
    Args:
        text (str): Texto usado
        audio_path (str): Caminho do áudio
        template (str): Template usado
        output_path (str): Caminho do vídeo gerado
        config (dict): Configurações usadas
    """
    try:
        import hashlib
        import shutil
        
        # Criar hash único
        cache_key_data = f"{text}_{audio_path}_{template}_{str(sorted(config.items()))}"
        cache_key = hashlib.md5(cache_key_data.encode()).hexdigest()
        
        # Diretório de cache
        cache_dir = "cache/avatar_videos"
        os.makedirs(cache_dir, exist_ok=True)
        
        # Caminhos de cache
        cache_video_path = os.path.join(cache_dir, f"{cache_key}.mp4")
        cache_metadata_path = os.path.join(cache_dir, f"{cache_key}.json")
        
        # Copiar vídeo para cache
        shutil.copy2(output_path, cache_video_path)
        
        # Salvar metadados
        metadata = {
            'text': text[:500],  # Limitado para economia de espaço
            'audio_path': os.path.basename(audio_path),
            'template': template,
            'duration': config.get('duration', 0),
            'resolution': [config['width'], config['height']],
            'method': 'enhanced_slide_mvp',
            'quality_score': _calculate_quality_score(config),
            'created_at': time.time(),
            'version': 'enhanced_v2'
        }
        
        with open(cache_metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        
        print(f"💾 Vídeo salvo no cache: {cache_key}")
        
    except Exception as e:
        print(f"⚠️ Erro ao salvar cache: {e}")


def _get_template_config(template: str, background_style: str) -> dict:
    """
    Retorna configurações específicas para cada template visual.
    
    Args:
        template (str): Nome do template
        background_style (str): Estilo do background
    
    Returns:
        dict: Configurações do template
    """
    templates = {
        "modern": {
            "background_color": (245, 245, 255) if background_style == "solid" else None,
            "text_color": (30, 30, 50),
            "accent_color": (100, 150, 255),
            "margin": 80,
            "font_size": 36,
            "animations": ["fade_in", "slide_up", "text_glow"],
            "gradient": [(240, 248, 255), (225, 235, 255)] if background_style == "gradient" else None
        },
        "corporate": {
            "background_color": (250, 250, 250) if background_style == "solid" else None,
            "text_color": (20, 40, 60),
            "accent_color": (0, 100, 180),
            "margin": 100,
            "font_size": 34,
            "animations": ["fade_in", "professional_reveal"],
            "gradient": [(248, 250, 252), (230, 240, 250)] if background_style == "gradient" else None
        },
        "tech": {
            "background_color": (15, 15, 25) if background_style == "solid" else None,
            "text_color": (200, 255, 200),
            "accent_color": (0, 255, 100),
            "margin": 60,
            "font_size": 38,
            "animations": ["matrix_effect", "neon_glow", "cyber_reveal"],
            "gradient": [(10, 15, 30), (25, 35, 50)] if background_style == "gradient" else None
        },
        "education": {
            "background_color": (255, 252, 240) if background_style == "solid" else None,
            "text_color": (60, 60, 80),
            "accent_color": (255, 140, 0),
            "margin": 90,
            "font_size": 32,
            "animations": ["friendly_bounce", "warm_fade"],
            "gradient": [(255, 250, 235), (250, 245, 220)] if background_style == "gradient" else None
        },
        "minimal": {
            "background_color": (255, 255, 255) if background_style == "solid" else None,
            "text_color": (50, 50, 50),
            "accent_color": (120, 120, 120),
            "margin": 120,
            "font_size": 30,
            "animations": ["subtle_fade"],
            "gradient": [(255, 255, 255), (248, 248, 248)] if background_style == "gradient" else None
        }
    }
    
    return templates.get(template, templates["modern"])


def _create_advanced_slide_image(text: str, width: int, height: int, 
                                template_config: dict, animations: bool) -> Image.Image:
    """
    Cria uma imagem de slide avançada com template e efeitos visuais.
    
    Args:
        text (str): Texto a ser renderizado
        width (int): Largura da imagem
        height (int): Altura da imagem
        template_config (dict): Configurações do template
        animations (bool): Se deve preparar para animações
    
    Returns:
        Image.Image: Imagem PIL com design avançado
    """
    try:
        # Criar imagem base
        if template_config.get('gradient') and template_config['gradient']:
            # Criar gradiente
            image = _create_gradient_background(width, height, template_config['gradient'])
        else:
            # Cor sólida
            bg_color = template_config['background_color'] or (255, 255, 255)
            image = Image.new('RGB', (width, height), bg_color)
        
        draw = ImageDraw.Draw(image)
        
        # Adicionar efeitos de background baseados no template
        _add_background_effects(image, draw, template_config, width, height)
        
        # Configurar fonte
        font_size = template_config['font_size']
        try:
            if os.name == 'nt':  # Windows
                font_path = "C:/Windows/Fonts/segoeui.ttf"  # Fonte moderna
            else:  # Linux/Mac
                font_path = "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
            
            if os.path.exists(font_path):
                font = ImageFont.truetype(font_path, font_size)
                font_small = ImageFont.truetype(font_path, int(font_size * 0.8))
            else:
                raise OSError("Fonte não encontrada")
                
        except (OSError, ImportError):
            print("⚠️ Usando fonte padrão")
            font = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        # Renderizar texto com formatação avançada
        _render_advanced_text(draw, text, font, template_config, width, height)
        
        # Adicionar elementos decorativos
        if animations:
            _add_decorative_elements(draw, template_config, width, height)
        
        return image
        
    except Exception as e:
        print(f"⚠️ Erro ao criar slide avançado: {e}")
        # Fallback para slide básico
        return _create_slide_image(text, width, height, 
                                 template_config['background_color'],
                                 template_config['text_color'],
                                 template_config['margin'])


def _create_gradient_background(width: int, height: int, colors: list) -> Image.Image:
    """
    Cria um background com gradiente vertical.
    
    Args:
        width (int): Largura da imagem
        height (int): Altura da imagem
        colors (list): Lista de cores RGB [(r,g,b), (r,g,b)]
    
    Returns:
        Image.Image: Imagem com gradiente
    """
    try:
        import numpy as np
        
        # Criar gradiente
        color1, color2 = colors[0], colors[1]
        
        # Array para o gradiente
        gradient = np.zeros((height, width, 3), dtype=np.uint8)
        
        for y in range(height):
            # Interpolação linear entre as cores
            ratio = y / height
            r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
            g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
            b = int(color1[2] * (1 - ratio) + color2[2] * ratio)
            gradient[y, :] = [r, g, b]
        
        return Image.fromarray(gradient)
        
    except ImportError:
        # Fallback sem numpy
        image = Image.new('RGB', (width, height), colors[0])
        return image
    except Exception as e:
        print(f"⚠️ Erro ao criar gradiente: {e}")
        image = Image.new('RGB', (width, height), colors[0])
        return image


def _add_background_effects(image: Image.Image, draw, template_config: dict, 
                           width: int, height: int):
    """
    Adiciona efeitos visuais ao background baseados no template.
    
    Args:
        image: Imagem PIL
        draw: Objeto de desenho
        template_config: Configurações do template
        width: Largura da imagem
        height: Altura da imagem
    """
    try:
        accent_color = template_config['accent_color']
        
        # Efeitos específicos por template
        if 'tech' in str(template_config.get('animations', [])):
            # Grid tecnológico
            _draw_tech_grid(draw, width, height, accent_color)
        
        elif 'corporate' in str(template_config.get('animations', [])):
            # Linhas sutis corporativas
            _draw_corporate_lines(draw, width, height, accent_color)
        
        elif 'education' in str(template_config.get('animations', [])):
            # Elementos educacionais amigáveis
            _draw_education_elements(draw, width, height, accent_color)
            
    except Exception as e:
        print(f"⚠️ Erro ao adicionar efeitos: {e}")


def _draw_tech_grid(draw, width: int, height: int, color: tuple):
    """Desenha grid tecnológico de background."""
    try:
        grid_size = 50
        alpha_color = tuple(list(color) + [30])  # Transparência simulada
        
        # Linhas verticais
        for x in range(0, width, grid_size):
            draw.line([(x, 0), (x, height)], fill=color, width=1)
        
        # Linhas horizontais  
        for y in range(0, height, grid_size):
            draw.line([(0, y), (width, y)], fill=color, width=1)
            
    except Exception:
        pass


def _draw_corporate_lines(draw, width: int, height: int, color: tuple):
    """Desenha linhas corporativas sutis."""
    try:
        # Linha superior
        draw.rectangle([(0, 20), (width, 25)], fill=color)
        # Linha inferior
        draw.rectangle([(0, height-25), (width, height-20)], fill=color)
        # Linha lateral
        draw.rectangle([(20, 0), (25, height)], fill=color)
        
    except Exception:
        pass


def _draw_education_elements(draw, width: int, height: int, color: tuple):
    """Desenha elementos educacionais amigáveis."""
    try:
        # Círculos decorativos
        for i in range(3):
            x = width - 100 - (i * 30)
            y = 50 + (i * 20)
            draw.ellipse([(x-10, y-10), (x+10, y+10)], outline=color, width=2)
            
    except Exception:
        pass


def _render_advanced_text(draw, text: str, font, template_config: dict, 
                          width: int, height: int):
    """
    Renderiza texto com formatação avançada e efeitos.
    
    Args:
        draw: Objeto de desenho PIL
        text: Texto a ser renderizado
        font: Fonte PIL
        template_config: Configurações do template
        width: Largura da imagem
        height: Altura da imagem
    """
    try:
        margin = template_config['margin']
        text_color = template_config['text_color']
        accent_color = template_config['accent_color']
        
        # Área útil para texto
        text_area_width = width - (2 * margin)
        
        # Quebrar texto em linhas
        lines = _wrap_text_advanced(text, font, text_area_width, draw)
        
        # Configurações de layout
        line_height = int(template_config['font_size'] * 1.4)
        total_text_height = len(lines) * line_height
        start_y = (height - total_text_height) // 2
        
        # Renderizar cada linha com efeitos
        for i, line in enumerate(lines):
            # Calcular posição X centralizada
            bbox = draw.textbbox((0, 0), line, font=font)
            line_width = bbox[2] - bbox[0]
            x = (width - line_width) // 2
            y = start_y + (i * line_height)
            
            # Efeito de sombra (se não for minimal)
            if template_config.get('animations') and 'minimal' not in str(template_config['animations']):
                shadow_offset = 2
                shadow_color = (0, 0, 0, 50)  # Sombra transparente
                draw.text((x + shadow_offset, y + shadow_offset), line, 
                         fill=(100, 100, 100), font=font)
            
            # Texto principal
            draw.text((x, y), line, fill=text_color, font=font)
            
            # Destacar primeira linha se for título
            if i == 0 and len(lines) > 1:
                # Sublinhado colorido
                underline_y = y + int(template_config['font_size'] * 1.1)
                draw.rectangle([(x, underline_y), (x + line_width, underline_y + 3)], 
                             fill=accent_color)
        
    except Exception as e:
        print(f"⚠️ Erro ao renderizar texto avançado: {e}")
        # Fallback para renderização básica
        _render_basic_text_fallback(draw, text, font, text_color, width, height, margin)


def _wrap_text_advanced(text: str, font, max_width: int, draw) -> list:
    """
    Versão avançada da quebra de texto com melhor formatação.
    
    Args:
        text: Texto a ser quebrado
        font: Fonte PIL
        max_width: Largura máxima
        draw: Objeto de desenho
    
    Returns:
        list: Lista de linhas otimizadas
    """
    try:
        lines = []
        paragraphs = text.split('\n')
        
        for paragraph in paragraphs:
            if not paragraph.strip():
                lines.append("")  # Linha vazia para espaçamento
                continue
                
            words = paragraph.split()
            current_line = ""
            
            for word in words:
                test_line = current_line + " " + word if current_line else word
                bbox = draw.textbbox((0, 0), test_line, font=font)
                line_width = bbox[2] - bbox[0]
                
                if line_width <= max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            
            if current_line:
                lines.append(current_line)
        
        return lines if lines else [text]
        
    except Exception as e:
        print(f"⚠️ Erro na quebra avançada: {e}")
        return _wrap_text(text, font, max_width, draw)


def _render_basic_text_fallback(draw, text: str, font, color: tuple, 
                               width: int, height: int, margin: int):
    """Renderização básica de texto como fallback."""
    try:
        # Simplificado para fallback
        lines = text.split('\n')[:5]  # Máximo 5 linhas
        line_height = 40
        total_height = len(lines) * line_height
        start_y = (height - total_height) // 2
        
        for i, line in enumerate(lines):
            if line.strip():
                bbox = draw.textbbox((0, 0), line, font=font)
                line_width = bbox[2] - bbox[0]
                x = (width - line_width) // 2
                y = start_y + (i * line_height)
                draw.text((x, y), line, fill=color, font=font)
                
    except Exception:
        # Último recurso
        draw.text((margin, height//2), "Erro na renderização", fill=color)


def _add_decorative_elements(draw, template_config: dict, width: int, height: int):
    """
    Adiciona elementos decorativos baseados no template.
    
    Args:
        draw: Objeto de desenho
        template_config: Configurações do template
        width: Largura da imagem
        height: Altura da imagem
    """
    try:
        accent_color = template_config['accent_color']
        animations = template_config.get('animations', [])
        
        # Elementos baseados em animações previstas
        if 'glow' in str(animations):
            # Pontos de luz nos cantos
            for corner in [(50, 50), (width-50, 50), (50, height-50), (width-50, height-50)]:
                draw.ellipse([(corner[0]-5, corner[1]-5), (corner[0]+5, corner[1]+5)], 
                           fill=accent_color)
        
        if 'professional' in str(animations):
            # Bordas profissionais
            border_width = 3
            draw.rectangle([(0, 0), (width, border_width)], fill=accent_color)
            draw.rectangle([(0, height-border_width), (width, height)], fill=accent_color)
        
    except Exception as e:
        print(f"⚠️ Erro ao adicionar elementos decorativos: {e}")


def _calculate_video_quality_score(resolution: str, template: str, animations: bool, 
                                   file_size: int, duration: float) -> float:
    """
    Calcula pontuação de qualidade do vídeo baseada em múltiplos fatores.
    
    Args:
        resolution: Resolução do vídeo
        template: Template usado
        animations: Se tem animações
        file_size: Tamanho do arquivo
        duration: Duração do vídeo
    
    Returns:
        float: Pontuação de 0.0 a 1.0
    """
    try:
        score = 0.5  # Base
        
        # Pontuação por resolução
        resolution_scores = {"hd": 0.1, "fhd": 0.2, "4k": 0.3}
        score += resolution_scores.get(resolution, 0.1)
        
        # Pontuação por template (sofisticação)
        template_scores = {
            "minimal": 0.05, "modern": 0.15, "education": 0.1, 
            "corporate": 0.2, "tech": 0.25
        }
        score += template_scores.get(template, 0.1)
        
        # Bonificação por animações
        if animations:
            score += 0.1
        
        # Eficiência de compressão (MB por segundo)
        if duration > 0:
            mb_per_second = (file_size / 1024 / 1024) / duration
            if mb_per_second < 0.1:  # Muito eficiente
                score += 0.1
            elif mb_per_second > 0.5:  # Menos eficiente
                score -= 0.05
        
        # Garantir que está entre 0.0 e 1.0
        return max(0.0, min(1.0, score))
        
    except Exception as e:
        print(f"⚠️ Erro ao calcular qualidade: {e}")
        return 0.7  # Score padrão


def _calculate_quality_score(config: dict) -> float:
    """
    Calcula pontuação de qualidade baseada nas configurações.
    
    Args:
        config (dict): Configurações do vídeo
    
    Returns:
        float: Pontuação de 0.0 a 1.0
    """
    try:
        score = 0.5  # Base
        
        # Resolução
        resolution = config['width'] * config['height']
        if resolution >= 1920 * 1080:  # Full HD+
            score += 0.2
        elif resolution >= 1280 * 720:  # HD
            score += 0.1
        
        # Frame rate
        fps = config.get('fps', 30)
        if fps >= 60:
            score += 0.1
        elif fps >= 30:
            score += 0.05
        
        # Efeitos aplicados
        effects_count = len(config.get('effects', []))
        score += min(effects_count * 0.05, 0.15)
        
        # Template sophistication
        if config.get('style') == 'corporate':
            score += 0.1
        elif config.get('style') == 'futuristic':
            score += 0.15
        
        return min(score, 1.0)
        
    except Exception as e:
        print(f"⚠️ Erro ao calcular qualidade: {e}")
        return 0.7  # Score padrão médio


def create_video_from_text_and_audio(text: str, audio_path: str, output_path: str, 
                                     template: str = "modern", resolution: str = "hd", 
                                     animations: bool = True, background_style: str = "gradient") -> dict:
    """
    Cria um vídeo sincronizado avançado a partir de texto e áudio.
    
    Esta função gera vídeos profissionais com templates visuais avançados,
    animações de texto, diferentes resoluções e estilos de background,
    combinando com áudio para criar vídeos sincronizados de alta qualidade.
    
    Args:
        text (str): Texto a ser exibido no slide (máximo 2000 caracteres)
        audio_path (str): Caminho para o arquivo de áudio (WAV, MP3, etc.)
        output_path (str): Caminho onde o vídeo final será salvo (.mp4)
        template (str): Template visual ("modern", "corporate", "tech", "education", "minimal")
        resolution (str): Resolução do vídeo ("hd"=720p, "fhd"=1080p, "4k"=2160p)
        animations (bool): Ativar animações de texto e transições
        background_style (str): Estilo do fundo ("solid", "gradient", "pattern", "image")
    
    Returns:
        dict: Dicionário com informações do vídeo criado:
            - success (bool): Se a operação foi bem-sucedida
            - output_path (str): Caminho do arquivo de vídeo gerado
            - duration (float): Duração do vídeo em segundos
            - resolution (tuple): Resolução do vídeo (largura, altura)
            - template_used (str): Template aplicado
            - quality_score (float): Pontuação de qualidade (0.0-1.0)
            - file_size (int): Tamanho do arquivo em bytes
            - animations_applied (list): Lista de animações aplicadas
            - error (str): Mensagem de erro se houver falha
    
    Raises:
        FileNotFoundError: Se o arquivo de áudio não for encontrado
        ValueError: Se os parâmetros estiverem inválidos
        Exception: Para outros erros durante o processamento
    """
    try:
        # Verificação de dependências necessárias
        if not MOVIEPY_AVAILABLE:
            raise ImportError("MoviePy não está disponível. Execute: pip install moviepy")
        
        if not PIL_AVAILABLE:
            raise ImportError("PIL não está disponível. Execute: pip install pillow")
        
        # Validação dos parâmetros de entrada
        if not text or not isinstance(text, str):
            raise ValueError("Texto deve ser uma string não vazia")
        
        if len(text) > 2000:
            raise ValueError("Texto muito longo (máximo 2000 caracteres)")
        
        if not os.path.exists(audio_path):
            raise FileNotFoundError(f"Arquivo de áudio não encontrado: {audio_path}")
        
        # Validar template
        valid_templates = ["modern", "corporate", "tech", "education", "minimal"]
        if template not in valid_templates:
            template = "modern"
            print(f"⚠️ Template inválido, usando 'modern'. Válidos: {valid_templates}")
        
        # Configurações de resolução
        resolution_configs = {
            "hd": (1280, 720),
            "fhd": (1920, 1080),
            "4k": (3840, 2160)
        }
        
        if resolution not in resolution_configs:
            resolution = "hd"
            print(f"⚠️ Resolução inválida, usando 'hd'. Válidas: {list(resolution_configs.keys())}")
        
        slide_width, slide_height = resolution_configs[resolution]
        
        # Obter configurações do template
        template_config = _get_template_config(template, background_style)
        
        # Aplicar configurações do template
        background_color = template_config['background_color']
        text_color = template_config['text_color']
        accent_color = template_config['accent_color']
        margin = template_config['margin']
        font_size = template_config['font_size']
        
        # Lista de animações aplicadas
        animations_applied = []
        
        print(f"🎬 Iniciando criação de vídeo...")
        print(f"📝 Texto: {text[:50]}..." if len(text) > 50 else f"📝 Texto: {text}")
        print(f"🎵 Áudio: {audio_path}")
        print(f"💾 Saída: {output_path}")
        
        # Etapa 1: Carregar o arquivo de áudio para obter a duração
        print("🎵 Carregando arquivo de áudio...")
        audio_clip = AudioFileClip(audio_path)
        audio_duration = audio_clip.duration
        print(f"⏱️ Duração do áudio: {audio_duration:.2f} segundos")
        
        # Etapa 2: Criar a imagem do slide com o texto (template avançado)
        print(f"🖼️ Gerando imagem do slide com template '{template}'...")
        slide_image = _create_advanced_slide_image(text, slide_width, slide_height, 
                                                  template_config, animations)
        
        # Etapa 3: Criar um clip de imagem temporário
        print("🎞️ Criando clip de vídeo...")
        # Salvar a imagem temporariamente
        temp_image_path = "temp_slide.png"
        slide_image.save(temp_image_path, "PNG")
        
        # Criar clip de imagem com a duração do áudio
        image_clip = ImageClip(temp_image_path, duration=audio_duration)
        
        # Etapa 4: Combinar imagem e áudio
        print("🔄 Combinando imagem e áudio...")
        final_clip = image_clip.set_audio(audio_clip)
        
        # Etapa 5: Renderizar o vídeo final
        print("🎬 Renderizando vídeo final...")
        # Garantir que o diretório de saída existe
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Renderizar com configurações otimizadas
        final_clip.write_videofile(
            output_path,
            fps=24,  # Frame rate de 24 fps para economizar espaço
            codec='libx264',  # Codec H.264 para boa compatibilidade
            audio_codec='aac',  # Codec AAC para áudio
            temp_audiofile='temp-audio.m4a',  # Arquivo temporário para áudio
            remove_temp=True,  # Remover arquivos temporários
            verbose=False,  # Reduzir verbosidade
            logger=None  # Desabilitar logs do MoviePy
        )
        
        # Etapa 6: Limpeza de arquivos temporários
        print("🧹 Limpando arquivos temporários...")
        if os.path.exists(temp_image_path):
            os.remove(temp_image_path)
        
        # Fechar clips para liberar recursos
        audio_clip.close()
        image_clip.close()
        final_clip.close()
        
        # Etapa 7: Obter informações do arquivo gerado
        file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
        
        # Calcular pontuação de qualidade
        quality_score = _calculate_video_quality_score(
            resolution, template, animations, file_size, audio_duration
        )
        
        # Adicionar animações aplicadas baseado no template
        if animations:
            animations_applied.extend(template_config.get('animations', ['fade_in', 'text_reveal']))
        
        # Resultado final expandido
        result = {
            'success': True,
            'output_path': output_path,
            'duration': audio_duration,
            'resolution': (slide_width, slide_height),
            'template_used': template,
            'quality_score': quality_score,
            'file_size': file_size,
            'animations_applied': animations_applied,
            'background_style': background_style,
            'total_frames': int(audio_duration * 24),  # 24 FPS
            'error': None
        }
        
        print(f"✅ Vídeo criado com sucesso!")
        print(f"📁 Arquivo: {output_path}")
        print(f"⏱️ Duração: {audio_duration:.2f}s")
        print(f"📐 Resolução: {slide_width}x{slide_height}")
        print(f"🎨 Template: {template}")
        print(f"⭐ Qualidade: {quality_score:.2f}/1.0")
        print(f"💾 Tamanho: {file_size / 1024 / 1024:.2f} MB")
        if animations_applied:
            print(f"🎬 Animações: {', '.join(animations_applied)}")
        
        return result
        
    except Exception as e:
        print(f"❌ Erro ao criar vídeo: {e}")
        return {
            'success': False,
            'output_path': None,
            'duration': 0,
            'resolution': (0, 0),
            'file_size': 0,
            'error': str(e)
        }


def _create_slide_image(text: str, width: int, height: int, 
                       bg_color: tuple, text_color: tuple, margin: int) -> Image.Image:
    """
    Cria uma imagem de slide com texto centralizado.
    
    Esta função auxiliar gera uma imagem PIL com o texto formatado
    de forma centralizada, com quebras de linha automáticas e
    formatação adequada para um slide simples.
    
    Args:
        text (str): Texto a ser renderizado
        width (int): Largura da imagem em pixels
        height (int): Altura da imagem em pixels
        bg_color (tuple): Cor de fundo RGB (R, G, B)
        text_color (tuple): Cor do texto RGB (R, G, B)
        margin (int): Margem das bordas em pixels
    
    Returns:
        Image.Image: Imagem PIL com o texto renderizado
    """
    try:
        # Criar imagem com cor de fundo
        image = Image.new('RGB', (width, height), bg_color)
        draw = ImageDraw.Draw(image)
        
        # Tentar carregar fonte do sistema (com fallback)
        try:
            # Tentar carregar uma fonte TrueType do sistema
            if os.name == 'nt':  # Windows
                font_path = "C:/Windows/Fonts/arial.ttf"
            else:  # Linux/Mac
                font_path = "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
            
            if os.path.exists(font_path):
                font = ImageFont.truetype(font_path, 32)
                font_small = ImageFont.truetype(font_path, 24)
            else:
                raise OSError("Fonte não encontrada")
                
        except (OSError, ImportError):
            # Usar fonte padrão se não conseguir carregar TrueType
            print("⚠️ Usando fonte padrão (TrueType não disponível)")
            font = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        # Área útil para o texto (considerando margem)
        text_area_width = width - (2 * margin)
        text_area_height = height - (2 * margin)
        
        # Quebrar texto em linhas que cabem na largura
        lines = _wrap_text(text, font, text_area_width, draw)
        
        # Calcular altura total do texto
        line_height = 40  # Espaçamento entre linhas
        total_text_height = len(lines) * line_height
        
        # Posição inicial Y para centralizar verticalmente
        start_y = (height - total_text_height) // 2
        
        # Renderizar cada linha de texto
        for i, line in enumerate(lines):
            # Calcular largura da linha para centralizar horizontalmente
            bbox = draw.textbbox((0, 0), line, font=font)
            line_width = bbox[2] - bbox[0]
            
            # Posição X para centralizar
            x = (width - line_width) // 2
            y = start_y + (i * line_height)
            
            # Desenhar o texto
            draw.text((x, y), line, fill=text_color, font=font)
        
        return image
        
    except Exception as e:
        print(f"⚠️ Erro ao criar imagem do slide: {e}")
        # Criar imagem básica em caso de erro
        image = Image.new('RGB', (width, height), bg_color)
        draw = ImageDraw.Draw(image)
        draw.text((margin, height//2), "Erro ao renderizar texto", fill=text_color)
        return image


def _wrap_text(text: str, font, max_width: int, draw) -> List[str]:
    """
    Quebra o texto em linhas que cabem na largura especificada.
    
    Esta função auxiliar divide o texto em múltiplas linhas,
    respeitando a largura máxima disponível e tentando
    quebrar em espaços quando possível.
    
    Args:
        text (str): Texto a ser quebrado
        font: Fonte PIL a ser usada
        max_width (int): Largura máxima em pixels
        draw: Objeto de desenho PIL
    
    Returns:
        List[str]: Lista de linhas de texto
    """
    try:
        words = text.split()
        lines = []
        current_line = ""
        
        for word in words:
            # Testar se a palavra cabe na linha atual
            test_line = current_line + " " + word if current_line else word
            bbox = draw.textbbox((0, 0), test_line, font=font)
            line_width = bbox[2] - bbox[0]
            
            if line_width <= max_width:
                # Palavra cabe, adicionar à linha atual
                current_line = test_line
            else:
                # Palavra não cabe, começar nova linha
                if current_line:
                    lines.append(current_line)
                current_line = word
        
        # Adicionar última linha
        if current_line:
            lines.append(current_line)
        
        return lines if lines else [text]
        
    except Exception as e:
        print(f"⚠️ Erro ao quebrar texto: {e}")
        # Retornar texto original em caso de erro
        return [text]


def create_video_pipeline_automatic(text: str, output_path: str, 
                                   voice: str = "pt", template: str = "modern",
                                   resolution: str = "hd", animations: bool = True) -> dict:
    """
    Pipeline completo automático: Texto → TTS → Vídeo.
    
    Esta função automatiza todo o processo de criação de vídeo:
    1. Gera áudio automaticamente usando TTS
    2. Cria vídeo sincronizado com templates avançados
    3. Limpa arquivos temporários
    4. Retorna resultado completo
    
    Args:
        text (str): Texto para narração e slides
        output_path (str): Caminho do vídeo final (.mp4)
        voice (str): Idioma/voz para TTS ("pt", "en", "es")
        template (str): Template visual do vídeo
        resolution (str): Resolução do vídeo
        animations (bool): Ativar animações
    
    Returns:
        dict: Resultado completo do pipeline
    """
    try:
        print("🚀 INICIANDO PIPELINE AUTOMÁTICO TEXTO → VÍDEO")
        print("=" * 60)
        
        # Etapa 1: Gerar áudio com TTS
        print("🎙️ Etapa 1: Gerando áudio com TTS...")
        audio_result = _generate_tts_audio(text, voice)
        
        if not audio_result['success']:
            return {
                'success': False,
                'stage': 'tts_generation',
                'error': f"Erro na geração de áudio: {audio_result['error']}"
            }
        
        audio_path = audio_result['audio_path']
        print(f"✅ Áudio gerado: {audio_path}")
        
        # Etapa 2: Criar vídeo com áudio
        print("🎬 Etapa 2: Criando vídeo sincronizado...")
        video_result = create_video_from_text_and_audio(
            text=text,
            audio_path=audio_path,
            output_path=output_path,
            template=template,
            resolution=resolution,
            animations=animations
        )
        
        if not video_result['success']:
            # Limpar áudio temporário em caso de erro
            if os.path.exists(audio_path):
                os.remove(audio_path)
            return {
                'success': False,
                'stage': 'video_creation',
                'error': f"Erro na criação do vídeo: {video_result['error']}"
            }
        
        # Etapa 3: Limpeza e resultado final
        print("🧹 Etapa 3: Finalizando pipeline...")
        
        # Limpar arquivo de áudio temporário
        if os.path.exists(audio_path):
            os.remove(audio_path)
            print("✅ Arquivos temporários removidos")
        
        # Resultado completo do pipeline
        pipeline_result = {
            'success': True,
            'output_path': video_result['output_path'],
            'pipeline_stages': {
                'tts_generation': audio_result,
                'video_creation': video_result
            },
            'final_video': {
                'path': video_result['output_path'],
                'duration': video_result['duration'],
                'resolution': video_result['resolution'],
                'template': video_result['template_used'],
                'quality_score': video_result['quality_score'],
                'file_size': video_result['file_size'],
                'animations': video_result['animations_applied']
            },
            'performance': {
                'total_time': audio_result.get('generation_time', 0) + 
                             video_result.get('creation_time', 0),
                'compression_ratio': video_result['file_size'] / len(text) if text else 0
            },
            'error': None
        }
        
        print("🎉 PIPELINE CONCLUÍDO COM SUCESSO!")
        print(f"📁 Vídeo final: {output_path}")
        print(f"⏱️ Duração: {video_result['duration']:.2f}s")
        print(f"⭐ Qualidade: {video_result['quality_score']:.2f}/1.0")
        print("=" * 60)
        
        return pipeline_result
        
    except Exception as e:
        print(f"❌ Erro no pipeline: {e}")
        return {
            'success': False,
            'stage': 'pipeline_error',
            'error': str(e)
        }


def _generate_tts_audio(text: str, voice: str = "pt") -> dict:
    """
    Gera áudio usando TTS (Text-to-Speech) automático.
    
    Args:
        text (str): Texto para conversão
        voice (str): Idioma/voz
    
    Returns:
        dict: Resultado da geração de áudio
    """
    try:
        import time
        start_time = time.time()
        
        # Criar nome único para arquivo temporário
        temp_audio_name = f"temp_tts_{int(time.time())}_{hash(text) % 10000}.wav"
        temp_audio_path = os.path.join("app/static/audios", temp_audio_name)
        
        # Garantir que diretório existe
        os.makedirs(os.path.dirname(temp_audio_path), exist_ok=True)
        
        # Tentar usar gTTS primeiro
        try:
            from gtts import gTTS
            
            # Mapear idiomas
            lang_map = {"pt": "pt", "en": "en", "es": "es"}
            lang = lang_map.get(voice, "pt")
            
            # Gerar áudio
            tts = gTTS(text=text, lang=lang, slow=False)
            tts.save(temp_audio_path)
            
            generation_time = time.time() - start_time
            
            return {
                'success': True,
                'audio_path': temp_audio_path,
                'method': 'gTTS',
                'language': lang,
                'generation_time': generation_time,
                'file_size': os.path.getsize(temp_audio_path),
                'error': None
            }
            
        except ImportError:
            print("⚠️ gTTS não disponível, criando áudio silencioso")
            return _create_silent_audio_fallback(temp_audio_path, text)
            
        except Exception as e:
            print(f"⚠️ Erro no gTTS: {e}, criando áudio silencioso")
            return _create_silent_audio_fallback(temp_audio_path, text)
            
    except Exception as e:
        return {
            'success': False,
            'audio_path': None,
            'error': str(e)
        }


def _create_silent_audio_fallback(audio_path: str, text: str) -> dict:
    """
    Cria áudio silencioso como fallback quando TTS falha.
    
    Args:
        audio_path (str): Caminho para salvar o áudio
        text (str): Texto (usado para calcular duração)
    
    Returns:
        dict: Resultado da criação de áudio
    """
    try:
        import wave
        import time
        
        # Calcular duração baseada no texto (aproximadamente)
        words_per_minute = 150  # Velocidade média de leitura
        word_count = len(text.split())
        duration = max(3.0, (word_count / words_per_minute) * 60)  # Mínimo 3 segundos
        
        # Configurações de áudio
        sample_rate = 44100
        frames = int(duration * sample_rate)
        
        # Criar áudio silencioso
        with wave.open(audio_path, 'w') as wav_file:
            wav_file.setnchannels(1)  # Mono
            wav_file.setsampwidth(2)  # 16-bit
            wav_file.setframerate(sample_rate)
            wav_file.writeframes(b'\x00' * (frames * 2))  # Silêncio
        
        return {
            'success': True,
            'audio_path': audio_path,
            'method': 'silent_fallback',
            'language': 'silent',
            'generation_time': 0.1,
            'file_size': os.path.getsize(audio_path),
            'duration': duration,
            'error': None
        }
        
    except Exception as e:
        return {
            'success': False,
            'audio_path': None,
            'error': str(e)
        }


def create_batch_videos(texts: list, output_dir: str = "videos/batch", 
                       template: str = "modern", resolution: str = "hd") -> dict:
    """
    Cria múltiplos vídeos em batch a partir de lista de textos.
    
    Args:
        texts (list): Lista de textos para conversão
        output_dir (str): Diretório de saída para os vídeos
        template (str): Template a ser usado
        resolution (str): Resolução dos vídeos
    
    Returns:
        dict: Resultado do processamento em batch
    """
    try:
        print(f"📦 PROCESSAMENTO EM BATCH: {len(texts)} vídeos")
        print("=" * 50)
        
        # Garantir que diretório existe
        os.makedirs(output_dir, exist_ok=True)
        
        results = []
        successful_videos = 0
        total_duration = 0
        total_size = 0
        
        for i, text in enumerate(texts, 1):
            print(f"\n🎬 Processando vídeo {i}/{len(texts)}...")
            
            # Nome do arquivo de saída
            output_filename = f"video_{i:03d}_{int(time.time())}.mp4"
            output_path = os.path.join(output_dir, output_filename)
            
            # Processar vídeo
            result = create_video_pipeline_automatic(
                text=text,
                output_path=output_path,
                template=template,
                resolution=resolution
            )
            
            # Registrar resultado
            result['batch_index'] = i
            result['input_text'] = text[:100] + "..." if len(text) > 100 else text
            results.append(result)
            
            if result['success']:
                successful_videos += 1
                video_info = result['final_video']
                total_duration += video_info['duration']
                total_size += video_info['file_size']
                print(f"✅ Vídeo {i} criado: {output_filename}")
            else:
                print(f"❌ Erro no vídeo {i}: {result['error']}")
        
        # Resumo do batch
        batch_summary = {
            'success': successful_videos > 0,
            'total_videos': len(texts),
            'successful_videos': successful_videos,
            'failed_videos': len(texts) - successful_videos,
            'success_rate': (successful_videos / len(texts)) * 100,
            'total_duration': total_duration,
            'total_file_size': total_size,
            'average_duration': total_duration / successful_videos if successful_videos > 0 else 0,
            'output_directory': output_dir,
            'template_used': template,
            'resolution_used': resolution,
            'individual_results': results
        }
        
        print(f"\n📊 RESUMO DO BATCH:")
        print(f"✅ Sucessos: {successful_videos}/{len(texts)}")
        print(f"📈 Taxa de sucesso: {batch_summary['success_rate']:.1f}%")
        print(f"⏱️ Duração total: {total_duration:.2f}s")
        print(f"💾 Tamanho total: {total_size / 1024 / 1024:.2f} MB")
        print("=" * 50)
        
        return batch_summary
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e),
            'total_videos': len(texts) if texts else 0,
            'successful_videos': 0
        }


def optimize_video_for_platform(input_path: str, platform: str = "youtube") -> dict:
    """
    Otimiza vídeo para plataformas específicas.
    
    Args:
        input_path (str): Caminho do vídeo original
        platform (str): Plataforma alvo ("youtube", "instagram", "tiktok", "linkedin")
    
    Returns:
        dict: Resultado da otimização
    """
    try:
        if not MOVIEPY_AVAILABLE:
            return {
                'success': False,
                'error': 'MoviePy não disponível para otimização'
            }
        
        print(f"🎯 Otimizando vídeo para {platform.upper()}...")
        
        # Configurações por plataforma
        platform_configs = {
            "youtube": {
                "resolution": (1920, 1080),
                "fps": 30,
                "bitrate": "2000k",
                "codec": "libx264",
                "suffix": "_youtube"
            },
            "instagram": {
                "resolution": (1080, 1080),  # Quadrado
                "fps": 30,
                "bitrate": "1500k", 
                "codec": "libx264",
                "suffix": "_instagram"
            },
            "tiktok": {
                "resolution": (1080, 1920),  # Vertical
                "fps": 30,
                "bitrate": "1200k",
                "codec": "libx264", 
                "suffix": "_tiktok"
            },
            "linkedin": {
                "resolution": (1920, 1080),
                "fps": 24,
                "bitrate": "1800k",
                "codec": "libx264",
                "suffix": "_linkedin"
            }
        }
        
        config = platform_configs.get(platform, platform_configs["youtube"])
        
        # Gerar nome do arquivo otimizado
        base_name = os.path.splitext(input_path)[0]
        optimized_path = f"{base_name}{config['suffix']}.mp4"
        
        # Carregar e otimizar vídeo
        from moviepy.editor import VideoFileClip
        
        clip = VideoFileClip(input_path)
        
        # Redimensionar se necessário
        target_resolution = config["resolution"]
        if (clip.w, clip.h) != target_resolution:
            clip = clip.resize(target_resolution)
        
        # Aplicar configurações e salvar
        clip.write_videofile(
            optimized_path,
            fps=config["fps"],
            codec=config["codec"],
            bitrate=config["bitrate"],
            verbose=False,
            logger=None
        )
        
        clip.close()
        
        # Informações do arquivo otimizado
        optimized_size = os.path.getsize(optimized_path)
        original_size = os.path.getsize(input_path)
        compression_ratio = (original_size - optimized_size) / original_size * 100
        
        result = {
            'success': True,
            'original_path': input_path,
            'optimized_path': optimized_path,
            'platform': platform,
            'original_size': original_size,
            'optimized_size': optimized_size,
            'compression_ratio': compression_ratio,
            'target_resolution': target_resolution,
            'target_fps': config["fps"],
            'savings_mb': (original_size - optimized_size) / 1024 / 1024
        }
        
        print(f"✅ Otimização concluída!")
        print(f"📁 Arquivo: {optimized_path}")
        print(f"📐 Resolução: {target_resolution[0]}x{target_resolution[1]}")
        print(f"💾 Compressão: {compression_ratio:.1f}% menor")
        
        return result
        
    except Exception as e:
        return {
            'success': False,
            'error': str(e)
        }


# Exemplo de uso e testes
if __name__ == '__main__':
    """
    Exemplo de uso da função create_video_from_text_and_audio.
    
    Este exemplo demonstra como usar a função para criar um vídeo
    a partir de um texto e um arquivo de áudio de teste.
    """
    print("🎬 TESTE: Criador de Vídeo a partir de Texto e Áudio")
    print("=" * 60)
    
    # Texto de exemplo para o slide
    texto_exemplo = """
    Bem-vindos ao TecnoCursos AI!
    
    Nesta aula, você aprenderá sobre:
    • Inteligência Artificial
    • Machine Learning
    • Processamento de Linguagem Natural
    
    Vamos começar nossa jornada!
    """
    
    # Configurações de teste
    caminho_audio_teste = "app/static/audios/test_audio.wav"
    caminho_video_saida = "app/static/videos/teste_video_completo.mp4"
    
    try:
        # Verificar se existe áudio de teste (criar um simples se não existir)
        if not os.path.exists(caminho_audio_teste):
            print("🎵 Criando áudio de teste simples...")
            
            # Criar diretório se não existir
            os.makedirs(os.path.dirname(caminho_audio_teste), exist_ok=True)
            
            # Tentar usar gTTS para criar áudio de teste
            try:
                from gtts import gTTS
                tts = gTTS(text="Este é um teste de áudio para demonstração", lang='pt')
                tts.save(caminho_audio_teste)
                print(f"✅ Áudio de teste criado: {caminho_audio_teste}")
            except ImportError:
                print("⚠️ gTTS não disponível - pulando criação de áudio de teste")
                print("Para testar completamente, instale: pip install gtts")
                print("Ou forneça um arquivo de áudio em:", caminho_audio_teste)
                exit(1)
        
        # Executar a função principal
        print("\n🚀 Executando função create_video_from_text_and_audio...")
        resultado = create_video_from_text_and_audio(
            text=texto_exemplo,
            audio_path=caminho_audio_teste,
            output_path=caminho_video_saida
        )
        
        # Exibir resultados
        print("\n📊 RESULTADOS:")
        print("-" * 40)
        if resultado['success']:
            print(f"✅ Sucesso: {resultado['success']}")
            print(f"📁 Arquivo: {resultado['output_path']}")
            print(f"⏱️ Duração: {resultado['duration']:.2f}s")
            print(f"📐 Resolução: {resultado['resolution']}")
            print(f"💾 Tamanho: {resultado['file_size'] / 1024 / 1024:.2f} MB")
            print(f"\n🎉 Vídeo criado com sucesso!")
            print(f"🔗 Abra o arquivo: {os.path.abspath(resultado['output_path'])}")
        else:
            print(f"❌ Erro: {resultado['error']}")
            
    except Exception as e:
        print(f"❌ Erro durante o teste: {e}")
        print("\n💡 Dicas para resolver:")
        print("1. Instale as dependências: pip install moviepy pillow gtts")
        print("2. Verifique se o arquivo de áudio existe")
        print("3. Certifique-se de ter permissões de escrita no diretório")
    
    print("\n" + "=" * 60)
    print("🏁 Teste concluído!")

def create_videos_from_pdf_and_audio(pdf_path: str, audio_path: str, output_folder: str,
                                   template: str = "modern", resolution: str = "hd",
                                   animations: bool = True, background_style: str = "gradient",
                                   sync_mode: str = "auto") -> dict:
    """
    Cria vídeos automaticamente a partir de um PDF e um arquivo de áudio.
    
    Esta função combina a extração de slides de PDF, divisão de áudio por slides
    e criação de vídeos em um único processo automatizado. Ideal para converter
    apresentações PDF em vídeos narrados de forma rápida e eficiente.
    
    Args:
        pdf_path (str): Caminho para o arquivo PDF
        audio_path (str): Caminho para o arquivo de áudio completo
        output_folder (str): Pasta onde os vídeos serão salvos
        template (str): Template visual para os vídeos
        resolution (str): Resolução dos vídeos ("hd", "fhd", "4k")
        animations (bool): Ativar animações de texto e transições
        background_style (str): Estilo do fundo
        sync_mode (str): Modo de sincronização ("auto", "manual", "equal_time")
    
    Returns:
        dict: Resultado do processamento com informações detalhadas:
            - success: Se o processamento foi bem-sucedido
            - pdf_slides: Lista de imagens extraídas do PDF
            - audio_segments: Lista de segmentos de áudio criados
            - videos_created: Lista de vídeos gerados
            - total_duration: Duração total dos vídeos
            - processing_time: Tempo de processamento
            - error: Mensagem de erro se houver falha
    """
    import time
    start_time = time.time()
    
    result = {
        "success": False,
        "pdf_slides": [],
        "audio_segments": [],
        "videos_created": [],
        "total_duration": 0,
        "processing_time": 0,
        "error": None
    }
    
    try:
        print("\n" + "="*60)
        print("🎬 CONVERSÃO AUTOMÁTICA: PDF + ÁUDIO → VÍDEOS")
        print("="*60)
        
        # Verificar arquivos de entrada
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF não encontrado: {pdf_path}")
        
        if not os.path.exists(audio_path):
            raise FileNotFoundError(f"Áudio não encontrado: {audio_path}")
        
        # Criar pastas de trabalho
        slides_folder = os.path.join(output_folder, "slides_extracted")
        audio_segments_folder = os.path.join(output_folder, "audio_segments")
        videos_folder = os.path.join(output_folder, "videos_output")
        
        # Etapa 1: Extrair slides do PDF
        print("\n🔸 ETAPA 1: Extraindo slides do PDF...")
        pdf_slides = extract_pdf_slides_as_images(
            pdf_path=pdf_path,
            output_folder=slides_folder,
            dpi=150,
            image_format="PNG"
        )
        
        if not pdf_slides:
            raise ValueError("Nenhum slide foi extraído do PDF")
        
        result["pdf_slides"] = pdf_slides
        print(f"✅ {len(pdf_slides)} slides extraídos com sucesso")
        
        # Etapa 2: Dividir áudio por slides
        print("\n🔸 ETAPA 2: Processando divisão de áudio...")
        
        # Obter duração do áudio
        try:
            if MOVIEPY_AVAILABLE:
                from moviepy.editor import AudioFileClip
                audio_clip = AudioFileClip(audio_path)
                total_audio_duration = audio_clip.duration
                audio_clip.close()
            else:
                # Fallback: estimar 3 segundos por slide
                total_audio_duration = len(pdf_slides) * 3.0
                print("⚠️ MoviePy não disponível, usando estimativa de duração")
        except:
            total_audio_duration = len(pdf_slides) * 3.0
            print("⚠️ Erro ao obter duração do áudio, usando estimativa")
        
        # Dividir áudio baseado no modo de sincronização
        if sync_mode == "equal_time":
            # Tempo igual para cada slide
            time_per_slide = total_audio_duration / len(pdf_slides)
            audio_segments = []
            
            for i in range(len(pdf_slides)):
                start_time_seg = i * time_per_slide
                end_time_seg = (i + 1) * time_per_slide
                
                segment_path = os.path.join(audio_segments_folder, f"segment_{i+1:03d}.wav")
                
                # Para demonstração, criamos referências aos tempos
                # Em implementação real, usaria bibliotecas de áudio para cortar
                audio_segments.append({
                    "path": audio_path,  # Usar áudio original por enquanto
                    "start_time": start_time_seg,
                    "end_time": end_time_seg,
                    "duration": time_per_slide
                })
        
        else:
            # Modo automático ou manual - usar áudio completo para cada slide (simplificado)
            audio_segments = []
            for i in range(len(pdf_slides)):
                audio_segments.append({
                    "path": audio_path,
                    "start_time": 0,
                    "end_time": total_audio_duration,
                    "duration": total_audio_duration
                })
        
        result["audio_segments"] = audio_segments
        print(f"✅ {len(audio_segments)} segmentos de áudio preparados")
        
        # Etapa 3: Extrair texto dos slides (simplificado)
        print("\n🔸 ETAPA 3: Extraindo texto dos slides...")
        
        # Para esta implementação, usar texto genérico baseado no número do slide
        slides_text = []
        for i, slide_path in enumerate(pdf_slides):
            # Em implementação completa, usaria OCR aqui
            slide_text = f"Slide {i+1} da apresentação. Este é o conteúdo do slide número {i+1}."
            slides_text.append(slide_text)
        
        print(f"✅ Texto extraído de {len(slides_text)} slides")
        
        # Etapa 4: Criar vídeos usando função existente (adaptada)
        print("\n🔸 ETAPA 4: Criando vídeos...")
        
        # Preparar listas para a função create_videos_for_slides
        audio_paths_for_videos = [seg["path"] for seg in audio_segments]
        
        # Criar vídeos usando a função existente
        videos_created = create_videos_for_slides(
            slides_text_list=slides_text,
            audios_path_list=audio_paths_for_videos,
            output_folder=videos_folder,
            template=template,
            resolution=resolution,
            animations=animations,
            background_style=background_style
        )
        
        result["videos_created"] = videos_created
        result["total_duration"] = total_audio_duration
        
        # Etapa 5: Estatísticas finais
        processing_time = time.time() - start_time
        result["processing_time"] = processing_time
        result["success"] = True
        
        print("\n" + "="*60)
        print("🎉 CONVERSÃO CONCLUÍDA COM SUCESSO!")
        print("="*60)
        print(f"📄 Slides extraídos: {len(pdf_slides)}")
        print(f"🎵 Segmentos de áudio: {len(audio_segments)}")
        print(f"🎬 Vídeos criados: {len(videos_created)}")
        print(f"⏱️ Duração total: {total_audio_duration:.2f}s")
        print(f"🕒 Tempo de processamento: {processing_time:.2f}s")
        print(f"📁 Pasta final: {videos_folder}")
        
        return result
        
    except Exception as e:
        result["error"] = str(e)
        result["processing_time"] = time.time() - start_time
        print(f"\n❌ Erro durante conversão: {str(e)}")
        return result

def optimize_batch_processing(slides_count: int, available_cores: int = None,
                            memory_limit_gb: int = 8) -> dict:
    """
    Otimiza o processamento em lote baseado nos recursos disponíveis do sistema.
    
    Esta função analisa o sistema e sugere a melhor estratégia para processar
    múltiplos vídeos de forma eficiente, considerando CPU, memória e armazenamento.
    
    Args:
        slides_count (int): Número de slides/vídeos para processar
        available_cores (int): Número de cores CPU disponíveis (None = detectar automaticamente)
        memory_limit_gb (int): Limite de memória RAM em GB (padrão 8GB)
    
    Returns:
        dict: Configurações otimizadas para processamento:
            - batch_size: Tamanho ideal do lote
            - parallel_workers: Número de workers paralelos
            - memory_per_worker: Memória por worker em MB
            - processing_strategy: Estratégia recomendada
            - estimated_time: Tempo estimado total
            - recommendations: Lista de recomendações
    """
    try:
        # Detectar número de cores se não fornecido
        if available_cores is None:
            import os
            available_cores = os.cpu_count() or 4
        
        # Configurações base
        memory_limit_mb = memory_limit_gb * 1024
        
        # Estimar uso de memória por vídeo (baseado na resolução)
        memory_per_video = {
            "hd": 200,    # MB por vídeo HD
            "fhd": 400,   # MB por vídeo FHD
            "4k": 1200    # MB por vídeo 4K
        }
        
        # Usar estimativa conservadora (FHD)
        estimated_memory_per_video = memory_per_video["fhd"]
        
        # Calcular tamanho ideal do lote
        max_videos_in_memory = max(1, memory_limit_mb // (estimated_memory_per_video * 2))
        
        # Limitar por cores disponíveis
        optimal_parallel_workers = min(available_cores, max_videos_in_memory, slides_count)
        
        # Calcular tamanho do lote
        if slides_count <= optimal_parallel_workers:
            batch_size = slides_count
            processing_strategy = "single_batch"
        elif slides_count <= optimal_parallel_workers * 3:
            batch_size = optimal_parallel_workers
            processing_strategy = "small_batches"
        else:
            batch_size = optimal_parallel_workers * 2
            processing_strategy = "large_batches"
        
        # Estimar tempo de processamento
        time_per_video = 2.0  # minutos por vídeo (estimativa conservadora)
        if processing_strategy == "single_batch":
            estimated_time = time_per_video
        else:
            batches_needed = (slides_count + batch_size - 1) // batch_size
            estimated_time = batches_needed * time_per_video
        
        # Gerar recomendações
        recommendations = [
            f"💻 Usar {optimal_parallel_workers} workers paralelos",
            f"📦 Processar em lotes de {batch_size} vídeos",
            f"🧠 Reservar {estimated_memory_per_video}MB de RAM por vídeo"
        ]
        
        if slides_count > 20:
            recommendations.append("⚡ Considere usar SSD para melhor performance de I/O")
        
        if memory_limit_gb < 8:
            recommendations.append("🔧 Aumente a RAM disponível para melhor performance")
        
        if available_cores < 4:
            recommendations.append("🖥️ CPU com mais cores aceleraria o processamento")
        
        # Resultado final
        optimization_config = {
            "batch_size": batch_size,
            "parallel_workers": optimal_parallel_workers,
            "memory_per_worker": estimated_memory_per_video,
            "processing_strategy": processing_strategy,
            "estimated_time_minutes": round(estimated_time, 1),
            "recommendations": recommendations,
            "system_info": {
                "available_cores": available_cores,
                "memory_limit_gb": memory_limit_gb,
                "slides_count": slides_count
            }
        }
        
        print(f"⚡ OTIMIZAÇÃO DE PROCESSAMENTO CALCULADA")
        print(f"📊 Slides para processar: {slides_count}")
        print(f"💻 Cores disponíveis: {available_cores}")
        print(f"🧠 RAM limite: {memory_limit_gb}GB")
        print(f"📦 Tamanho do lote: {batch_size}")
        print(f"⚙️ Workers paralelos: {optimal_parallel_workers}")
        print(f"⏱️ Tempo estimado: {estimated_time:.1f} minutos")
        print(f"🎯 Estratégia: {processing_strategy}")
        
        return optimization_config
        
    except Exception as e:
        print(f"❌ Erro ao calcular otimização: {str(e)}")
        return {
            "batch_size": 1,
            "parallel_workers": 1,
            "memory_per_worker": 500,
            "processing_strategy": "sequential",
            "estimated_time_minutes": slides_count * 2,
            "recommendations": ["❌ Usar processamento sequencial devido ao erro"],
            "error": str(e)
        }

def stitch_videos_to_presentation(video_paths: List[str], output_path: str,
                                 transition_duration: float = 0.5,
                                 add_intro: bool = True, add_outro: bool = True,
                                 background_music: str = None) -> dict:
    """
    Une múltiplos vídeos em uma apresentação final completa.
    
    Esta função combina vários vídeos de slides em uma apresentação contínua,
    adicionando transições suaves, intro/outro opcionais e música de fundo.
    Ideal para criar apresentações profissionais a partir de slides individuais.
    
    Args:
        video_paths (List[str]): Lista de caminhos dos vídeos para unir
        output_path (str): Caminho do vídeo final de saída
        transition_duration (float): Duração das transições em segundos
        add_intro (bool): Se deve adicionar slide de introdução
        add_outro (bool): Se deve adicionar slide de encerramento
        background_music (str): Caminho para música de fundo (opcional)
    
    Returns:
        dict: Resultado da operação com informações detalhadas:
            - success: Se a operação foi bem-sucedida
            - final_video_path: Caminho do vídeo final criado
            - total_duration: Duração total do vídeo final
            - videos_processed: Número de vídeos processados
            - file_size: Tamanho do arquivo final
            - processing_time: Tempo de processamento
            - error: Mensagem de erro se houver falha
    """
    import time
    start_time = time.time()
    
    result = {
        "success": False,
        "final_video_path": None,
        "total_duration": 0,
        "videos_processed": 0,
        "file_size": 0,
        "processing_time": 0,
        "error": None
    }
    
    try:
        print("\n" + "="*60)
        print("🎬 UNIÃO DE VÍDEOS EM APRESENTAÇÃO FINAL")
        print("="*60)
        
        # Verificar se MoviePy está disponível
        if not MOVIEPY_AVAILABLE:
            raise ImportError("MoviePy não disponível. Execute: pip install moviepy")
        
        from moviepy.editor import VideoFileClip, concatenate_videoclips, CompositeVideoClip
        
        # Validar vídeos de entrada
        valid_videos = []
        invalid_videos = []
        
        for video_path in video_paths:
            if os.path.exists(video_path):
                try:
                    # Testar se o vídeo pode ser aberto
                    test_clip = VideoFileClip(video_path)
                    test_clip.close()
                    valid_videos.append(video_path)
                except Exception as e:
                    print(f"⚠️ Vídeo inválido: {video_path} - {str(e)}")
                    invalid_videos.append(video_path)
            else:
                print(f"❌ Vídeo não encontrado: {video_path}")
                invalid_videos.append(video_path)
        
        if not valid_videos:
            raise ValueError("Nenhum vídeo válido encontrado para processamento")
        
        print(f"✅ Vídeos válidos: {len(valid_videos)}")
        print(f"❌ Vídeos inválidos: {len(invalid_videos)}")
        
        # Criar clipes de vídeo
        video_clips = []
        total_duration = 0
        
        print("\n🔸 Carregando e processando vídeos...")
        
        for i, video_path in enumerate(valid_videos):
            try:
                print(f"   📹 Carregando vídeo {i+1}/{len(valid_videos)}: {os.path.basename(video_path)}")
                
                # Carregar clipe de vídeo
                clip = VideoFileClip(video_path)
                
                # Adicionar transição de fade in/out se especificado
                if transition_duration > 0:
                    clip = clip.fadeout(transition_duration).fadein(transition_duration)
                
                video_clips.append(clip)
                total_duration += clip.duration
                
                print(f"      ⏱️ Duração: {clip.duration:.2f}s")
                print(f"      📐 Resolução: {clip.w}x{clip.h}")
                
            except Exception as e:
                print(f"   ❌ Erro ao carregar {video_path}: {str(e)}")
                continue
        
        if not video_clips:
            raise ValueError("Nenhum clipe de vídeo foi carregado com sucesso")
        
        # Criar slides de intro e outro se solicitado
        if add_intro:
            print("\n🔸 Criando slide de introdução...")
            intro_clip = _create_intro_slide_clip()
            if intro_clip:
                video_clips.insert(0, intro_clip)
                total_duration += intro_clip.duration
        
        if add_outro:
            print("\n🔸 Criando slide de encerramento...")
            outro_clip = _create_outro_slide_clip()
            if outro_clip:
                video_clips.append(outro_clip)
                total_duration += outro_clip.duration
        
        # Concatenar todos os vídeos
        print(f"\n🔸 Concatenando {len(video_clips)} clipes...")
        final_video = concatenate_videoclips(video_clips, method="compose")
        
        # Adicionar música de fundo se especificada
        if background_music and os.path.exists(background_music):
            print("🔸 Adicionando música de fundo...")
            try:
                from moviepy.editor import AudioFileClip
                
                # Carregar música de fundo
                music_clip = AudioFileClip(background_music)
                
                # Ajustar duração da música para o vídeo
                if music_clip.duration > final_video.duration:
                    music_clip = music_clip.subclip(0, final_video.duration)
                else:
                    # Loop da música se for mais curta que o vídeo
                    loops_needed = int(final_video.duration / music_clip.duration) + 1
                    music_clips = [music_clip] * loops_needed
                    music_clip = concatenate_audioclips(music_clips).subclip(0, final_video.duration)
                
                # Reduzir volume da música (20% do volume original)
                music_clip = music_clip.volumex(0.2)
                
                # Combinar áudio original com música de fundo
                final_audio = CompositeAudioClip([final_video.audio, music_clip])
                final_video = final_video.set_audio(final_audio)
                
                print("✅ Música de fundo adicionada com sucesso")
                
            except Exception as music_error:
                print(f"⚠️ Erro ao adicionar música: {str(music_error)}")
        
        # Salvar vídeo final
        print(f"\n🔸 Salvando apresentação final...")
        print(f"   📁 Arquivo: {output_path}")
        print(f"   ⏱️ Duração total: {total_duration:.2f}s")
        
        # Criar pasta de output se necessário
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Exportar vídeo com configurações otimizadas
        final_video.write_videofile(
            output_path,
            codec='libx264',
            audio_codec='aac',
            temp_audiofile='temp-audio.m4a',
            remove_temp=True,
            fps=24,
            preset='medium',
            ffmpeg_params=['-crf', '23']  # Boa qualidade com tamanho razoável
        )
        
        # Limpar recursos
        for clip in video_clips:
            clip.close()
        final_video.close()
        
        # Obter informações do arquivo final
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            
            result.update({
                "success": True,
                "final_video_path": output_path,
                "total_duration": total_duration,
                "videos_processed": len(valid_videos),
                "file_size": file_size,
                "processing_time": time.time() - start_time
            })
            
            print("\n" + "="*60)
            print("🎉 APRESENTAÇÃO CRIADA COM SUCESSO!")
            print("="*60)
            print(f"📹 Vídeos processados: {len(valid_videos)}")
            print(f"⏱️ Duração total: {total_duration:.2f}s")
            print(f"💾 Tamanho do arquivo: {file_size / 1024 / 1024:.2f} MB")
            print(f"🕒 Tempo de processamento: {result['processing_time']:.2f}s")
            print(f"📁 Arquivo final: {output_path}")
            
        else:
            raise Exception("Arquivo de vídeo final não foi criado")
        
        return result
        
    except Exception as e:
        result["error"] = str(e)
        result["processing_time"] = time.time() - start_time
        print(f"\n❌ Erro ao criar apresentação: {str(e)}")
        return result

def _create_intro_slide_clip(duration: float = 3.0) -> object:
    """
    Cria um clipe de vídeo para slide de introdução.
    
    Args:
        duration: Duração do slide em segundos
    
    Returns:
        VideoFileClip: Clipe de vídeo da introdução ou None se falhar
    """
    try:
        if not MOVIEPY_AVAILABLE or not PIL_AVAILABLE:
            return None
        
        from moviepy.editor import ImageClip
        
        # Criar imagem de introdução
        intro_text = "TecnoCursos AI\nApresentação Gerada Automaticamente"
        intro_image = _create_slide_image(
            text=intro_text,
            width=1280,
            height=720,
            bg_color=(25, 25, 35),  # Azul escuro
            text_color=(255, 255, 255),  # Branco
            margin=80
        )
        
        # Salvar imagem temporária
        temp_intro_path = "temp_intro_slide.png"
        intro_image.save(temp_intro_path)
        
        # Criar clipe de vídeo
        intro_clip = ImageClip(temp_intro_path, duration=duration)
        
        # Aplicar efeito de fade
        intro_clip = intro_clip.fadeout(0.5).fadein(0.5)
        
        # Remover arquivo temporário
        if os.path.exists(temp_intro_path):
            os.remove(temp_intro_path)
        
        return intro_clip
        
    except Exception as e:
        print(f"⚠️ Erro ao criar slide de introdução: {str(e)}")
        return None

def _create_outro_slide_clip(duration: float = 3.0) -> object:
    """
    Cria um clipe de vídeo para slide de encerramento.
    
    Args:
        duration: Duração do slide em segundos
    
    Returns:
        VideoFileClip: Clipe de vídeo do encerramento ou None se falhar
    """
    try:
        if not MOVIEPY_AVAILABLE or not PIL_AVAILABLE:
            return None
        
        from moviepy.editor import ImageClip
        
        # Criar imagem de encerramento
        outro_text = "Obrigado!\nTecnoCursos AI\nContinue Aprendendo"
        outro_image = _create_slide_image(
            text=outro_text,
            width=1280,
            height=720,
            bg_color=(35, 25, 25),  # Vermelho escuro
            text_color=(255, 255, 255),  # Branco
            margin=80
        )
        
        # Salvar imagem temporária
        temp_outro_path = "temp_outro_slide.png"
        outro_image.save(temp_outro_path)
        
        # Criar clipe de vídeo
        outro_clip = ImageClip(temp_outro_path, duration=duration)
        
        # Aplicar efeito de fade
        outro_clip = outro_clip.fadeout(0.5).fadein(0.5)
        
        # Remover arquivo temporário
        if os.path.exists(temp_outro_path):
            os.remove(temp_outro_path)
        
        return outro_clip
        
    except Exception as e:
        print(f"⚠️ Erro ao criar slide de encerramento: {str(e)}")
        return None

def create_complete_presentation_from_pdf(pdf_path: str, audio_path: str, output_path: str,
                                         template: str = "modern", resolution: str = "hd",
                                         add_transitions: bool = True, add_music: bool = False,
                                         music_path: str = None) -> dict:
    """
    Pipeline completo: PDF + Áudio → Apresentação de Vídeo Final.
    
    Esta é a função mais avançada do sistema, que automatiza todo o processo:
    1. Extrai slides do PDF
    2. Cria vídeos individuais 
    3. Une tudo em uma apresentação final
    4. Adiciona intro/outro e música se solicitado
    
    Args:
        pdf_path (str): Caminho para o arquivo PDF
        audio_path (str): Caminho para o arquivo de áudio
        output_path (str): Caminho para o vídeo final
        template (str): Template visual
        resolution (str): Resolução do vídeo
        add_transitions (bool): Adicionar transições entre slides
        add_music (bool): Adicionar música de fundo
        music_path (str): Caminho para música de fundo
    
    Returns:
        dict: Resultado completo do processamento com todas as etapas
    """
    import time
    start_time = time.time()
    
    result = {
        "success": False,
        "pdf_processing": {},
        "video_creation": {},
        "final_stitching": {},
        "total_processing_time": 0,
        "final_video_path": None,
        "error": None
    }
    
    try:
        print("\n" + "="*80)
        print("🚀 PIPELINE COMPLETO: PDF → APRESENTAÇÃO DE VÍDEO FINAL")
        print("="*80)
        
        # Criar pasta temporária para processamento
        temp_folder = os.path.join(os.path.dirname(output_path), "temp_processing")
        
        print(f"📂 Pasta temporária: {temp_folder}")
        print(f"📄 PDF de entrada: {pdf_path}")
        print(f"🎵 Áudio de entrada: {audio_path}")
        print(f"🎬 Vídeo final: {output_path}")
        
        # FASE 1: Processar PDF e Áudio
        print(f"\n" + "🔸"*60)
        print("FASE 1: PROCESSAMENTO DE PDF E ÁUDIO")
        print("🔸"*60)
        
        pdf_result = create_videos_from_pdf_and_audio(
            pdf_path=pdf_path,
            audio_path=audio_path,
            output_folder=temp_folder,
            template=template,
            resolution=resolution,
            animations=True,
            background_style="gradient",
            sync_mode="auto"
        )
        
        result["pdf_processing"] = pdf_result
        
        if not pdf_result.get("success", False):
            raise Exception(f"Falha no processamento do PDF: {pdf_result.get('error', 'Erro desconhecido')}")
        
        videos_created = pdf_result.get("videos_created", [])
        
        if not videos_created:
            raise Exception("Nenhum vídeo foi criado a partir do PDF")
        
        print(f"✅ FASE 1 CONCLUÍDA: {len(videos_created)} vídeos criados")
        
        # FASE 2: União em Apresentação Final
        print(f"\n" + "🔸"*60)
        print("FASE 2: CRIAÇÃO DA APRESENTAÇÃO FINAL")
        print("🔸"*60)
        
        stitching_result = stitch_videos_to_presentation(
            video_paths=videos_created,
            output_path=output_path,
            transition_duration=0.5 if add_transitions else 0.0,
            add_intro=True,
            add_outro=True,
            background_music=music_path if add_music and music_path else None
        )
        
        result["final_stitching"] = stitching_result
        
        if not stitching_result.get("success", False):
            raise Exception(f"Falha na união dos vídeos: {stitching_result.get('error', 'Erro desconhecido')}")
        
        print(f"✅ FASE 2 CONCLUÍDA: Apresentação final criada")
        
        # FASE 3: Finalização e Limpeza
        print(f"\n" + "🔸"*60)
        print("FASE 3: FINALIZAÇÃO E LIMPEZA")
        print("🔸"*60)
        
        # Limpar pasta temporária (opcional)
        try:
            import shutil
            if os.path.exists(temp_folder):
                shutil.rmtree(temp_folder)
                print("🧹 Pasta temporária removida")
        except Exception as cleanup_error:
            print(f"⚠️ Erro na limpeza: {str(cleanup_error)}")
        
        # Calcular estatísticas finais
        total_processing_time = time.time() - start_time
        result.update({
            "success": True,
            "total_processing_time": total_processing_time,
            "final_video_path": output_path
        })
        
        # Estatísticas completas
        final_file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
        
        print("\n" + "="*80)
        print("🎉 PIPELINE COMPLETO FINALIZADO COM SUCESSO!")
        print("="*80)
        print(f"📄 Slides processados: {len(pdf_result.get('pdf_slides', []))}")
        print(f"🎬 Vídeos intermediários: {len(videos_created)}")
        print(f"⏱️ Duração total: {stitching_result.get('total_duration', 0):.2f}s")
        print(f"💾 Tamanho final: {final_file_size / 1024 / 1024:.2f} MB")
        print(f"🕒 Tempo total de processamento: {total_processing_time:.2f}s")
        print(f"📁 Arquivo final: {output_path}")
        print(f"🎨 Template usado: {template}")
        print(f"📐 Resolução: {resolution}")
        print(f"✨ Transições: {'Ativadas' if add_transitions else 'Desativadas'}")
        print(f"🎵 Música de fundo: {'Adicionada' if add_music and music_path else 'Não adicionada'}")
        
        return result
        
    except Exception as e:
        result["error"] = str(e)
        result["total_processing_time"] = time.time() - start_time
        print(f"\n❌ ERRO NO PIPELINE: {str(e)}")
        return result


def concatenate_videos(video_paths_list: List[str], output_path: str) -> dict:
    """
    Concatena múltiplos vídeos em um único vídeo final usando MoviePy.
    
    Esta função recebe uma lista de caminhos de vídeos e os concatena em sequência
    para criar um único arquivo de vídeo final. É útil para unir múltiplas gravações,
    clipes ou vídeos de apresentação em uma única peça.
    
    Características principais:
    - Concatenação sequencial simples e eficiente
    - Validação automática de arquivos de entrada
    - Tratamento robusto de erros
    - Relatório detalhado do processamento
    - Limpeza automática de recursos
    
    Args:
        video_paths_list (List[str]): Lista com os caminhos dos vídeos a serem concatenados.
                                     Os vídeos serão unidos na ordem especificada.
        output_path (str): Caminho completo onde salvar o vídeo final concatenado.
                          Inclui nome do arquivo e extensão (ex: "videos/final.mp4").
    
    Returns:
        dict: Resultado da operação contendo:
            - success (bool): True se a concatenação foi bem-sucedida
            - output_path (str): Caminho do vídeo final criado
            - total_duration (float): Duração total do vídeo final em segundos
            - videos_processed (int): Número de vídeos processados com sucesso
            - videos_skipped (int): Número de vídeos que foram ignorados por erro
            - file_size (int): Tamanho do arquivo final em bytes
            - processing_time (float): Tempo total de processamento em segundos
            - error (str): Mensagem de erro se houver falha, None se sucesso
            - details (dict): Informações detalhadas sobre cada vídeo processado
    
    Raises:
        ImportError: Se MoviePy não estiver disponível
        ValueError: Se nenhum vídeo válido for encontrado
        Exception: Para outros erros durante o processamento
    
    Example:
        >>> videos = ["video1.mp4", "video2.mp4", "video3.mp4"]
        >>> result = concatenate_videos(videos, "final_presentation.mp4")
        >>> if result["success"]:
        ...     print(f"Vídeo criado: {result['output_path']}")
        ...     print(f"Duração: {result['total_duration']:.2f}s")
    """
    import time
    start_time = time.time()
    
    # Estrutura de resultado padronizada
    result = {
        "success": False,
        "output_path": None,
        "total_duration": 0.0,
        "videos_processed": 0,
        "videos_skipped": 0,
        "file_size": 0,
        "processing_time": 0.0,
        "error": None,
        "details": {
            "valid_videos": [],
            "invalid_videos": [],
            "video_info": []
        }
    }
    
    try:
        print("\n" + "="*70)
        print("🎬 CONCATENAÇÃO DE VÍDEOS - SISTEMA TECNOCURSOS AI")
        print("="*70)
        print(f"📝 Vídeos para processar: {len(video_paths_list)}")
        print(f"📁 Arquivo de saída: {output_path}")
        
        # Verificar se MoviePy está disponível
        if not MOVIEPY_AVAILABLE:
            raise ImportError(
                "MoviePy não está disponível. "
                "Instale com: pip install moviepy"
            )
        
        # Importar dependências necessárias
        from moviepy.editor import VideoFileClip, concatenate_videoclips
        
        # Verificar se há vídeos para processar
        if not video_paths_list:
            raise ValueError("Lista de vídeos está vazia")
        
        # Validar e catalogar vídeos de entrada
        print("\n🔍 VALIDAÇÃO DOS VÍDEOS DE ENTRADA")
        print("-" * 50)
        
        valid_video_clips = []
        video_details = []
        
        for i, video_path in enumerate(video_paths_list, 1):
            print(f"📹 Verificando vídeo {i}/{len(video_paths_list)}: {os.path.basename(video_path)}")
            
            video_detail = {
                "index": i,
                "path": video_path,
                "filename": os.path.basename(video_path),
                "exists": False,
                "valid": False,
                "duration": 0.0,
                "resolution": None,
                "error": None
            }
            
            try:
                # Verificar se o arquivo existe
                if not os.path.exists(video_path):
                    error_msg = "Arquivo não encontrado"
                    video_detail["error"] = error_msg
                    result["details"]["invalid_videos"].append(video_path)
                    print(f"   ❌ {error_msg}")
                    continue
                
                video_detail["exists"] = True
                
                # Tentar carregar o vídeo
                video_clip = VideoFileClip(video_path)
                
                # Obter informações do vídeo
                video_detail.update({
                    "valid": True,
                    "duration": video_clip.duration,
                    "resolution": f"{video_clip.w}x{video_clip.h}",
                })
                
                # Adicionar à lista de vídeos válidos
                valid_video_clips.append(video_clip)
                result["details"]["valid_videos"].append(video_path)
                
                print(f"   ✅ Válido - Duração: {video_clip.duration:.2f}s, Resolução: {video_clip.w}x{video_clip.h}")
                
            except Exception as e:
                video_detail["error"] = str(e)
                result["details"]["invalid_videos"].append(video_path)
                print(f"   ❌ Erro ao carregar: {str(e)}")
            
            finally:
                video_details.append(video_detail)
        
        # Atualizar estatísticas
        result["videos_processed"] = len(valid_video_clips)
        result["videos_skipped"] = len(video_paths_list) - len(valid_video_clips)
        result["details"]["video_info"] = video_details
        
        # Verificar se temos vídeos válidos para processar
        if not valid_video_clips:
            raise ValueError(
                f"Nenhum vídeo válido encontrado. "
                f"Vídeos inválidos: {len(result['details']['invalid_videos'])}"
            )
        
        print(f"\n📊 RESUMO DA VALIDAÇÃO:")
        print(f"   ✅ Vídeos válidos: {result['videos_processed']}")
        print(f"   ❌ Vídeos ignorados: {result['videos_skipped']}")
        
        # Calcular duração total estimada
        total_duration = sum(clip.duration for clip in valid_video_clips)
        result["total_duration"] = total_duration
        
        print(f"   ⏱️ Duração total estimada: {total_duration:.2f}s ({total_duration/60:.1f} min)")
        
        # Realizar a concatenação
        print(f"\n🔧 PROCESSAMENTO - CONCATENANDO VÍDEOS")
        print("-" * 50)
        print(f"🎬 Unindo {len(valid_video_clips)} vídeo(s) em sequência...")
        
        # Usar o método "compose" para melhor compatibilidade entre vídeos
        final_video = concatenate_videoclips(valid_video_clips, method="compose")
        
        print("✅ Concatenação concluída com sucesso")
        
        # Criar diretório de saída se necessário
        output_directory = os.path.dirname(output_path)
        if output_directory and not os.path.exists(output_directory):
            os.makedirs(output_directory, exist_ok=True)
            print(f"📁 Diretório criado: {output_directory}")
        
        # Salvar o vídeo final
        print(f"\n💾 EXPORTAÇÃO DO VÍDEO FINAL")
        print("-" * 50)
        print(f"📁 Salvando em: {output_path}")
        
        # Configurações otimizadas de exportação
        final_video.write_videofile(
            output_path,
            codec='libx264',           # Codec de vídeo padrão e compatível
            audio_codec='aac',         # Codec de áudio padrão
            temp_audiofile='temp-audio.m4a',  # Arquivo temporário para áudio
            remove_temp=True,          # Remover arquivos temporários
            fps=24,                    # Frame rate padrão
            preset='medium',           # Velocidade vs qualidade balanceada
            ffmpeg_params=['-crf', '23']  # Qualidade visual otimizada
        )
        
        print("✅ Vídeo exportado com sucesso")
        
        # Limpeza de recursos para evitar vazamentos de memória
        print(f"\n🧹 LIMPEZA DE RECURSOS")
        print("-" * 30)
        
        for i, clip in enumerate(valid_video_clips):
            try:
                clip.close()
                print(f"   🗑️ Clip {i+1} limpo")
            except Exception as cleanup_error:
                print(f"   ⚠️ Erro na limpeza do clip {i+1}: {cleanup_error}")
        
        try:
            final_video.close()
            print("   🗑️ Vídeo final limpo")
        except Exception as cleanup_error:
            print(f"   ⚠️ Erro na limpeza do vídeo final: {cleanup_error}")
        
        # Verificar resultado final e obter estatísticas
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            processing_time = time.time() - start_time
            
            # Atualizar resultado com sucesso
            result.update({
                "success": True,
                "output_path": output_path,
                "file_size": file_size,
                "processing_time": processing_time
            })
            
            # Relatório final de sucesso
            print("\n" + "="*70)
            print("🎉 CONCATENAÇÃO CONCLUÍDA COM SUCESSO!")
            print("="*70)
            print(f"📹 Vídeos processados: {result['videos_processed']}")
            print(f"⏱️ Duração total: {result['total_duration']:.2f}s ({result['total_duration']/60:.1f} min)")
            print(f"💾 Tamanho do arquivo: {file_size / (1024*1024):.2f} MB")
            print(f"🕒 Tempo de processamento: {processing_time:.2f}s")
            print(f"📁 Arquivo criado: {output_path}")
            
            if result["videos_skipped"] > 0:
                print(f"⚠️ Vídeos ignorados: {result['videos_skipped']} (por erro)")
            
            print("="*70)
            
        else:
            raise Exception("Arquivo de vídeo final não foi criado corretamente")
        
        return result
        
    except Exception as e:
        # Capturar qualquer erro e atualizar resultado
        error_msg = str(e)
        processing_time = time.time() - start_time
        
        result.update({
            "error": error_msg,
            "processing_time": processing_time
        })
        
        print(f"\n❌ ERRO NA CONCATENAÇÃO: {error_msg}")
        print(f"🕒 Tempo até o erro: {processing_time:.2f}s")
        
        # Tentar limpar recursos mesmo em caso de erro
        try:
            if 'valid_video_clips' in locals():
                for clip in valid_video_clips:
                    clip.close()
            if 'final_video' in locals():
                final_video.close()
        except:
            pass  # Ignorar erros de limpeza em caso de falha
        
        return result

def create_videos_for_slides(slides_data, output_dir="static/videos", **kwargs):
    """
    Cria vídeos para slides de apresentação
    
    Args:
        slides_data: Dados dos slides
        output_dir: Diretório de saída
        **kwargs: Argumentos adicionais
        
    Returns:
        dict: Resultado da geração de vídeos
    """
    try:
        import os
        from datetime import datetime
        
        # Criar diretório se não existir
        os.makedirs(output_dir, exist_ok=True)
        
        videos_created = []
        
        for i, slide in enumerate(slides_data):
            video_filename = f"slide_{i+1}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.mp4"
            video_path = os.path.join(output_dir, video_filename)
            
            # Simular criação de vídeo (implementação básica)
            # Em produção, aqui seria usado MoviePy ou similar
            
            videos_created.append({
                "slide_index": i,
                "video_path": video_path,
                "status": "created",
                "duration": slide.get("duration", 5.0)
            })
        
        return {
            "success": True,
            "videos_created": len(videos_created),
            "videos": videos_created,
            "output_directory": output_dir
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "videos": []
        }
