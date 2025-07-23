from pptx import Presentation
import os
from typing import List, Dict
from PIL import Image

def parse_pptx(file_path: str, output_dir: str) -> List[Dict]:
    """
    Extrai texto e imagens de cada slide de um PPTX.
    Retorna lista de slides: [{ 'texto': str, 'imagens': [str] }]
    Salva imagens extraídas em output_dir.
    """
    os.makedirs(output_dir, exist_ok=True)
    slides = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        texto = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text])
        imagens = []
        for img_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                ext = image.ext
                img_path = os.path.join(output_dir, f"slide_{i+1}_img_{img_idx+1}.{ext}")
                with open(img_path, "wb") as f:
                    f.write(image.blob)
                imagens.append(img_path)
        slides.append({"texto": texto, "imagens": imagens})
    return slides

# Para adicionar novos formatos, crie um novo arquivo parser e siga a mesma interface: parse_<formato>(file_path, output_dir) -> List[Dict]
